"""
ETRA Design System tokens for PowerPoint builders.
Source: ETRA-Design-System.pdf (v1.0)
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

_FORMULA_SUB = re.compile(r"_\{([^}]+)\}")

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "public" / "assets" / "etra-wordmark.png"
LOGO_FALLBACK = ROOT / "public" / "assets" / "etra-logo.png"
DIAGRAMS = ROOT / "public" / "assets" / "session1-diagrams"
FONT_DIR = ROOT / "public" / "font" / "din-next"

# ── Colors (page 06 / 07) ──────────────────────────────────────────────────
PRIMARY = RGBColor(0x52, 0x34, 0xB7)       # --etra-primary
SECONDARY = RGBColor(0x9E, 0x59, 0xCD)     # --etra-secondary
SURFACE = RGBColor(0xFA, 0xF8, 0xFF)       # --etra-surface / light content bg
MUTED = RGBColor(0x5A, 0x54, 0x70)         # --etra-muted
INK = RGBColor(0x12, 0x10, 0x18)           # --etra-ink
BLACK = RGBColor(0x0A, 0x06, 0x14)         # --etra-black
SOFT = RGBColor(0xEE, 0xE8, 0xFA)          # --etra-soft
SOFT_2 = RGBColor(0xF5, 0xF1, 0xFC)        # --etra-soft-2
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE4, 0xDC, 0xF4)

# ── Layout (pages 09 / 13) ─────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)  # 0.6" all sides
RAIL_W = Inches(0.06)  # thin accent bar on right edge

# Font family: DIN Next LT W23 (ETRA Sans). Falls back to Helvetica name if needed.
FONT = "DIN Next LT W23"
FONT_FALLBACK = "Helvetica"


def _font_name() -> str:
    # PowerPoint resolves installed fonts by family name; DIN files are present on this machine.
    return FONT


def no_shadow(shape) -> None:
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith("effectLst"):
            spPr.remove(child)
    spPr.append(spPr.makeelement(qn("a:effectLst"), {}))


# Formulas use a widely available face so subscripts/baselines render reliably.
FORMULA_FONT = "Helvetica Neue"
FORMULA_CACHE = ROOT / "public" / "assets" / "formula-cache"


def set_run(run, size, *, bold=False, color=INK, font=None, subscript=False):
    fname = font or _font_name()
    run.font.name = fname
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", fname)
    # Remove incorrect child-element baselines from older builds
    for el in rPr.findall(qn("a:baseline")):
        rPr.remove(el)
    # OOXML: baseline is an attribute on a:rPr (−25000 = subscript)
    if subscript:
        rPr.set("baseline", "-25000")
    elif "baseline" in rPr.attrib:
        del rPr.attrib["baseline"]


def fill_formula_paragraph(paragraph, value, *, size=18, bold=True, color=PRIMARY, font=FORMULA_FONT):
    """Render formula markup. Use _{sub} for subscripts, e.g. x_{min}."""
    p_elem = paragraph._p
    for child in list(p_elem):
        if child.tag.endswith("}r"):
            p_elem.remove(child)

    pos = 0
    for match in _FORMULA_SUB.finditer(value):
        if match.start() > pos:
            run = paragraph.add_run()
            run.text = value[pos : match.start()]
            set_run(run, size, bold=bold, color=color, font=font)
        run = paragraph.add_run()
        run.text = match.group(1)
        set_run(run, max(size - 4, 12), bold=bold, color=color, font=font, subscript=True)
        pos = match.end()
    if pos < len(value):
        run = paragraph.add_run()
        run.text = value[pos:]
        set_run(run, size, bold=bold, color=color, font=font)


_FRAC_RE = re.compile(
    r"^(?P<lhs>.+?)\s*=\s*(?:(?P<prefix>1\s*[−\-]\s*))?\((?P<num>.+?)\)\s*/\s*\((?P<den>.+?)\)$"
)

# Average-sum: lhs = (1 / den) Σ_{limits} term
_AVG_SUM_RE = re.compile(
    r"^(?P<lhs>.+?)\s*=\s*\(\s*1\s*/\s*(?P<den>[^)]+?)\s*\)\s*"
    r"(?P<sum>Σ|∑)(?:_\{(?P<lo>[^}]*)\})?(?:\^\{(?P<hi>[^}]*)\})?\s*"
    r"(?P<term>.+)$"
)

# Inline division with optional 1− prefix: lhs = [1 −] (num) / (den)
_INLINE_DIV_RE = re.compile(
    r"^(?P<lhs>.+?)\s*=\s*(?P<prefix>1\s*[−\-]\s*)?\((?P<num>.+?)\)\s*/\s*\((?P<den>.+?)\)$"
)

# Product-style adj R²: lhs = 1 − (a)(b)/(c)  or  1 − (a)(b) / (c)
_ADJ_R2_RE = re.compile(
    r"^(?P<lhs>.+?)\s*=\s*1\s*[−\-]\s*"
    r"\((?P<a>.+?)\)\s*\((?P<b>.+?)\)\s*/\s*\((?P<c>.+?)\)$"
)


def formula_to_latex(value: str) -> str:
    """Convert presentation formula markup / unicode math into mathtext LaTeX."""
    s = (value or "").strip()
    if not s:
        return s
    if s.startswith("$") and s.endswith("$"):
        return s[1:-1]
    # Already looks like LaTeX
    if "\\" in s and any(tok in s for tok in (r"\frac", r"\sum", r"\hat", r"\min", r"\lambda")):
        return s

    # Prefer structured patterns before generic unicode swaps
    m = _AVG_SUM_RE.match(s)
    if m:
        lhs = _token_to_latex(m.group("lhs"))
        den = _token_to_latex(m.group("den").strip())
        lo = _token_to_latex((m.group("lo") or "").strip())
        hi = _token_to_latex((m.group("hi") or "").strip())
        term = _token_to_latex(m.group("term").strip())
        # Infer common RF limits when missing
        if not lo and not hi and "N" in den:
            lo, hi = "t=1", den
        if lo and hi:
            limits = "_{%s}^{%s}" % (lo, hi)
        elif lo:
            limits = "_{%s}" % lo
        elif hi:
            limits = "^{%s}" % hi
        else:
            limits = ""
        return rf"{lhs} = \dfrac{{1}}{{{den}}} \sum{limits} {term}"

    m = _ADJ_R2_RE.match(s)
    if m:
        lhs = _token_to_latex(m.group("lhs"))
        a = _token_to_latex(m.group("a"))
        b = _token_to_latex(m.group("b"))
        c = _token_to_latex(m.group("c"))
        return rf"{lhs} = 1 - \dfrac{{({a})({b})}}{{{c}}}"

    m = _INLINE_DIV_RE.match(s)
    if m:
        lhs = _token_to_latex(m.group("lhs"))
        prefix = (m.group("prefix") or "").replace("-", "−").strip()
        num = _token_to_latex(m.group("num"))
        den = _token_to_latex(m.group("den"))
        if prefix:
            return rf"{lhs} = 1 - \dfrac{{{num}}}{{{den}}}"
        return rf"{lhs} = \dfrac{{{num}}}{{{den}}}"

    # Generic unicode / markup conversion
    return _token_to_latex(s)


def _token_to_latex(token: str) -> str:
    s = token.strip()
    if not s:
        return s

    # Before unicode greek swaps (β → \beta would break these)
    s = s.replace("min_β", r"\min_{\beta}")
    s = s.replace("min_{w}", r"\min_{w}")
    s = re.sub(r"\bmin\b", lambda _: r"\min", s)

    # Protect existing latex-ish braces content later via placeholders
    reps = [
        ("ŷ", r"\hat{y}"),
        ("ȳ", r"\bar{y}"),
        ("Σ", r"\sum"),
        ("∑", r"\sum"),
        ("∥", r"\|"),
        ("≤", r"\leq"),
        ("≥", r"\geq"),
        ("⇒", r"\Rightarrow"),
        ("⋯", r"\cdots"),
        ("…", r"\ldots"),
        ("−", r"-"),
        ("½", r"\tfrac{1}{2}"),
        ("ᵀ", r"^{\top}"),
        ("φ", r"\varphi"),
        ("ε", r"\varepsilon"),
        ("ξ", r"\xi"),
        ("λ", r"\lambda"),
        ("β", r"\beta"),
        ("Ω", r"\Omega"),
        ("∈", r"\in"),
        ("ᵢ", r"_i"),
        ("ⱼ", r"_j"),
        ("ₜ", r"_t"),
        ("²", r"^2"),
        ("³", r"^3"),
    ]
    # Handle R² before generic ² if still present
    s = s.replace("R²", r"R^2")
    for a, b in reps:
        s = s.replace(a, b)
    # Trailing star (slack variables): ξᵢ* → \xi_i^{*}
    s = re.sub(r"(?<=[\w}])\*", r"^{*}", s)

    # N_trees / S_leaf style identifiers → N_{trees}
    s = re.sub(r"\bN_trees\b", "N_{trees}", s)
    s = re.sub(r"\bS_leaf\b", "S_{leaf}", s)
    s = re.sub(r"\bMSE_split\b", lambda _: r"\mathrm{MSE}_{split}", s)
    s = re.sub(r"\bMSE\b", lambda _: r"\mathrm{MSE}", s)
    s = re.sub(r"\bSS_res\b", "SS_{res}", s)
    s = re.sub(r"\bSS_tot\b", "SS_{tot}", s)

    # Σ_{i = 1 … n} or sum_{i = 1 … n} already partially converted
    s = re.sub(
        r"\\sum_\{i\s*=\s*1\s*\\ldots\s*n\}",
        lambda _: r"\sum_{i=1}^{n}",
        s,
    )
    s = re.sub(
        r"\\sum_\{i\s*=\s*1\s*\.\.\.\s*n\}",
        lambda _: r"\sum_{i=1}^{n}",
        s,
    )
    s = re.sub(r"\\sum_\{i\s*\\in\s*S\}", lambda _: r"\sum_{i \in S}", s)
    s = re.sub(
        r"\\sum_\{i\s*\\in\s*S_\{leaf\}\}",
        lambda _: r"\sum_{i \in S_{leaf}}",
        s,
    )

    # Bare sum without limits after average: leave as \sum
    # Spaces around operators
    s = re.sub(r"\s+", " ", s).strip()
    return s


def render_formula_png(
    latex: str,
    *,
    fontsize: float = 30,
    color: RGBColor | None = None,
    cache_key: str | None = None,
) -> Path | None:
    """Render mathtext LaTeX to a transparent PNG (serif math look)."""
    import hashlib

    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    tex = latex.strip().strip("$")
    if not tex:
        return None

    ink = color or INK
    hex_color = f"#{ink[0]:02x}{ink[1]:02x}{ink[2]:02x}"
    key_src = f"{tex}|{fontsize}|{hex_color}"
    digest = hashlib.sha1((cache_key or key_src).encode("utf-8")).hexdigest()[:16]
    FORMULA_CACHE.mkdir(parents=True, exist_ok=True)
    out = FORMULA_CACHE / f"f-{digest}.png"
    if out.is_file() and out.stat().st_size > 200:
        return out

    fig = plt.figure(figsize=(10, 1.6))
    fig.patch.set_alpha(0.0)
    try:
        fig.text(
            0.5,
            0.5,
            f"${tex}$",
            fontsize=fontsize,
            ha="center",
            va="center",
            color=hex_color,
        )
        fig.savefig(
            out,
            dpi=220,
            transparent=True,
            bbox_inches="tight",
            pad_inches=0.18,
        )
    except Exception:
        plt.close(fig)
        return None
    plt.close(fig)
    return out if out.is_file() else None


def to_inches(val) -> float:
    """
    Normalize a length to inches.
    Accepts pptx Length (has .inches), EMU ints from Length arithmetic, or float inches.
    """
    if hasattr(val, "inches"):
        return float(val.inches)
    n = float(val)
    # python-pptx Length + Length yields EMU ints (~914400 per inch)
    if abs(n) >= 1000:
        return n / 914400.0
    return n


def _add_formula_image(slide, left, top, width, height, path: Path):
    """Center a formula PNG inside the given box."""
    from PIL import Image

    with Image.open(path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return None
    aspect = px_w / px_h
    max_w = to_inches(width)
    max_h = to_inches(height)
    left_in = to_inches(left)
    top_in = to_inches(top)
    fit_w = min(max_w, max_h * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h:
        fit_h = max_h
        fit_w = fit_h * aspect
    x = left_in + (max_w - fit_w) / 2
    y = top_in + (max_h - fit_h) / 2
    return slide.shapes.add_picture(
        str(path),
        Inches(x),
        Inches(y),
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def add_formula(
    slide,
    left,
    top,
    width,
    height,
    value,
    *,
    size=18,
    bold=True,
    color=PRIMARY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    formula_tex: str | None = None,
    math_image: bool = True,
):
    """
    Add a display formula.
    Prefer LaTeX-style math PNG (stacked fractions, Σ limits) over plain text.
    Markup: x_{min}; optional formula_tex for exact mathtext.
    """
    text = value if isinstance(value, str) else str(value)
    tex = (formula_tex or "").strip() or formula_to_latex(text)

    if math_image and tex:
        # Larger type for short formulas; slightly smaller for long ones
        fs = 34 if len(tex) < 40 else 28 if len(tex) < 70 else 24
        path = render_formula_png(tex, fontsize=fs, color=INK)
        if path is not None:
            pic = _add_formula_image(slide, left, top, width, height, path)
            if pic is not None:
                return pic

    # Fallback: stacked-fraction shapes or plain text runs
    match = _FRAC_RE.match(text.strip())
    if match:
        return add_fraction_formula(
            slide,
            left,
            top,
            width,
            height,
            match.group("lhs"),
            match.group("num"),
            match.group("den"),
            size=size,
            bold=bold,
            color=color,
            prefix=(match.group("prefix") or "").replace("-", "−").strip(),
        )

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    fill_formula_paragraph(p, text, size=size, bold=bold, color=color)
    return box


def add_fraction_formula(
    slide,
    left,
    top,
    width,
    height,
    lhs,
    numerator,
    denominator,
    *,
    size=18,
    bold=True,
    color=PRIMARY,
    prefix="",
):
    """Visual stacked fraction: lhs = [prefix] num / den (text fallback)."""
    lhs_w = width * 0.34 if prefix else width * 0.28
    lhs_box = slide.shapes.add_textbox(left, top, lhs_w, height)
    lhs_tf = lhs_box.text_frame
    lhs_tf.word_wrap = False
    try:
        lhs_tf._txBody.bodyPr.set("anchor", "ctr")
    except Exception:
        pass
    lhs_p = lhs_tf.paragraphs[0]
    lhs_p.alignment = PP_ALIGN.RIGHT
    lhs_text = f"{lhs} = {prefix}" if prefix else f"{lhs} ="
    fill_formula_paragraph(lhs_p, lhs_text, size=size, bold=bold, color=color)

    frac_left = left + lhs_w + width * 0.02
    frac_w = width - lhs_w - width * 0.02
    num_h = height * 0.38
    den_h = height * 0.38
    gap = height * 0.12

    num_box = slide.shapes.add_textbox(frac_left, top, frac_w, num_h)
    num_tf = num_box.text_frame
    num_tf.word_wrap = False
    try:
        num_tf._txBody.bodyPr.set("anchor", "b")
    except Exception:
        pass
    num_p = num_tf.paragraphs[0]
    num_p.alignment = PP_ALIGN.CENTER
    fill_formula_paragraph(num_p, numerator, size=size, bold=bold, color=color)

    line_y = top + num_h + gap * 0.35
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        frac_left + frac_w * 0.05,
        line_y,
        frac_w * 0.90,
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    no_shadow(line)

    den_box = slide.shapes.add_textbox(frac_left, top + num_h + gap, frac_w, den_h)
    den_tf = den_box.text_frame
    den_tf.word_wrap = False
    try:
        den_tf._txBody.bodyPr.set("anchor", "t")
    except Exception:
        pass
    den_p = den_tf.paragraphs[0]
    den_p.alignment = PP_ALIGN.CENTER
    fill_formula_paragraph(den_p, denominator, size=size, bold=bold, color=color)

    return lhs_box


def add_text(
    slide,
    left,
    top,
    width,
    height,
    value,
    *,
    size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    set_run(run, size, bold=bold, color=color)
    return box


def rect(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    no_shadow(shape)
    return shape


def soft_card(slide, left, top, width, height, *, fill=SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        # ~18px radius feel on 1920 canvas
        shape.adjustments[0] = 0.1
    except Exception:
        pass
    no_shadow(shape)
    return shape


def gradient_fill(shape, c1: RGBColor, c2: RGBColor, angle: float = 135.0):
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].color.rgb = c1
    stops[1].color.rgb = c2


def paint_light(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = SURFACE


def paint_dark(slide):
    """Flat dark hero — no decorative circles/orbs over content."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BLACK


def right_rail(slide, *, dark=False):
    """Brand accent strip on the right edge (content slides)."""
    color = SECONDARY if dark else PRIMARY
    return rect(slide, SLIDE_W - RAIL_W, 0, RAIL_W, SLIDE_H, color)


def logo(slide, *, dark=False, height=Inches(0.35)):
    path = LOGO if LOGO.is_file() else LOGO_FALLBACK
    if not path.is_file():
        return
    # Keep logo in top-right safe area (page 13)
    pic = slide.shapes.add_picture(str(path), Inches(11.55), Inches(0.45), height=height)
    return pic


def add_diagram(slide, name, left, top, width, max_height):
    """
    Embed a Session 1 classroom diagram PNG fitted inside a width × max_height box.
    Preserves aspect ratio and centers horizontally within the box.
    """
    from PIL import Image

    path = DIAGRAMS / name
    if not path.is_file() and not str(name).endswith(".png"):
        path = DIAGRAMS / f"{name}.png"
    if not path.is_file():
        return None

    with Image.open(path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return None

    aspect = px_w / px_h
    max_w_in = width.inches if hasattr(width, "inches") else float(width)
    max_h_in = max_height.inches if hasattr(max_height, "inches") else float(max_height)
    left_in = left.inches if hasattr(left, "inches") else float(left)
    top_in = top.inches if hasattr(top, "inches") else float(top)

    fit_w_in = min(max_w_in, max_h_in * aspect)
    fit_h_in = fit_w_in / aspect
    if fit_h_in > max_h_in:
        fit_h_in = max_h_in
        fit_w_in = fit_h_in * aspect

    x_in = left_in + (max_w_in - fit_w_in) / 2
    return slide.shapes.add_picture(
        str(path),
        Inches(x_in),
        Inches(top_in),
        width=Inches(fit_w_in),
        height=Inches(fit_h_in),
    )


def content_header(slide, kicker: str, slide_num: str):
    """Light content chrome: logo, kicker, hairline — no heavy bars."""
    logo(slide)
    add_text(slide, MARGIN, Inches(0.48), Inches(8.5), Inches(0.3), kicker, size=12, color=MUTED)
    rect(slide, MARGIN, Inches(0.9), Inches(12.1), Inches(0.012), LINE)


def content_footer(slide, page: int, total: int):
    """Footer: ETRA left · page right, thin primary line above."""
    y = Inches(7.05)
    rect(slide, MARGIN, y, Inches(12.1), Inches(0.012), PRIMARY)
    add_text(slide, MARGIN, Inches(7.12), Inches(3), Inches(0.28), "ETRA", size=11, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(10.5),
        Inches(7.12),
        Inches(2.2),
        Inches(0.28),
        f"{page:02d}  /  {total:02d}",
        size=11,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def title_block(slide, title: str, subtitle: str | None = None, *, top=Inches(1.15)):
    add_text(slide, MARGIN, top, Inches(11.2), Inches(0.65), title, size=32, bold=True, color=PRIMARY)
    if subtitle:
        add_text(slide, MARGIN, top + Inches(0.6), Inches(11.2), Inches(0.4), subtitle, size=15, color=MUTED)


def is_fraction_formula(value: str) -> bool:
    """True when the formula needs a taller display box (fraction / sum / LaTeX)."""
    s = (value or "").strip()
    if not s:
        return False
    if _FRAC_RE.match(s) or _AVG_SUM_RE.match(s) or _ADJ_R2_RE.match(s):
        return True
    if any(tok in s for tok in ("Σ", "∑", "/", r"\frac", r"\sum", r"\dfrac")):
        return True
    return False


def bullets(
    slide,
    items: list[str],
    *,
    top=Inches(2.2),
    size=18,
    left=None,
    width=None,
    pitch=0.68,
    item_height=None,
):
    """
    Render a bullet list.
    Optional left/width keep bullets in a column (e.g. beside a plot).
    pitch is vertical spacing in inches between items.
    """
    left = MARGIN if left is None else left
    width = Inches(11.2) if width is None else width
    item_h = Inches(0.45) if item_height is None else item_height
    text_w = max(width.inches - 0.4, 1.5)
    for i, item in enumerate(items):
        y = top + Inches(i * pitch)
        add_text(slide, left, y, Inches(0.35), item_h, "–", size=size, color=SECONDARY)
        add_text(
            slide,
            left + Inches(0.4),
            y,
            Inches(text_w),
            item_h,
            item,
            size=size,
            color=INK,
        )
