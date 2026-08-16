#!/usr/bin/env python3
"""Merge all ETRA session PPTX decks into one bootcamp file."""

from __future__ import annotations

import copy
import io
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parent.parent
EXPORTS = ROOT / "pdf-exports"
OUT = EXPORTS / "ETRA-AI-ML-Bootcamp.pptx"

DECKS = [
    "Week1-Session1-Foundations.pptx",
    "Week1-Session2-Regression-Models.pptx",
    "Week1-Session3-Classification-Basics.pptx",
    "Week1-Session4-Naive-Bayes-Trees-Evaluation.pptx",
    "Week1-Session5-SVM-Kernel-Methods.pptx",
    "Week2-Session1-Deep-Learning.pptx",
    "Week3-Session1-NLP.pptx",
    "Week4-Session1-GenAI.pptx",
]


def _clear_shapes(slide) -> None:
    sp_tree = slide.shapes._spTree
    for child in list(sp_tree):
        tag = child.tag.rpartition("}")[-1]
        if tag in {"sp", "pic", "grpSp", "graphicFrame", "cxnSp", "contentPart"}:
            sp_tree.remove(child)


def _copy_picture(shape, dest_slide) -> None:
    stream = io.BytesIO(shape.image.blob)
    dest_slide.shapes.add_picture(
        stream,
        shape.left,
        shape.top,
        shape.width,
        shape.height,
    )


def _copy_shape(shape, dest_slide) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        _copy_picture(shape, dest_slide)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _copy_shape(child, dest_slide)
        return
    dest_slide.shapes._spTree.append(copy.deepcopy(shape.element))


def _copy_notes(src_slide, dest_slide) -> None:
    src_notes = src_slide.has_notes_slide and src_slide.notes_slide.notes_text_frame
    if not src_notes:
        return
    text = src_notes.text.strip()
    if not text:
        return
    dest_slide.notes_slide.notes_text_frame.text = text


def append_deck(src_path: Path, dest_prs: Presentation) -> int:
    src = Presentation(src_path)
    blank = dest_prs.slide_layouts[6]
    n = 0
    for src_slide in src.slides:
        dest_slide = dest_prs.slides.add_slide(blank)
        _clear_shapes(dest_slide)
        for shape in src_slide.shapes:
            _copy_shape(shape, dest_slide)
        _copy_notes(src_slide, dest_slide)
        n += 1
    return n


def main() -> None:
    missing = [name for name in DECKS if not (EXPORTS / name).is_file()]
    if missing:
        raise SystemExit(f"Missing decks: {', '.join(missing)}")

    first = EXPORTS / DECKS[0]
    shutil.copyfile(first, OUT)
    dest = Presentation(OUT)
    dest.slide_width = Presentation(first).slide_width
    dest.slide_height = Presentation(first).slide_height

    print(f"Base: {DECKS[0]} ({len(dest.slides)} slides)")
    for name in DECKS[1:]:
        added = append_deck(EXPORTS / name, dest)
        print(f"  + {name} ({added})")
    dest.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(dest.slides)}")


if __name__ == "__main__":
    main()
