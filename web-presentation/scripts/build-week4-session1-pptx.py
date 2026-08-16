#!/usr/bin/env python3
"""
Build Week 4 · Session 1 presentation using official ETRA Design System.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week4-Session1-GenAI.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session-w4s1-diagrams"

SESSION = {
    "eyebrow": "Week 4  ·  Session 1",
    "section_number": "14",
    "section_tag": "S14",
    "section_title": "GenAI",
    "subtitle": "GenAI: BERT, T5, GPT, Attention",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "BERT · T5 · GPT · Attention",
    "topics": [
        "Core Concepts",
        "Transfer Learning",
        "BERT",
        "T5",
        "GPT",
        "Attention",
        "Scaling & Practice",
    ],
}

TOPIC_CONTENT: dict = {
    "Core Concepts": {
        "title": "Generative AI: Core Concepts",
        "kicker": "BERT, T5, GPT, LLMs, and Attention (Focused Edition)",
        "bullets": [
            "This section extracts the most important ideas into a compact teaching flow.",
            "Focus areas: transfer learning, model families, attention mechanics, scaling, and practical deployment.",
            "Designed for bootcamp delivery within roughly 20 slides.",
        ],
    },
    "Transfer Learning": {
        "title": "Transfer Learning in Generative AI",
        "layout": "table",
        "table": {
            "headers": ["Stage", "Data Type", "Goal"],
            "rows": [
                ["Pretraining", "Large mostly unlabeled corpora", "Learn general language representations"],
                ["Fine-tuning", "Task-labeled dataset", "Adapt model to specific downstream objective"],
                ["Inference", "User prompt/context", "Generate or classify outputs for real tasks"],
            ],
        },
        "note": "Transfer learning reduces labeled data needs and accelerates convergence.",
        "extra_slides": [
            {
                "title": "Feature-Based Transfer vs Fine-Tuning",
                "layout": "table",
                "table": {
                    "headers": ["Strategy", "What Changes", "Pros", "Tradeoff"],
                    "rows": [
                        [
                            "Feature-based",
                            "Freeze pretrained backbone; use embeddings as input features",
                            "Fast and lightweight",
                            "Less task adaptation",
                        ],
                        [
                            "Fine-tuning",
                            "Update model weights (often with a task head)",
                            "Higher task performance",
                            "More compute and tuning effort",
                        ],
                    ],
                },
            },
        ],
    },
    "BERT": {
        "title": "BERT: Encoder-Only Bidirectional Understanding",
        "layout": "diagram",
        "plot_path": "transformer-full.png",
        "extra_slides": [
            {
                "title": "BERT: Encoder-Only Bidirectional Understanding",
                "bullets": [
                    "BERT learns contextual token representations using both left and right context.",
                    "Strong for understanding tasks: classification, NER, QA, and retrieval-oriented use cases.",
                    "Pretrained at scale (Wikipedia + BooksCorpus), then adapted to downstream tasks.",
                ],
            },
            {
                "title": "BERT Pretraining Objectives",
                "layout": "table",
                "table": {
                    "headers": ["Objective", "How It Works", "What It Teaches"],
                    "rows": [
                        [
                            "MLM (Masked Language Modeling)",
                            "Mask subset of tokens and predict originals",
                            "Bidirectional contextual semantics",
                        ],
                        [
                            "NSP (Next Sentence Prediction)",
                            "Predict whether sentence B follows sentence A",
                            "Inter-sentence coherence signals",
                        ],
                    ],
                },
                "note": "MLM and NSP were jointly used in original BERT training setup.",
            },
            {
                "title": "BERT Special Tokens: [CLS] and [SEP]",
                "bullets": [
                    "[CLS] is prepended and its final embedding is used as a sequence-level summary.",
                    "[SEP] separates sentence segments and marks boundaries.",
                    "In classification tasks, [CLS] passes through a lightweight prediction head.",
                ],
                "note": "Token design helps unify single- and pair-sentence tasks.",
            },
            {
                "title": "Arabic Adaptation Example: AraBERT",
                "bullets": [
                    "AraBERT follows BERT-base style architecture with Arabic-specific preprocessing.",
                    "Segmentation and Arabic-tailored tokenization improve vocabulary coverage.",
                    "Domain/language adaptation quality can matter more than raw parameter count.",
                ],
            },
            {
                "title": "Arabic Adaptation Example: AraBERT",
                "layout": "table",
                "table": {
                    "headers": ["Component", "Adaptation Benefit"],
                    "rows": [
                        [
                            "Farasa-style segmentation",
                            "Better handling of prefixes/suffixes and morphology",
                        ],
                        [
                            "Arabic-focused vocabulary",
                            "Lower fragmentation and stronger lexical coverage",
                        ],
                        [
                            "Arabic corpus curation",
                            "Improved contextual understanding in target language",
                        ],
                    ],
                },
            },
        ],
    },
    "T5": {
        "title": "T5: Unified Text-to-Text Framework",
        "bullets": [
            "T5 reframes every NLP task as text input -> text output.",
            "Task instruction is included as a prefix (e.g., \"summarize:\", \"translate:\").",
            "One shared architecture/training recipe serves multiple tasks.",
        ],
        "note": "This design strongly influenced modern prompt-based LLM usage.",
        "extra_slides": [
            {
                "title": "T5 Pretraining Objective: Span Corruption",
                "bullets": [
                    "Remove contiguous text spans from input and replace with sentinel tokens.",
                    "Train model to reconstruct missing spans autoregressively.",
                    "Encourages stronger generative behavior than token-only masking.",
                ],
            },
            {
                "title": "T5 Pretraining Objective: Span Corruption",
                "layout": "table",
                "table": {
                    "headers": ["Model", "Corruption Style", "Typical Strength"],
                    "rows": [
                        [
                            "BERT",
                            "Random masked individual tokens",
                            "Language understanding tasks",
                        ],
                        [
                            "T5",
                            "Contiguous span corruption",
                            "Generation + multitask transfer",
                        ],
                    ],
                },
            },
            {
                "title": "T5 Attention Strategies (Encoder/Decoder)",
                "layout": "diagram",
                "plot_path": "encoding-comparison.png",
            },
            {
                "title": "T5 Attention Strategies (Encoder/Decoder)",
                "layout": "table",
                "table": {
                    "headers": ["Attention Type", "Visibility", "Used In"],
                    "rows": [
                        [
                            "Fully visible",
                            "All tokens attend to all tokens",
                            "Encoder representations",
                        ],
                        [
                            "Causal",
                            "Token attends to past only",
                            "Autoregressive generation",
                        ],
                        [
                            "Prefix-causal hybrid",
                            "Full source access + causal target decoding",
                            "Seq2Seq decoding in T5",
                        ],
                    ],
                },
            },
        ],
    },
    "GPT": {
        "title": "GPT: Decoder-Only Autoregressive Generation",
        "layout": "diagram",
        "plot_path": "gpt-architecture.png",
        "extra_slides": [
            {
                "title": "GPT: Decoder-Only Autoregressive Generation",
                "bullets": [
                    "GPT predicts next token from left context only (causal modeling).",
                    "Uses masked self-attention in decoder stack for no-lookahead generation.",
                    "Excellent for open-ended completion, instruction following, and synthesis.",
                ],
            },
            {
                "title": "GPT Evolution and Foundation Model Scale",
                "layout": "table",
                "table": {
                    "headers": ["Generation", "Approx. Parameters", "Key Theme"],
                    "rows": [
                        [
                            "GPT-1",
                            "117M",
                            "Proof of transfer-learning viability",
                        ],
                        [
                            "GPT-2",
                            "1.5B",
                            "Large-scale web pretraining for generation",
                        ],
                        [
                            "GPT-3",
                            "175B",
                            "Few-shot in-context capabilities at scale",
                        ],
                    ],
                },
                "note": "Model utility scales with data, architecture, and training compute quality.",
            },
        ],
    },
    "Attention": {
        "title": "Attention: Why It Was a Breakthrough",
        "layout": "diagram",
        "plot_path": "attention-architecture.png",
        "extra_slides": [
            {
                "title": "Attention: Why It Was a Breakthrough",
                "bullets": [
                    "Attention solves fixed-vector bottlenecks in sequence transduction.",
                    "At each step, model dynamically focuses on most relevant source positions.",
                    "Entire mechanism is differentiable and learned end-to-end.",
                ],
            },
            {
                "title": "Attention Computation Pipeline",
                "layout": "formula_example",
                "formula": "Attention(Q,K,V)=softmax(QK⊤ / √d_k)V",
                "formula_tex": r"\mathrm{Attention}(Q,K,V)=\mathrm{softmax}\left(\frac{QK^{\top}}{\sqrt{d_k}}\right)V",
                "bullets": [
                    "Compute relevance scores between query state and source states.",
                    "Apply softmax to obtain attention weights (probability distribution).",
                    "Return weighted sum of value/source states as contextual output.",
                ],
            },
            {
                "title": "Bahdanau vs Luong Attention (Classic RNN Era)",
                "layout": "table",
                "table": {
                    "headers": ["Variant", "Scoring Function", "Typical Placement"],
                    "rows": [
                        [
                            "Luong",
                            "Bilinear/dot-style",
                            "Attention after decoder state update",
                        ],
                        [
                            "Bahdanau",
                            "MLP/additive score",
                            "Attention integrated before decoder update",
                        ],
                    ],
                },
                "note": "Both approximate alignment; design choice affects speed and accuracy tradeoffs.",
            },
            {
                "title": "Self-Attention, Q/K/V, and Multi-Head",
                "layout": "diagram",
                "plot_path": "attention-heatmap.png",
            },
            {
                "title": "Self-Attention, Q/K/V, and Multi-Head",
                "bullets": [
                    "Query asks for relevant context; Key indexes available context; Value carries content.",
                    "Self-attention lets each token aggregate information from other tokens in sequence.",
                    "Multi-head attention captures different linguistic relations in parallel.",
                ],
                "note": "Different heads can specialize in syntax, agreement, locality, or semantics.",
            },
            {
                "title": "Masked Self-Attention in Decoders",
                "layout": "diagram",
                "plot_path": "transformer-block.png",
            },
            {
                "title": "Masked Self-Attention in Decoders",
                "bullets": [
                    "Decoder cannot access future tokens during generation.",
                    "Causal mask enforces left-to-right consistency and prevents information leakage.",
                    "Enables parallel training over full target sequence while preserving autoregressive objective.",
                ],
            },
            {
                "title": "Long-Sequence Challenge in Transformers",
                "layout": "diagram",
                "plot_path": "transformer-block.png",
            },
            {
                "title": "Long-Sequence Challenge in Transformers",
                "layout": "table",
                "table": {
                    "headers": ["Bottleneck", "Complexity", "Practical Impact"],
                    "rows": [
                        [
                            "Attention map computation",
                            "O(L²)",
                            "Time/memory explode as context grows",
                        ],
                        [
                            "Activation storage (training)",
                            "O(N*L*d_model)",
                            "High VRAM demand for deep long-context models",
                        ],
                    ],
                },
                "note": "Long-context efficiency is now a central LLM engineering topic.",
            },
        ],
    },
    "Scaling & Practice": {
        "title": "Scaling Laws and Compute-Optimal Training",
        "bullets": [
            "Performance improves with data, model size, and compute, but with diminishing returns.",
            "Compute-optimal training balances parameter count and token budget.",
            "Chinchilla-style insight: many large models are under-trained relative to their size.",
        ],
        "extra_slides": [
            {
                "title": "Scaling Laws and Compute-Optimal Training",
                "layout": "table",
                "table": {
                    "headers": ["Principle", "Implication"],
                    "rows": [
                        [
                            "Fixed compute budget",
                            "Tune model size and data volume jointly",
                        ],
                        [
                            "Data-quality filtering",
                            "Small fraction of tokens may dominate useful learning",
                        ],
                        [
                            "Chinchilla ratio (rule of thumb)",
                            "Roughly ~20 training tokens per parameter",
                        ],
                    ],
                },
            },
            {
                "title": "Generative AI Practical Takeaways",
                "bullets": [
                    "Choose architecture by task type: encoder-only, encoder-decoder, or decoder-only.",
                    "Prefer transfer learning and fine-tuning over training from scratch when possible.",
                    "Use attention-aware designs and decoding strategy (greedy/beam) based on product constraints.",
                    "Plan for memory/compute early: quantization, context limits, and deployment targets.",
                    "For domain-heavy use cases, adapt with specialized corpora and evaluation protocols.",
                ],
                "note": "Strong GenAI systems come from balanced choices across data, architecture, compute, and evaluation.",
            },
        ],
    },
}

TAKEAWAYS = [
    "Choose architecture by task: encoder-only (BERT), encoder-decoder (T5), or decoder-only (GPT).",
    "Prefer transfer learning and fine-tuning over training from scratch.",
    "Attention, decoding strategy, and compute limits shape what you can ship.",
]

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "Week 4 Session 1 — GenAI: BERT, T5, GPT, Attention.",
    "current": "S14",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.05), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.5),
        Inches(11.5),
        Inches(0.85),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.5), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(3.85), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    topics = s["topics"]
    n_cols = 3
    chip_w = Inches(3.7)
    chip_h = Inches(0.48)
    gap_x = Inches(0.22)
    start_y = Inches(4.7)
    for i, topic in enumerate(topics):
        row, col = divmod(i, n_cols)
        x = MARGIN + col * (chip_w + gap_x)
        y = start_y + Inches(row * 0.68)
        soft_card(slide, x, y, chip_w, chip_h, fill=SOFT)
        add_text(slide, x, y + Inches(0.06), chip_w, Inches(0.36), topic, size=12, color=PRIMARY, align=PP_ALIGN.CENTER)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def content_title(slide, title: str, subtitle: str | None = None):
    """Title that shrinks on long headlines so it does not cover the kicker."""
    long = len(title or "") > 46
    size = 24 if long else 32
    title_h = 0.80 if long else 0.62
    add_text(
        slide,
        MARGIN,
        Inches(1.12),
        Inches(11.2),
        Inches(title_h),
        title,
        size=size,
        bold=True,
        color=PRIMARY,
    )
    if subtitle:
        add_text(
            slide,
            MARGIN,
            Inches(1.12 + title_h),
            Inches(11.2),
            Inches(0.34),
            subtitle,
            size=14 if long else 15,
            color=MUTED,
        )


def content_top(title: str, kicker: str | None = None) -> float:
    long = len(title or "") > 46
    title_h = 0.80 if long else 0.62
    kicker_h = 0.38 if kicker else 0.0
    return 1.12 + title_h + kicker_h + 0.16


def _chars_per_line(width_in: float, size: float) -> int:
    return max(22, int(width_in * (108 / max(size, 10))))


def _line_count(text: str, width_in: float, size: float) -> int:
    cpl = _chars_per_line(width_in, size)
    words = (text or "").split()
    if not words:
        return 1
    lines, cur = 1, 0
    for word in words:
        need = len(word) + (1 if cur else 0)
        if cur + need > cpl:
            lines += 1
            cur = len(word)
        else:
            cur += need
    return lines


def _bullet_box(items: list[str], width_in: float, size: float) -> tuple[float, float]:
    lines = max((_line_count(t, width_in, size) for t in items), default=1)
    item_h = max(0.42, min(1.15, lines * (size / 72) * 1.35 + 0.10))
    return item_h, item_h + 0.10


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
        if px_w <= 0 or px_h <= 0:
            return None
        aspect = px_w / px_h
        src = path
        if aspect < 0.5:
            from io import BytesIO

            crop_h = max(1, min(px_h, int(px_w / 2.4)))
            cropped = im.crop((0, 0, px_w, crop_h))
            buf = BytesIO()
            cropped.save(buf, format="PNG")
            buf.seek(0)
            src = buf
            px_w, px_h = cropped.size
            aspect = px_w / max(px_h, 1)
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(src) if src is path else src,
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    bsize = 14 if dense else 16

    if content.get("body"):
        top = content_top(content["title"], content.get("kicker"))
        soft_card(slide, MARGIN, Inches(top), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(top + 0.20),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(top + 1.35)
        bsize = 14 if dense else 15
    else:
        bullet_top = Inches(content_top(content["title"], content.get("kicker")))
        bsize = 14 if dense else 15

    item_h, pitch = _bullet_box(items, 11.2, bsize)
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(item_h),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    has_body = bool((content.get("body") or "").strip())
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)

    y_after = content_top(content["title"], content.get("kicker"))
    if has_body:
        soft_card(slide, MARGIN, Inches(y_after), col_w, Inches(1.2), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(y_after + 0.15),
            text_w,
            Inches(0.95),
            content.get("body") or "",
            size=13,
            color=INK,
        )
        y_after = y_after + 1.35

    if has_formula:
        formula_card_h = 1.65 if is_frac else 1.35
        formula_box_h = 0.95 if is_frac else 0.65
        soft_card(slide, MARGIN, Inches(y_after), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(y_after + 0.10),
            text_w,
            Inches(0.22),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(y_after + 0.35),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        if content.get("formula_note"):
            add_text(
                slide,
                Inches(MARGIN.inches + 0.3),
                Inches(y_after + 0.35 + formula_box_h + 0.02),
                text_w,
                Inches(0.28),
                content["formula_note"],
                size=10,
                color=MUTED,
            )
        bullet_top = y_after + formula_card_h + 0.12
        bsize = 12
    else:
        bullet_top = y_after
        bsize = 13

    item_h, pitch = _bullet_box(items, bullet_w.inches, bsize)
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch)) if items else 0
    if has_formula and has_plot:
        max_items = min(max_items, 2)
    if items:
        bullets(
            slide,
            items[:max_items],
            top=Inches(bullet_top),
            size=bsize,
            left=MARGIN,
            width=bullet_w,
            pitch=pitch,
            item_height=Inches(item_h),
        )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(content_top(content["title"], content.get("kicker")))
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(body_top.inches + 0.13),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(body_top.inches + 0.85)
    else:
        formula_top = body_top

    formula_h = 1.55 if is_frac else 1.25
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.12),
        Inches(11.2),
        Inches(0.95 if is_frac else 0.7),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )
    if content.get("formula_note"):
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(formula_top.inches + formula_h - 0.35),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    bullet_top = formula_top.inches + formula_h + 0.12
    bsize = 13 if dense else 15
    item_h, pitch = _bullet_box(items, 11.2, bsize)
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch)) if items else 0
    if has_note:
        max_items = min(max_items, 3)
    if items:
        bullets(
            slide,
            items[:max_items],
            top=Inches(bullet_top),
            size=bsize,
            pitch=pitch,
            width=Inches(11.2),
            item_height=Inches(item_h),
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = content_top(content["title"], content.get("kicker"))
    if show_body:
        soft_card(slide, MARGIN, Inches(table_top), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(table_top + 0.10),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = table_top + 0.70

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    note_top = table_top + table_h + 0.08
    render_note = has_note and (note_top + 0.45) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    kicker = content.get("kicker")
    content_title(slide, content["title"], kicker)

    diagram_top = Inches(content_top(content["title"], kicker))
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(diagram_top.inches - 0.05),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(diagram_top.inches + 0.30)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    footer_y = 7.05
    has_note = bool(content.get("note"))
    note_band = 0.40 if has_note else 0.10
    card_h = max(2.6, footer_y - note_band - diagram_top.inches)
    pad = 0.12
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(card_h), fill=SOFT)
    pic = _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + pad),
        Inches(11.6),
        Inches(card_h - 2 * pad),
        folder=folder,
    )
    if pic is not None:
        extra = (card_h - 2 * pad) - pic.height.inches
        if extra > 0.04:
            pic.top = Inches(diagram_top.inches + pad + extra / 2)

    if has_note:
        add_text(
            slide,
            MARGIN,
            Inches(footer_y - 0.36),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def slide_takeaways(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Takeaways", f"{index:02d}")
    title_block(slide, "3 Ideas to Keep", "Leave with these")
    for i, text in enumerate(TAKEAWAYS):
        y = Inches(2.35) + Inches(i * 1.25)
        soft_card(slide, MARGIN, y, Inches(12.0), Inches(1.08), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            y + Inches(0.28),
            Inches(0.5),
            Inches(0.5),
            str(i + 1),
            size=22,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            MARGIN + Inches(0.95),
            y + Inches(0.28),
            Inches(10.6),
            Inches(0.6),
            text,
            size=16,
            color=INK,
        )
    content_footer(slide, index, total)


def slide_close(prs, total, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))
    add_text(
        slide,
        MARGIN,
        Inches(2.9),
        Inches(12),
        Inches(0.8),
        "Thank you",
        size=40,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN,
        Inches(3.8),
        Inches(12),
        Inches(0.4),
        "Labs: code/17- GenAI",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    content_footer(slide, index, total)


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 6 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    n += 1
    slide_takeaways(prs, total, n)
    n += 1
    slide_close(prs, total, n)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
