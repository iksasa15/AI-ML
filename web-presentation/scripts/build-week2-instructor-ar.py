#!/usr/bin/env python3
"""
Week 2 · Session 1 instructor deck (Arabic say-this + clarify).

Student slides stay in Week2-Session1-Deep-Learning.pptx.
This file is a separate RTL guide: one instructor slide per student slide.
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    FONT,
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    add_text,
    content_footer,
    content_header,
    paint_light,
    right_rail,
    set_run,
    soft_card,
)
from week2_instructor_scripts import SCRIPTS  # noqa: E402

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week2-Session1-Deep-Learning-Instructor-AR.pptx"
STUDENT_BUILDER = Path(__file__).resolve().parent / "build-week2-session1-pptx.py"

# DIN Next LT W23 is the Arabic cut (name ID 16). Arial if a machine lacks it.
AR_FONT = FONT
AR_FALLBACK = "Arial"


def _load_student():
    spec = importlib.util.spec_from_file_location("w2_student", STUDENT_BUILDER)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def _kind(content: dict) -> str:
    layout = content.get("layout")
    if layout == "diagram" or (
        content.get("plot_path")
        and not content.get("bullets")
        and not content.get("formula")
        and not content.get("table")
    ):
        return "diagram"
    if layout == "table" or (
        content.get("table") and layout not in ("formula_example", "diagram")
    ):
        return "table"
    if layout == "formula_example" or content.get("formula"):
        return "formula"
    if content.get("plot_path"):
        return "mixed"
    return "bullets"


def student_catalog(module) -> list[dict]:
    """Slide order/titles matching build-week2-session1-pptx.py main()."""
    s = module.SESSION
    rows = [
        {
            "title": s["section_title"],
            "kicker": s["eyebrow"],
            "kind": "chrome",
        },
        {
            "title": module.BIG_PICTURE["title"],
            "kicker": module.BIG_PICTURE["focus"],
            "kind": "chrome",
        },
        {
            "title": s["section_title"],
            "kicker": s["focus"],
            "kind": "chrome",
        },
        {
            "title": s["focus"],
            "kicker": "What we will cover in this session",
            "kind": "chrome",
        },
    ]
    for topic in s["topics"]:
        content = module.TOPIC_CONTENT[topic]
        rows.append(
            {
                "title": content["title"],
                "kicker": content.get("kicker") or "",
                "kind": _kind(content),
            }
        )
        for extra in content.get("extra_slides") or []:
            rows.append(
                {
                    "title": extra["title"],
                    "kicker": extra.get("kicker") or "",
                    "kind": _kind(extra),
                }
            )
    rows.append(
        {
            "title": "3 Ideas to Keep",
            "kicker": "Leave with these",
            "kind": "chrome",
        }
    )
    rows.append(
        {
            "title": "Thank you",
            "kicker": "Labs: code/15- Deep Learning",
            "kind": "chrome",
        }
    )
    return rows


def set_paragraph_rtl(paragraph) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")
    pPr.set("algn", "r")


def add_rtl_text(
    slide,
    left,
    top,
    width,
    height,
    value,
    *,
    size=16,
    bold=False,
    color=INK,
    font=None,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[
                anchor
            ],
        )
    except Exception:
        pass
    chunks = [part.strip() for part in (value or "").split("\n") if part.strip()]
    if not chunks:
        chunks = [""]
    face = font or AR_FONT
    for i, chunk in enumerate(chunks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        set_paragraph_rtl(p)
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = chunk
        set_run(run, size, bold=bold, color=color, font=face)
        rPr = run._r.get_or_add_rPr()
        rPr.set("rtl", "1")
    return box


def rtl_card(slide, left, top, width, height, label: str, body: str, *, fill=SOFT):
    soft_card(slide, left, top, width, height, fill=fill)
    pad = Inches(0.28)
    add_rtl_text(
        slide,
        left + pad,
        top + Inches(0.12),
        width - pad * 2,
        Inches(0.36),
        label,
        size=14,
        bold=True,
        color=PRIMARY,
    )
    add_rtl_text(
        slide,
        left + pad,
        top + Inches(0.48),
        width - pad * 2,
        height - Inches(0.62),
        body,
        size=15,
        color=INK,
    )


def instructor_header(slide, page: int, total: int) -> None:
    content_header(slide, "S07 · دليل المحاضر", f"{page:02d}")
    add_text(
        slide,
        Inches(8.9),
        Inches(0.48),
        Inches(2.55),
        Inches(0.3),
        f"{page} / {total}",
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def instructor_title(slide, title: str, kicker: str | None) -> float:
    long = len(title or "") > 46
    size = 24 if long else 30
    title_h = 0.78 if long else 0.55
    add_text(
        slide,
        MARGIN,
        Inches(1.08),
        Inches(12.0),
        Inches(title_h),
        title,
        size=size,
        bold=True,
        color=PRIMARY,
    )
    top = 1.08 + title_h
    if kicker:
        add_text(
            slide,
            MARGIN,
            Inches(top),
            Inches(12.0),
            Inches(0.32),
            kicker,
            size=13,
            color=MUTED,
        )
        top += 0.32
    return top + 0.10


def slide_instructor(prs, page: int, total: int, meta: dict, script: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    instructor_header(slide, page, total)
    cards_top = instructor_title(slide, meta["title"], meta.get("kicker") or None)

    footer_y = 7.05
    available = footer_y - cards_top - 0.08
    gap = 0.14
    say_h = available * 0.58
    clarify_h = available - say_h - gap
    card_w = 12.0

    rtl_card(
        slide,
        MARGIN,
        Inches(cards_top),
        Inches(card_w),
        Inches(say_h),
        "قل للطلاب:",
        script["say"],
        fill=SOFT,
    )
    rtl_card(
        slide,
        MARGIN,
        Inches(cards_top + say_h + gap),
        Inches(card_w),
        Inches(clarify_h),
        "توضيح:",
        script["clarify"],
        fill=SOFT_2,
    )
    content_footer(slide, page, total)


def main() -> None:
    student = _load_student()
    catalog = student_catalog(student)
    if len(catalog) != len(SCRIPTS):
        raise SystemExit(
            f"Script count {len(SCRIPTS)} does not match student slides {len(catalog)}"
        )
    for i, (meta, script) in enumerate(zip(catalog, SCRIPTS), start=1):
        expected = script.get("title")
        if expected and expected != meta["title"]:
            raise SystemExit(
                f"Slide {i}: script title {expected!r} != student title {meta['title']!r}"
            )

    total = len(catalog)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for i, (meta, script) in enumerate(zip(catalog, SCRIPTS), start=1):
        slide_instructor(prs, i, total, meta, script)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Instructor RTL guide — student deck unchanged")


if __name__ == "__main__":
    main()
