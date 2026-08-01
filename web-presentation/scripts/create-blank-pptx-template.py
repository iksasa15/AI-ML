#!/usr/bin/env python3
"""
Generate a polished empty ETRA-branded PowerPoint template (16:9).

Layouts:
  0. Guide
  1. Title / Hero
  2. Section divider (dark)
  3. Content bullets
  4. Two columns
  5. Title + body / note
  6. Three cards
  7. Quote / takeaway
  8. Blank canvas

Usage:
  python3 scripts/create-blank-pptx-template.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
LOGO = ROOT / "public" / "assets" / "etra-logo.png"
OUT = ROOT / "pdf-exports" / "ETRA-Presentation-Blank-Template.pptx"

# Design tokens — Academic Precision / ETRA
BG = RGBColor(0xF8, 0xF7, 0xFC)
BG_SUBTLE = RGBColor(0xF0, 0xEE, 0xF8)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x52, 0x34, 0xB7)
ACCENT_2 = RGBColor(0x9E, 0x59, 0xCD)
TEXT = RGBColor(0x1A, 0x15, 0x33)
MUTED = RGBColor(0x6B, 0x64, 0x90)
DIVIDER_BG = RGBColor(0x0D, 0x0D, 0x1A)
DIVIDER_SURFACE = RGBColor(0x1A, 0x1A, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_PURPLE = RGBColor(0xED, 0xE8, 0xFA)
PLACEHOLDER_OK = RGBColor(0xB0, 0xA8, 0xD4)
CARD_EDGE = RGBColor(0xE4, 0xDE, 0xF4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.7)


def _set_run(run, *, size_pt: float, bold: bool = False, color=TEXT, font="Helvetica"):
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", font)


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size_pt=size, bold=bold, color=color)
    return box


def _round_rect(slide, left, top, width, height, fill, *, line=None, line_w=1.0, adj=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    try:
        shape.adjustments[0] = adj
    except Exception:
        pass
    return shape


def _rect(slide, left, top, width, height, fill, *, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def _ellipse(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _set_gradient(shape, c1: RGBColor, c2: RGBColor, angle: float = 135.0):
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    stops = fill.gradient_stops
    stops[0].color.rgb = c1
    stops[1].color.rgb = c2


def _paint_bg(slide, *, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DIVIDER_BG if dark else BG


def _atmosphere(slide):
    """Soft purple blobs for light slides — subtle depth without clutter."""
    a = _ellipse(slide, Inches(10.2), Inches(-1.4), Inches(5.2), Inches(5.2), SOFT_PURPLE)
    b = _ellipse(slide, Inches(-1.8), Inches(4.8), Inches(4.5), Inches(4.5), BG_SUBTLE)
    # Keep blobs behind content by drawing them first (caller order).
    return a, b


def _add_header(slide, section_tag="SECTION", section_label="Topic title", slide_num="1 / N"):
    _rect(slide, 0, 0, SLIDE_W, Inches(0.78), SURFACE)
    # accent underline under header
    accent_line = _rect(slide, 0, Inches(0.78), SLIDE_W, Inches(0.035), ACCENT)

    # tag pill
    pill = _round_rect(
        slide,
        MARGIN_X,
        Inches(0.2),
        Inches(1.55),
        Inches(0.38),
        SOFT_PURPLE,
        adj=0.5,
    )
    _add_textbox(
        slide,
        MARGIN_X,
        Inches(0.22),
        Inches(1.55),
        Inches(0.35),
        section_tag,
        size=11,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        Inches(2.45),
        Inches(0.22),
        Inches(7.5),
        Inches(0.35),
        section_label,
        size=13,
        color=MUTED,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    _add_textbox(
        slide,
        Inches(10.4),
        Inches(0.22),
        Inches(2.2),
        Inches(0.35),
        slide_num,
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return accent_line


def _add_footer(slide, progress=0.25, *, dark=False):
    track_color = RGBColor(0x2A, 0x2A, 0x4A) if dark else CARD_EDGE
    fill_color = ACCENT_2 if dark else ACCENT
    _round_rect(
        slide,
        MARGIN_X,
        Inches(7.18),
        Inches(11.9),
        Inches(0.1),
        track_color,
        adj=0.5,
    )
    fill_w = max(Inches(0.35), int(Inches(11.9) * min(max(progress, 0.02), 1.0)))
    _round_rect(
        slide,
        MARGIN_X,
        Inches(7.18),
        fill_w,
        Inches(0.1),
        fill_color,
        adj=0.5,
    )


def _add_logo(slide, *, corner="content"):
    if not LOGO.is_file():
        return
    if corner == "hero":
        slide.shapes.add_picture(str(LOGO), Inches(11.05), Inches(0.4), height=Inches(1.05))
    elif corner == "divider":
        slide.shapes.add_picture(str(LOGO), Inches(11.15), Inches(0.4), height=Inches(0.95))
    else:
        slide.shapes.add_picture(str(LOGO), Inches(11.35), Inches(1.05), height=Inches(0.72))


def _title_block(slide, title: str, subtitle: str | None = None, *, top=Inches(1.05)):
    _add_textbox(slide, MARGIN_X, top, Inches(10.4), Inches(0.65), title, size=30, bold=True)
    # title rule — matches web .slide-title-rule
    rule = _round_rect(
        slide,
        MARGIN_X,
        top + Inches(0.68),
        Inches(1.35),
        Inches(0.07),
        ACCENT,
        adj=0.5,
    )
    if subtitle:
        _add_textbox(
            slide,
            MARGIN_X,
            top + Inches(0.88),
            Inches(10.4),
            Inches(0.4),
            subtitle,
            size=15,
            color=MUTED,
        )
    return rule


def _content_card(slide, left, top, width, height, *, title=None, body=None, accent_bar=True):
    # soft shadow layer
    _round_rect(
        slide,
        left + Inches(0.04),
        top + Inches(0.05),
        width,
        height,
        BG_SUBTLE,
        adj=0.07,
    )
    card = _round_rect(slide, left, top, width, height, SURFACE, line=CARD_EDGE, line_w=1.1, adj=0.07)
    if accent_bar:
        bar = _rect(slide, left, top, Inches(0.09), height, ACCENT)
        # clip look: overlay rounded corners approx by covering with tiny fills — skip, bar is fine
        _ = bar
    if title:
        _add_textbox(
            slide,
            left + Inches(0.35),
            top + Inches(0.28),
            width - Inches(0.55),
            Inches(0.4),
            title,
            size=16,
            bold=True,
            color=ACCENT,
        )
    if body:
        y = top + (Inches(0.75) if title else height / 2 - Inches(0.2))
        _add_textbox(
            slide,
            left + Inches(0.35),
            y,
            width - Inches(0.55),
            height - (Inches(1.0) if title else Inches(0.5)),
            body,
            size=15,
            color=PLACEHOLDER_OK if not title else MUTED,
            align=PP_ALIGN.LEFT if title else PP_ALIGN.CENTER,
        )
    return card


def _bullet_row(slide, left, top, width, text, *, index=1):
    # card row
    _round_rect(
        slide,
        left,
        top,
        width,
        Inches(0.78),
        SURFACE,
        line=CARD_EDGE,
        line_w=1.0,
        adj=0.12,
    )
    # number circle
    circle = _ellipse(slide, left + Inches(0.18), top + Inches(0.16), Inches(0.46), Inches(0.46), SOFT_PURPLE)
    _ = circle
    _add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(0.2),
        Inches(0.46),
        Inches(0.4),
        str(index),
        size=13,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        left + Inches(0.8),
        top + Inches(0.2),
        width - Inches(1.1),
        Inches(0.42),
        text,
        size=17,
        color=TEXT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# ─── Layouts ───────────────────────────────────────────────────────────────


def slide_guide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "TEMPLATE", "How to use", "01")
    _add_logo(slide)
    _title_block(slide, "ETRA presentation template", "Empty layouts · brand-ready · 16:9")

    tips = [
        ("01", "Duplicate a layout slide, then replace the placeholder copy."),
        ("02", "Keep the header tag + topic label; update the slide number."),
        ("03", "Accent #5234B7 — use for rules, tags, and emphasis only."),
        ("04", "Prefer short titles and one idea per bullet or card."),
    ]
    for i, (num, tip) in enumerate(tips):
        y = Inches(2.55) + Inches(i * 0.95)
        _round_rect(
            slide,
            MARGIN_X,
            y,
            Inches(11.9),
            Inches(0.82),
            SURFACE,
            line=CARD_EDGE,
            adj=0.12,
        )
        _ellipse(slide, MARGIN_X + Inches(0.22), y + Inches(0.18), Inches(0.46), Inches(0.46), SOFT_PURPLE)
        _add_textbox(
            slide,
            MARGIN_X + Inches(0.22),
            y + Inches(0.22),
            Inches(0.46),
            Inches(0.4),
            num,
            size=12,
            bold=True,
            color=ACCENT,
            align=PP_ALIGN.CENTER,
        )
        _add_textbox(
            slide,
            MARGIN_X + Inches(0.9),
            y + Inches(0.22),
            Inches(10.5),
            Inches(0.42),
            tip,
            size=16,
            color=TEXT,
        )
    _add_footer(slide, 0.08)


def slide_title(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)

    # Full left gradient panel
    panel = _rect(slide, 0, 0, Inches(0.22), SLIDE_H, ACCENT)
    _set_gradient(panel, ACCENT, ACCENT_2, 90)

    # Atmospheric shapes
    _ellipse(slide, Inches(9.5), Inches(-2.0), Inches(6.5), Inches(6.5), SOFT_PURPLE)
    _ellipse(slide, Inches(11.2), Inches(4.2), Inches(3.8), Inches(3.8), BG_SUBTLE)

    # Bottom brand strip
    strip = _rect(slide, 0, Inches(6.55), SLIDE_W, Inches(0.95), SURFACE)
    _ = strip
    _rect(slide, 0, Inches(6.55), SLIDE_W, Inches(0.04), ACCENT)

    _add_logo(slide, corner="hero")

    _round_rect(
        slide,
        Inches(0.9),
        Inches(1.85),
        Inches(2.4),
        Inches(0.38),
        SOFT_PURPLE,
        adj=0.5,
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(1.87),
        Inches(2.4),
        Inches(0.35),
        "ETRA  ·  TRAINING",
        size=11,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.5),
        Inches(10.8),
        Inches(1.15),
        "Presentation title",
        size=46,
        bold=True,
        color=TEXT,
    )
    _round_rect(slide, Inches(0.9), Inches(3.7), Inches(1.6), Inches(0.08), ACCENT, adj=0.5)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(4.0),
        Inches(9.5),
        Inches(0.55),
        "Subtitle — session focus in one clear line",
        size=20,
        color=MUTED,
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(6.78),
        Inches(8.5),
        Inches(0.4),
        "Trainer name   ·   Date   ·   Organization",
        size=13,
        color=MUTED,
    )
    _add_footer(slide, 0.05)


def slide_section_divider(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide, dark=True)

    # Glow orbs
    orb1 = _ellipse(slide, Inches(-1.2), Inches(-1.5), Inches(5.5), Inches(5.5), RGBColor(0x1E, 0x14, 0x45))
    orb2 = _ellipse(slide, Inches(9.8), Inches(3.5), Inches(5.0), Inches(5.0), RGBColor(0x22, 0x16, 0x48))
    _ = orb1, orb2

    # Large faint section number
    _add_textbox(
        slide,
        Inches(0.5),
        Inches(1.4),
        Inches(4),
        Inches(2.2),
        "01",
        size=120,
        bold=True,
        color=RGBColor(0x2A, 0x22, 0x55),
    )

    _add_logo(slide, corner="divider")

    _round_rect(
        slide,
        Inches(0.9),
        Inches(2.35),
        Inches(3.1),
        Inches(0.38),
        RGBColor(0x2A, 0x1F, 0x55),
        adj=0.5,
    )
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(2.37),
        Inches(3.1),
        Inches(0.35),
        "WEEK X  ·  SESSION Y",
        size=11,
        bold=True,
        color=ACCENT_2,
        align=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide,
        Inches(0.9),
        Inches(3.0),
        Inches(11),
        Inches(1.1),
        "Section title",
        size=42,
        bold=True,
        color=WHITE,
    )
    _round_rect(slide, Inches(0.9), Inches(4.2), Inches(1.5), Inches(0.07), ACCENT_2, adj=0.5)
    _add_textbox(
        slide,
        Inches(0.9),
        Inches(4.5),
        Inches(9.5),
        Inches(0.55),
        "One sentence describing what this section covers",
        size=17,
        color=PLACEHOLDER_OK,
    )

    # topic chips
    chips = ["Topic A", "Topic B", "Topic C"]
    x = Inches(0.9)
    for chip in chips:
        w = Inches(1.55)
        _round_rect(slide, x, Inches(5.5), w, Inches(0.4), DIVIDER_SURFACE, line=RGBColor(0x3A, 0x30, 0x68), adj=0.5)
        _add_textbox(slide, x, Inches(5.52), w, Inches(0.36), chip, size=12, color=PLACEHOLDER_OK, align=PP_ALIGN.CENTER)
        x += Inches(1.75)

    _add_footer(slide, 0.22, dark=True)


def slide_content_bullets(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "SECTION", "Topic title", "03")
    _add_logo(slide)
    _title_block(slide, "Slide title", "Optional supporting line")

    bullets = [
        "First key point — replace with your content",
        "Second key point — keep one idea per line",
        "Third key point — short and concrete",
        "Fourth key point — optional detail",
    ]
    for i, text in enumerate(bullets):
        _bullet_row(slide, MARGIN_X, Inches(2.55) + Inches(i * 0.95), Inches(11.9), text, index=i + 1)
    _add_footer(slide, 0.38)


def slide_two_columns(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "SECTION", "Topic title", "04")
    _add_logo(slide)
    _title_block(slide, "Compare two ideas", "Side-by-side structure")

    _content_card(
        slide,
        MARGIN_X,
        Inches(2.45),
        Inches(5.75),
        Inches(4.2),
        title="Left idea",
        body="Add definition, steps, or bullet points here.",
    )
    _content_card(
        slide,
        Inches(6.85),
        Inches(2.45),
        Inches(5.75),
        Inches(4.2),
        title="Right idea",
        body="Mirror structure for a clean comparison.",
    )
    _add_footer(slide, 0.5)


def slide_title_body(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "SECTION", "Topic title", "05")
    _add_logo(slide)
    _title_block(slide, "Concept title", "Definition or explanation")

    # Main body card
    _content_card(
        slide,
        MARGIN_X,
        Inches(2.45),
        Inches(11.9),
        Inches(2.7),
        title="Core explanation",
        body="Write 2–4 sentences. Keep the language concrete and trainee-friendly.",
    )
    # Note strip (matches web note-box)
    note = _round_rect(
        slide,
        MARGIN_X,
        Inches(5.45),
        Inches(11.9),
        Inches(1.15),
        SOFT_PURPLE,
        adj=0.1,
    )
    _ = note
    _rect(slide, MARGIN_X, Inches(5.45), Inches(0.1), Inches(1.15), ACCENT)
    _add_textbox(
        slide,
        MARGIN_X + Inches(0.35),
        Inches(5.6),
        Inches(11.2),
        Inches(0.3),
        "KEY TAKEAWAY",
        size=11,
        bold=True,
        color=ACCENT,
    )
    _add_textbox(
        slide,
        MARGIN_X + Inches(0.35),
        Inches(5.95),
        Inches(11.2),
        Inches(0.4),
        "One memorable sentence the trainee should leave with.",
        size=15,
        color=TEXT,
    )
    _add_footer(slide, 0.62)


def slide_three_cards(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "SECTION", "Topic title", "06")
    _add_logo(slide)
    _title_block(slide, "Three pillars", "Equal-weight cards")

    labels = [("01", "Pillar one"), ("02", "Pillar two"), ("03", "Pillar three")]
    card_w = Inches(3.75)
    gap = Inches(0.3)
    start = MARGIN_X
    for i, (num, title) in enumerate(labels):
        x = start + i * (card_w + gap)
        _round_rect(
            slide,
            x + Inches(0.04),
            Inches(2.55),
            card_w,
            Inches(4.05),
            BG_SUBTLE,
            adj=0.08,
        )
        card = _round_rect(slide, x, Inches(2.5), card_w, Inches(4.05), SURFACE, line=CARD_EDGE, adj=0.08)
        # top accent gradient bar
        topbar = _rect(slide, x, Inches(2.5), card_w, Inches(0.12), ACCENT)
        if i == 1:
            _set_gradient(topbar, ACCENT, ACCENT_2, 0)
        elif i == 2:
            topbar.fill.solid()
            topbar.fill.fore_color.rgb = ACCENT_2
        _add_textbox(
            slide,
            x + Inches(0.3),
            Inches(2.9),
            Inches(3.1),
            Inches(0.4),
            num,
            size=22,
            bold=True,
            color=ACCENT,
        )
        _add_textbox(
            slide,
            x + Inches(0.3),
            Inches(3.45),
            Inches(3.1),
            Inches(0.45),
            title,
            size=18,
            bold=True,
            color=TEXT,
        )
        _add_textbox(
            slide,
            x + Inches(0.3),
            Inches(4.1),
            Inches(3.1),
            Inches(1.8),
            "Short supporting copy for this pillar.",
            size=14,
            color=MUTED,
        )
        _ = card
    _add_footer(slide, 0.74)


def slide_quote(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "SECTION", "Takeaway", "07")
    _add_logo(slide)

    panel = _round_rect(
        slide,
        Inches(1.4),
        Inches(1.7),
        Inches(10.5),
        Inches(4.6),
        SURFACE,
        line=CARD_EDGE,
        adj=0.08,
    )
    _ = panel
    accent = _rect(slide, Inches(1.4), Inches(1.7), Inches(0.14), Inches(4.6), ACCENT)
    _set_gradient(accent, ACCENT, ACCENT_2, 90)

    _add_textbox(
        slide,
        Inches(2.0),
        Inches(2.2),
        Inches(9.2),
        Inches(0.4),
        "REMEMBER",
        size=12,
        bold=True,
        color=ACCENT,
    )
    _add_textbox(
        slide,
        Inches(2.0),
        Inches(2.9),
        Inches(9.2),
        Inches(2.2),
        "“Write the one sentence you want every trainee to repeat.”",
        size=28,
        bold=True,
        color=TEXT,
    )
    _round_rect(slide, Inches(2.0), Inches(5.3), Inches(1.2), Inches(0.06), ACCENT_2, adj=0.5)
    _add_textbox(
        slide,
        Inches(2.0),
        Inches(5.55),
        Inches(9.0),
        Inches(0.35),
        "Optional attribution or session cue",
        size=14,
        color=MUTED,
    )
    _add_footer(slide, 0.86)


def slide_blank_canvas(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(slide)
    _atmosphere(slide)
    _add_header(slide, "SECTION", "Topic title", "08")
    _add_logo(slide)

    outer = _round_rect(
        slide,
        MARGIN_X,
        Inches(1.15),
        Inches(11.9),
        Inches(5.55),
        SURFACE,
        line=CARD_EDGE,
        line_w=1.5,
        adj=0.06,
    )
    _ = outer

    def corner(x, y, *, flip_x=False, flip_y=False):
        hx = -Inches(0.35) if flip_x else Inches(0.35)
        hy = -Inches(0.35) if flip_y else Inches(0.35)
        _rect(slide, x if not flip_x else x + hx, y, abs(hx), Inches(0.04), ACCENT)
        _rect(slide, x, y if not flip_y else y + hy, Inches(0.04), abs(hy), ACCENT)

    corner(MARGIN_X + Inches(0.28), Inches(1.4))
    corner(Inches(12.25), Inches(1.4), flip_x=True)
    corner(MARGIN_X + Inches(0.28), Inches(6.4), flip_y=True)
    corner(Inches(12.25), Inches(6.4), flip_x=True, flip_y=True)

    _ellipse(slide, Inches(6.15), Inches(3.15), Inches(1.0), Inches(1.0), SOFT_PURPLE)
    _add_textbox(
        slide,
        Inches(5.7),
        Inches(3.35),
        Inches(1.9),
        Inches(0.6),
        "+",
        size=28,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        slide,
        Inches(2.5),
        Inches(4.4),
        Inches(8.3),
        Inches(0.5),
        "Drop a diagram, screenshot, or freeform visual here",
        size=16,
        color=PLACEHOLDER_OK,
        align=PP_ALIGN.CENTER,
    )
    _add_footer(slide, 0.95)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_guide(prs)
    slide_title(prs)
    slide_section_divider(prs)
    slide_content_bullets(prs)
    slide_two_columns(prs)
    slide_title_body(prs)
    slide_three_cards(prs)
    slide_quote(prs)
    slide_blank_canvas(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
