#!/usr/bin/env python3
"""
Build Week 1 · Session 5 presentation using official ETRA Design System.
Content is filled incrementally as teaching blocks are pasted in.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week1-Session5-SVM-Kernel-Methods.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session5-diagrams"

SESSION = {
    "eyebrow": "Week 1  ·  Session 5",
    "section_number": "05",
    "section_tag": "S05",
    "section_title": "Classical ML",
    "subtitle": "SVM and kernel methods for classification",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "SVM · Kernel Methods",
    "topics": [
        "Margin maximization",
        "Kernel trick",
        "Support vectors",
        "Linear vs RBF",
    ],
}

TOPIC_CONTENT: dict = {
    "Margin maximization": {
        "title": "Support Vector Machine (SVM): Core Idea",
        "kicker": "Maximum Margin Classifier",
        "body": "SVM is a supervised method that finds the best separating boundary (hyperplane) between classes.",
        "bullets": [
            "It maximizes the margin between classes.",
            "Only support vectors (closest points) determine the boundary.",
            "A larger margin usually improves generalization.",
        ],
        "plot_path": "svm-maximum-margin.png",
        "note": "SVM Maximum Margin · widest gap between classes · support vectors define the margin.",
        "extra_slides": [
            {
                "title": "How SVM Classification Works",
                "kicker": "Decision Function",
                "body": "New points are classified by the side of the hyperplane they fall on.",
                "formula": "f(x) = wᵀx + b,  ŷ = sign(f(x))",
                "formula_tex": r"f(x)=w^{\top}x+b,\quad \hat{y}=\operatorname{sign}(f(x))",
                "formula_note": "f(x)>0 → class +1 · f(x)<0 → class −1",
                "layout": "formula_example",
                "bullets": [
                    "Hard margin: no misclassification allowed (strict separation).",
                    "Soft margin: allows violations for better robustness with overlap/noise.",
                ],
                "note": "Soft-margin SVM adds penalty C to control the margin–violation tradeoff.",
            },
            {
                "title": "Why SVM Is Powerful",
                "kicker": "When SVM Shines",
                "layout": "table",
                "table": {
                    "headers": ["Capability", "Why It Matters"],
                    "rows": [
                        [
                            "Maximum-margin principle",
                            "Better robustness to small perturbations",
                        ],
                        [
                            "Kernel support",
                            "Handles non-linear class boundaries",
                        ],
                        [
                            "High-dimensional performance",
                            "Works well in text and sparse feature spaces",
                        ],
                    ],
                },
                "note": "Common in text/image features · often effective on small-to-medium datasets. Support vectors define the margin.",
            },
        ],
    },
    "Kernel trick": {
        "title": "Kernel Trick: From Non-Linear to Linear Separation",
        "kicker": "Feature Mapping Without Explicit φ(x)",
        "body": "When data is not linearly separable in original space, SVM uses kernel functions to separate it in a transformed feature space.",
        "bullets": [
            "Avoids explicit high-dimensional mapping in many cases.",
            "Computes similarity using kernel function K(xᵢ, xⱼ).",
        ],
        "plot_path": "kernel-trick.png",
        "note": "Work with K(xᵢ, xⱼ) = ⟨φ(xᵢ), φ(xⱼ)⟩ instead of building φ(x) by hand.",
        "extra_slides": [
            {
                "title": "Kernel Similarity K(xᵢ, xⱼ)",
                "kicker": "Kernel Trick",
                "body": "Avoids explicit high-dimensional mapping · computes similarity with K(xᵢ, xⱼ).",
                "layout": "diagram",
                "plot_path": "kernel-similarity.png",
                "note": "SVM decision depends on kernel similarities to support vectors — not on building φ(x).",
            },
            {
                "title": "RBF Kernel in SVM",
                "kicker": "Gaussian Similarity",
                "body": "RBF measures how close x is to a landmark lᵢ — a common default for non-linear SVM.",
                "formula": "K(x, lᵢ) = exp(−∥x − lᵢ∥² / (2σ²))",
                "formula_tex": r"K(x,l_i)=\exp\!\left(-\dfrac{\|x-l_i\|^2}{2\sigma^2}\right)",
                "formula_note": "Close to lᵢ → K≈1 · far → K≈0",
                "layout": "formula_example",
                "bullets": [
                    "If x is close to landmark lᵢ, similarity is near 1.",
                    "If x is far, similarity approaches 0.",
                    "RBF is a common default for non-linear SVM.",
                ],
                "note": "Effect of σ: large → smoother, wider boundary (higher bias) · small → tighter, complex boundary (higher variance).",
            },
            {
                "title": "Effect of Sigma (σ) in RBF",
                "kicker": "Bias–Variance of the Kernel Width",
                "body": "σ controls how far a landmark’s influence reaches.",
                "bullets": [
                    "Large sigma: smoother, wider decision boundary (higher bias).",
                    "Small sigma: tighter, complex boundary (higher variance / overfitting risk).",
                ],
                "plot_path": "rbf-sigma-effect.png",
                "note": "Tune σ (or gamma) with cross-validation together with C.",
            },
            {
                "title": "Common SVM Kernels",
                "kicker": "Choosing K(xᵢ, xⱼ)",
                "layout": "table",
                "table": {
                    "headers": ["Kernel", "Typical Use", "Notes"],
                    "rows": [
                        [
                            "Linear",
                            "High-dimensional sparse data",
                            "Fast and interpretable margin",
                        ],
                        [
                            "RBF (Gaussian)",
                            "General non-linear patterns",
                            "Strong baseline in many tasks",
                        ],
                        [
                            "Polynomial",
                            "Polynomial-like interactions",
                            "Degree controls complexity",
                        ],
                        [
                            "Sigmoid",
                            "Neural-style boundary behavior",
                            "Less common in practice",
                        ],
                    ],
                },
                "note": "Start with Linear (text/sparse) or RBF (general). Support vectors still define the margin.",
            },
            {
                "title": "When to Use SVM / When Not",
                "kicker": "Practical Guidance",
                "layout": "table",
                "table": {
                    "headers": ["Use SVM When", "Avoid SVM When"],
                    "rows": [
                        [
                            "Data is high-dimensional or moderately sized",
                            "Dataset is very large (training can be expensive)",
                        ],
                        [
                            "Classes are reasonably separable",
                            "Data is extremely noisy with weak class structure",
                        ],
                        [
                            "You need non-linear boundaries via kernels",
                            "Kernel matrix becomes too costly (features × samples)",
                        ],
                    ],
                },
                "note": "Use: medium data, clear margin, high-dim text · Avoid: massive data, heavy noise without tuning.",
            },
        ],
    },
    "Support vectors": {
        "title": "Support Vectors — Why They Matter",
        "kicker": "Sparse Decision Boundary",
        "body": "Only the points on (or inside) the margin shape the SVM — the rest can be removed without changing the boundary.",
        "bullets": [
            "Support vectors define the margin edges.",
            "Decision depends on a sparse subset of training points.",
            "Removing non-support points usually leaves the model unchanged.",
        ],
        "plot_path": "svm-maximum-margin.png",
        "note": "Wider margin · fewer influential points · often better generalization.",
    },
    "Linear vs RBF": {
        "title": "SVM vs Logistic Regression vs K-NN",
        "kicker": "Model Comparison",
        "layout": "table",
        "table": {
            "headers": ["Aspect", "SVM", "Logistic Regression", "K-NN"],
            "rows": [
                [
                    "Primary task",
                    "Classification (and SVR)",
                    "Classification",
                    "Classification",
                ],
                [
                    "Non-linearity",
                    "Yes, with kernels",
                    "Limited in linear form",
                    "Yes (distance-based)",
                ],
                [
                    "Compute profile",
                    "Can be heavy on large data",
                    "Usually efficient",
                    "Heavy at prediction",
                ],
                [
                    "Best for",
                    "Complex margins, high-dim spaces",
                    "Fast interpretable baseline",
                    "Local neighborhood patterns",
                ],
            ],
        },
        "note": "LogReg: probabilities + τ=0.5 · SVM: max-margin · K-NN: local votes.",
        "extra_slides": [
            {
                "title": "SVM Practical Hyperparameters",
                "kicker": "What to Tune",
                "layout": "table",
                "table": {
                    "headers": ["Parameter", "Role", "Typical Tuning Direction"],
                    "rows": [
                        [
                            "C",
                            "Penalty for margin violations",
                            "Higher C → stricter fit · lower C → smoother margin",
                        ],
                        [
                            "Kernel",
                            "Similarity function",
                            "Start with RBF · compare linear/polynomial",
                        ],
                        [
                            "gamma (RBF)",
                            "Locality of influence",
                            "Higher gamma → tighter boundary",
                        ],
                        [
                            "degree (poly)",
                            "Polynomial complexity",
                            "Increase only when needed",
                        ],
                    ],
                },
                "note": "Scale features first · tune with CV · support vectors define the margin.",
            },
        ],
    },
}

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "After trees and metrics — maximum-margin classifiers and kernels.",
    "current": "S5",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.15), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.65),
        Inches(11.5),
        Inches(0.9),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.7), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.05), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    x = MARGIN
    for topic in s["topics"]:
        w = Inches(min(2.8, 0.12 * len(topic) + 1.1))
        if x + w > Inches(12.4):
            break
        soft_card(slide, x, Inches(5.25), w, Inches(0.42), fill=SOFT)
        add_text(slide, x, Inches(5.3), w, Inches(0.32), topic, size=11, color=PRIMARY, align=PP_ALIGN.CENTER)
        x += w + Inches(0.14)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    """Fit a plot PNG inside a box (same idea as add_diagram)."""
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return None
    aspect = px_w / px_h
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(path),
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    pitch = 0.50 if dense else 0.62
    bsize = 14 if dense else 16

    if content.get("body"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.45),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(3.6)
    else:
        bullet_top = Inches(2.35)
        bsize = 15 if dense else 17

    # Keep bullets above the note/footer band
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    """Definition + formula + optional scatter/line plot."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)

    soft_card(slide, MARGIN, Inches(2.2), col_w, Inches(1.2), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.3),
        Inches(2.35),
        text_w,
        Inches(0.95),
        content.get("body") or "",
        size=13,
        color=INK,
    )

    if has_formula:
        formula_card_h = 1.65 if is_frac else 1.35
        formula_box_h = 0.95 if is_frac else 0.65
        soft_card(slide, MARGIN, Inches(3.55), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(3.65),
            text_w,
            Inches(0.22),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(3.9),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        if content.get("formula_note"):
            note_y = 3.9 + formula_box_h + 0.05
            add_text(
                slide,
                Inches(MARGIN.inches + 0.3),
                Inches(note_y),
                text_w,
                Inches(0.3),
                content["formula_note"],
                size=10,
                color=MUTED,
            )
        bullet_top = 3.55 + formula_card_h + 0.12
        pitch = 0.46 if (has_note or len(items) >= 2) else 0.55
        bsize = 12
    else:
        bullet_top = 3.55
        pitch = 0.50 if (has_note or len(items) >= 3) else 0.58
        bsize = 13

    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    # With formula+plot keep at most 2 bullets to avoid crowding
    if has_formula and has_plot:
        max_items = min(max_items, 2)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        left=MARGIN,
        width=bullet_w,
        pitch=pitch,
        item_height=Inches(0.42),
    )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    """Full-width practical example with formula and interpretation."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(2.15)
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.28),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(3.0)
    else:
        formula_top = Inches(2.25)

    formula_h = 1.55 if is_frac else 1.25
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.12),
        Inches(11.2),
        Inches(0.95 if is_frac else 0.7),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )
    if content.get("formula_note"):
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(formula_top.inches + formula_h - 0.35),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    bullet_top = formula_top.inches + formula_h + 0.12
    pitch = 0.46 if dense else 0.56
    bsize = 13 if dense else 15
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    if has_note:
        max_items = min(max_items, 3)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(0.42),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    # Prefer kicker over a long body when the table is dense
    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = 2.25
    if show_body:
        soft_card(slide, MARGIN, Inches(2.15), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.25),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = 2.85

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    # If note would still collide, drop the note rather than overlap the table
    note_top = table_top + table_h + 0.08
    render_note = has_note and (note_top + 0.45) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    """Title + optional short body + full-width diagram."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    diagram_top = Inches(2.15)
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(2.1),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(2.45)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(4.35), fill=SOFT)
    _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + 0.15),
        Inches(11.6),
        Inches(4.05),
        folder=folder,
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.9),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    """Returns number of slides added."""
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(4.25),
        Inches(11.2),
        Inches(0.4),
        "Content coming next — paste the teaching block when ready.",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 4 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
