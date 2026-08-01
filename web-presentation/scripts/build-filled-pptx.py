#!/usr/bin/env python3
"""
Build Week 1 · Session 1 presentation using official ETRA Design System.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week1-Session1-Foundations.pptx"

SESSION = {
    "eyebrow": "Week 1  ·  Session 1",
    "section_number": "01",
    "section_tag": "S01",
    "section_title": "Foundations",
    "subtitle": "Data pre-processing for reliable machine learning",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "Data Pre-Processing",
    "topics": [
        "Data Pre-Processing",
        "Train/test splits",
        "Feature scaling",
        "Encoding & imputation",
        "Data leakage",
    ],
}

ML_PROCESS = {
    "title": "The Machine Learning Process",
    "subtitle": "The 3 main steps",
    "steps": [
        {
            "number": "1",
            "title": "Data Pre-Processing",
            "bullets": [
                "Import the data",
                "Clean the data",
                "Split into training and test sets",
                "Feature scaling",
            ],
        },
        {
            "number": "2",
            "title": "Modelling",
            "bullets": [
                "Build the model",
                "Train the model",
                "Make predictions",
            ],
        },
        {
            "number": "3",
            "title": "Evaluation",
            "bullets": [
                "Calculate performance metrics",
                "Make a final prediction",
            ],
        },
    ],
}

TOPIC_CONTENT = {
    "Data Pre-Processing": {
        "title": "Data Preprocessing Template",
        "kicker": "Standard Workflow",
        "headline": "Clean → Scale → Encode → Split",
        "body": "Data preprocessing before modeling — a repeatable checklist before you train.",
        "note": "Start by loading the dataset, then follow the workflow in order.",
        "detail": {
            "title": "Standard Workflow",
            "kicker": "Data Preprocessing Template",
            "bullets": [
                "Load the dataset.",
                "Clean — fix missing values, duplicates, and inconsistent records.",
                "Scale — put numeric features on a comparable range.",
                "Encode — convert categorical variables into model-ready form.",
                "Split — create train/test sets before fitting the model.",
            ],
        },
        "workflow": {
            "title": "Standard Workflow",
            "kicker": "Data Preprocessing Template",
            "lead": "Data preprocessing before modeling",
            "steps": [
                ("1", "Clean", "Handle missing values, errors, and noise."),
                ("2", "Scale", "Normalize or standardize numeric features."),
                ("3", "Encode", "Turn categories into numeric representations."),
                ("4", "Split", "Separate training and test data."),
            ],
        },
    },
    "Train/test splits": {
        "title": "Training Set and Test Set",
        "kicker": "Key Idea",
        "headline": "Train / Test Split",
        "body": "Hold-out partition — train wide, test narrow",
        "stats": [
            ("Train", "80%"),
            ("Test", "20%"),
        ],
        "note": "Typical split: 80% train · 20% test",
        "detail": {
            "title": "Training Set and Test Set",
            "kicker": "Definitions",
            "cards": [
                (
                    "Training Set",
                    "Data used by the model to learn patterns.",
                ),
                (
                    "Test Set",
                    "New, unseen data used to check model performance.",
                ),
            ],
            "bullets": [
                "Never tune hyperparameters on the test set — it becomes part of training indirectly.",
                "Stratified splits preserve class balance in classification tasks.",
                "We train on one part of the data, then test on another part to make sure the model generalizes well.",
            ],
        },
    },
    "Feature scaling": {
        "title": "Feature Scaling",
        "kicker": "What is Feature Scaling?",
        "body": "Feature scaling means adjusting feature values to a similar range using methods like standardization and normalization.",
        "detail": {
            "title": "Feature Scaling",
            "kicker": "Why it matters",
            "cards": [
                (
                    "Comparable range",
                    "Features brought to a comparable range.",
                ),
                (
                    "After scaling",
                    "Heights are comparable.",
                ),
            ],
            "bullets": [
                "It prevents large-value features from dominating smaller-value features.",
                "It improves performance for distance-based and gradient-based algorithms (e.g., KNN, SVM, Gradient Descent).",
                "It helps training converge faster.",
            ],
        },
        "table": {
            "title": "Example (Before vs After Scaling)",
            "kicker": "Worked example",
            "headers": [
                "Student",
                "Height (cm)",
                "Weight (kg)",
                "Height (scaled)",
                "Weight (scaled)",
            ],
            "rows": [
                ["A", "150", "50", "0.25", "0.22"],
                ["B", "170", "70", "0.58", "0.67"],
                ["C", "185", "85", "0.83", "1.00"],
            ],
            "note": "Before scaling, height and weight have different numeric ranges. After scaling, both features are on a comparable scale and contribute more fairly.",
        },
        "methods": {
            "title": "Feature Scaling Methods",
            "kicker": "Feature Scaling",
            "cards": [
                (
                    "Comparable range",
                    "Features brought to a comparable range.",
                ),
                (
                    "After scaling",
                    "Heights are comparable.",
                ),
            ],
            "methods": [
                (
                    "Normalization",
                    "Scales to [0, 1] using min–max — sensitive to outliers.",
                ),
                (
                    "Standardization",
                    "Uses mean 0 and std 1 — preferred when features are roughly Gaussian.",
                ),
                (
                    "Robust scaling",
                    "Uses median and IQR when outliers dominate.",
                ),
            ],
        },
        "method_details": [
            {
                "title": "1) Normalization (Min-Max Scaling)",
                "kicker": "Feature Scaling Methods",
                "lead": "Normalization transforms each value to a range between 0 and 1.",
                "formula": "x_{norm} = (x − x_{min}) / (x_{max} − x_{min})",
                "bullets": [
                    "Useful for models that depend on distances or absolute magnitudes.",
                    "Common choice for algorithms such as KNN and K-Means.",
                ],
            },
            {
                "title": "2) Standardization (Z-score Scaling)",
                "kicker": "Feature Scaling Methods",
                "lead": "Standardization centers data around 0 with a standard deviation of 1.",
                "formula": "z = (x − μ) / σ",
                "formula_note": "μ = mean of the feature · σ = standard deviation of the feature",
                "bullets": [
                    "Preferred when features are close to a Gaussian distribution.",
                    "Works well with Logistic Regression, SVM, and Linear Regression.",
                ],
            },
        ],
        "choosing": {
            "title": "Choosing the Right Scaling Method",
            "kicker": "Practical guidance",
            "bullets": [
                "Use standardization for SVM, logistic regression, PCA, and gradient-based training.",
                "Tree-based models are scale-invariant — scaling is optional.",
                "When in doubt, fit the scaler on train only, then transform train and test.",
            ],
            "focus_card": {
                "title": "Normalization",
                "body": "Use Normalization when feature values have very different ranges and you want all values mapped to [0, 1].",
                "note": "Typical method: Min-Max Scaling.",
            },
            "follow_up": {
                "title": "Normalization vs Standardization",
                "kicker": "When to use each",
                "cards": [
                    {
                        "title": "Normalization",
                        "body": "Good for distance-based models where magnitude differences strongly affect results.",
                        "note": "Maps values to [0, 1] · typical method: Min-Max Scaling.",
                    },
                    {
                        "title": "Standardization",
                        "body": "Use when features are expected to follow (or be close to) a Gaussian distribution, or when the dataset contains outliers.",
                        "note": "Keeps data centered around 0 and scales variability to 1.",
                    },
                ],
                "bullets": [
                    "Often preferred for linear models and optimization-based algorithms.",
                ],
            },
        },
        "norm_example": {
            "title": "Feature Scaling Example (Normalization)",
            "kicker": "Worked example",
            "lead": "Given feature values: [20, 40, 60, 80, 100]",
            "cards": [
                (
                    "Comparable range",
                    "Features brought to a comparable range.",
                ),
                (
                    "After scaling",
                    "Heights are comparable.",
                ),
            ],
            "formula": "x_{norm} = (x − x_{min}) / (x_{max} − x_{min})",
            "formula_note": "x_{min} = 20 · x_{max} = 100",
            "values": [
                ("20", "0.00"),
                ("40", "0.25"),
                ("60", "0.50"),
                ("80", "0.75"),
                ("100", "1.00"),
            ],
            "table": {
                "title": "Before vs After Normalization",
                "kicker": "Feature Scaling Example",
                "formula": "x_{norm} = (x − x_{min}) / (x_{max} − x_{min})",
                "formula_note": "x_{min} = 20 · x_{max} = 100",
                "headers": [
                    "Original value (x)",
                    "Normalized value (xₙₒᵣₘ)",
                ],
                "rows": [
                    ["20", "0.00"],
                    ["40", "0.25"],
                    ["60", "0.50"],
                    ["80", "0.75"],
                    ["100", "1.00"],
                ],
                "note": "This table shows how all values are scaled to the range [0, 1].",
            },
        },
        "std_example": {
            "title": "Feature Scaling Example (Standardization)",
            "kicker": "Worked example",
            "lead": "Given feature values: [20, 40, 60, 80, 100]",
            "cards": [
                (
                    "Comparable range",
                    "Features brought to a comparable range.",
                ),
                (
                    "After scaling",
                    "Heights are comparable.",
                ),
            ],
            "formula": "z = (x − μ) / σ",
            "formula_note": "Mean: μ = 60 · Standard deviation: σ = 28.28",
            "values": [
                ("20", "−1.41"),
                ("40", "−0.71"),
                ("60", "0.00"),
                ("80", "0.71"),
                ("100", "1.41"),
            ],
            "table": {
                "title": "Before vs After Standardization",
                "kicker": "Feature Scaling Example",
                "formula": "z = (x − μ) / σ",
                "formula_note": "Mean: μ = 60 · Standard deviation: σ = 28.28",
                "headers": [
                    "Original value (x)",
                    "Standardized value (z)",
                ],
                "rows": [
                    ["20", "−1.41"],
                    ["40", "−0.71"],
                    ["60", "0.00"],
                    ["80", "0.71"],
                    ["100", "1.41"],
                ],
                "note": "This table shows that standardized values are centered around 0 and measured in units of standard deviation.",
            },
        },
    },
    "Encoding & imputation": {
        "title": "Categorical Data",
        "kicker": "What is Categorical Data?",
        "body": "Categorical data describes values that belong to named groups rather than continuous numbers.",
        "table": {
            "title": "Category Types",
            "kicker": "Ordinal vs Nominal",
            "headers": ["Type", "Description", "Example"],
            "rows": [
                [
                    "Ordinal",
                    "Categories have a meaningful order.",
                    "Education Level: High School < Bachelor's < Master's < PhD",
                ],
                [
                    "Nominal",
                    "Categories have no natural order.",
                    "Payment Method: Cash, Card, Bank Transfer",
                ],
            ],
            "note": "Choose encoding based on whether order matters.",
        },
        "extra_tables": [
            {
                "title": "Example 1: Label Encoding",
                "kicker": "Encoding categorical data",
                "headers": ["Raw category", "Encoded value"],
                "rows": [
                    ["Cash", "0"],
                    ["Card", "1"],
                    ["Bank Transfer", "2"],
                ],
                "note": "Each category is mapped to a unique integer.",
            },
            {
                "title": "Example 2: One-Hot Encoding",
                "kicker": "Encoding categorical data",
                "headers": ["Raw category", "Cash", "Card", "Bank Transfer"],
                "rows": [
                    ["Cash", "1", "0", "0"],
                    ["Card", "0", "1", "0"],
                    ["Bank Transfer", "0", "0", "1"],
                ],
                "note": "Most machine learning algorithms work with numerical inputs only, so categorical values must be converted into numbers.",
            },
        ],
        "ordinal": {
            "title": "How to Deal with Categorical Data (Ordinal Encoding)",
            "kicker": "What is Ordinal Encoding?",
            "body": "Ordinal Encoding converts ordered categories into integers that preserve their ranking.",
            "comparison": {
                "title": "Encoding Comparison",
                "kicker": "Choosing an encoding",
                "lead": "Three ways to encode categorical variables",
                "cards": [
                    (
                        "Label Encoding",
                        "Maps each category to a unique integer. Simple, but can invent false order for nominal data.",
                    ),
                    (
                        "One-Hot Encoding",
                        "Creates a binary column per category. Safe for nominal data; increases feature count.",
                    ),
                    (
                        "Ordinal Encoding",
                        "Assigns integers that preserve a meaningful ranking for ordered categories.",
                    ),
                ],
            },
            "howto": {
                "title": "How to Apply Ordinal Encoding",
                "kicker": "Practical steps",
                "lead": "Use Ordinal Encoding only when categories have a true order.",
                "steps": [
                    "Step 1: Identify the correct category order.",
                    "Step 2: Assign an integer to each level based on that order.",
                    "Step 3: Replace the original text values with encoded numbers.",
                ],
            },
            "example_table": {
                "title": "Example",
                "kicker": "Ordinal Encoding",
                "headers": ["Satisfaction level", "Encoded value"],
                "rows": [
                    ["Low", "1"],
                    ["Medium", "2"],
                    ["High", "3"],
                ],
                "note": "This method keeps the order information, which is important for many machine learning models.",
            },
        },
        "one_hot": {
            "title": "How to Deal with Categorical Data (One-Hot Encoding)",
            "kicker": "What is One-Hot Encoding?",
            "body": "One-Hot Encoding converts each category into a separate binary column (0 or 1).",
            "comparison_note": {
                "title": "Encoding Comparison",
                "kicker": "Choosing an encoding",
                "lead": "Three ways to encode categorical variables",
                "cards": [
                    (
                        "Label Encoding",
                        "Maps each category to a unique integer. Simple, but can invent false order for nominal data.",
                    ),
                    (
                        "One-Hot Encoding",
                        "Creates a binary column per category. Safe for nominal data; increases feature count.",
                    ),
                    (
                        "Ordinal Encoding",
                        "Assigns integers that preserve a meaningful ranking for ordered categories.",
                    ),
                ],
            },
            "howto": {
                "title": "How to Apply One-Hot Encoding",
                "kicker": "Practical steps",
                "lead": "Use One-Hot Encoding for nominal categories (no natural order).",
                "steps": [
                    "Step 1: List all unique categories.",
                    "Step 2: Create one column for each category.",
                    "Step 3: Put 1 in the matching category column and 0 in all others.",
                ],
            },
            "example_table": {
                "title": "Example",
                "kicker": "One-Hot Encoding",
                "headers": ["Payment method", "Cash", "Card", "Bank Transfer"],
                "rows": [
                    ["Cash", "1", "0", "0"],
                    ["Card", "0", "1", "0"],
                    ["Bank Transfer", "0", "0", "1"],
                ],
                "note": "This method avoids creating a false ranking between categories.",
            },
            "dummy_trap": {
                "title": "Avoiding the Dummy Variable Trap",
                "kicker": "One-Hot Encoding",
                "lead": "It helps prevent the dummy variable trap (perfect multicollinearity), especially in linear models.",
                "steps": [
                    "Step 1: Start with One-Hot encoded columns.",
                    "Step 2: Drop one category as a reference (baseline).",
                    "Step 3: Keep the remaining binary columns.",
                ],
                "context": "Original categories: Payment Method = [Cash, Card, Bank Transfer].\nReference category (dropped): Cash.",
                "example_table": {
                    "title": "Example (Reference = Cash)",
                    "kicker": "Dummy variable trap",
                    "headers": ["Payment method", "Card", "Bank Transfer"],
                    "rows": [
                        ["Cash", "0", "0"],
                        ["Card", "1", "0"],
                        ["Bank Transfer", "0", "1"],
                    ],
                    "note": "In this setup, Cash is represented when all remaining columns are 0.",
                },
            },
        },
        "advanced_encoding": {
            "overview_table": {
                "title": "Encoding Methods Overview",
                "kicker": "Beyond label / one-hot / ordinal",
                "headers": ["Encoding type", "Idea", "Best use case"],
                "rows": [
                    [
                        "Frequency Encoding",
                        "Replace each category with how often it appears.",
                        "Large datasets with many categories",
                    ],
                    [
                        "Count Encoding",
                        "Similar to frequency, but uses raw counts directly.",
                        "Tree-based models with high-cardinality features",
                    ],
                    [
                        "Target Encoding",
                        "Replace category with target mean for that category.",
                        "Supervised tasks (use with care to avoid leakage)",
                    ],
                    [
                        "Binary Encoding",
                        "Convert category index to binary digits across columns.",
                        "High-cardinality data with fewer columns than one-hot",
                    ],
                    [
                        "Hash Encoding",
                        "Use a hash function to map categories into fixed columns.",
                        "Very large and dynamic category sets",
                    ],
                ],
            },
            "extra_tables": [
                {
                    "title": "1) Frequency Encoding (City)",
                    "kicker": "Encoding methods",
                    "headers": ["City", "Frequency"],
                    "rows": [
                        ["Riyadh", "0.50"],
                        ["Jeddah", "0.33"],
                        ["Dammam", "0.17"],
                    ],
                },
                {
                    "title": "2) Count Encoding (City)",
                    "kicker": "Encoding methods",
                    "headers": ["City", "Count"],
                    "rows": [
                        ["Riyadh", "3"],
                        ["Jeddah", "2"],
                        ["Dammam", "1"],
                    ],
                },
                {
                    "title": "3) Target Encoding (City)",
                    "kicker": "Encoding methods",
                    "headers": ["City", "Mean target"],
                    "rows": [
                        ["Riyadh", "0.80"],
                        ["Jeddah", "0.40"],
                        ["Dammam", "0.20"],
                    ],
                },
                {
                    "title": "4) Binary Encoding (City)",
                    "kicker": "Encoding methods",
                    "headers": ["City", "Category ID", "Binary code"],
                    "rows": [
                        ["Riyadh", "1", "01"],
                        ["Jeddah", "2", "10"],
                        ["Dammam", "3", "11"],
                    ],
                },
                {
                    "title": "5) Hash Encoding (4 Buckets)",
                    "kicker": "Encoding methods",
                    "headers": ["City", "Hash bucket"],
                    "rows": [
                        ["Riyadh", "2"],
                        ["Jeddah", "0"],
                        ["Dammam", "3"],
                    ],
                },
            ],
        },
        "missing_data": {
            "title": "Handling Missing Data",
            "kicker": "What is Missing Data?",
            "body": "Missing data refers to values that are absent or incomplete in one or more features of a dataset.",
            "bullets": [
                "It can introduce bias into model training.",
                "It may reduce prediction accuracy.",
                "It can cause errors in analysis if not handled properly.",
            ],
            "types_table": {
                "title": "Types of Missing Data",
                "kicker": "MCAR · MAR · NMAR",
                "headers": ["Type", "Meaning"],
                "rows": [
                    [
                        "MCAR (Missing Completely at Random)",
                        "Missingness happens randomly and is unrelated to any variable in the dataset.",
                    ],
                    [
                        "MAR (Missing at Random)",
                        "Missingness depends on other observed variables, but not on the missing value itself.",
                    ],
                    [
                        "NMAR (Not Missing at Random)",
                        "Missingness depends on the missing value itself or unobserved factors.",
                    ],
                ],
            },
            "causes": {
                "title": "Why Do We Have Missing Data?",
                "kicker": "Handling Missing Data",
                "lead": "Missing values can appear for many practical reasons during data collection and processing.",
                "table": {
                    "title": "Common Causes",
                    "kicker": "Why values go missing",
                    "headers": ["Cause", "Description", "Example"],
                    "rows": [
                        [
                            "Human Error",
                            "Data is skipped, entered incorrectly, or forgotten.",
                            "A user leaves the Age field empty in a form.",
                        ],
                        [
                            "Device/Sensor Failure",
                            "Measurement tools fail or disconnect.",
                            "A medical sensor stops recording heart rate temporarily.",
                        ],
                        [
                            "Data Integration Issues",
                            "Missing fields appear when combining data from multiple sources.",
                            "One database has salary, another does not.",
                        ],
                        [
                            "Privacy or Refusal",
                            "Participants choose not to share sensitive information.",
                            "A customer does not provide income details.",
                        ],
                        [
                            "System/Transmission Errors",
                            "Data is lost during storage, transfer, or export.",
                            "Network interruption causes missing rows in logs.",
                        ],
                        [
                            "Conditional Missingness",
                            "Some fields are only relevant for specific groups.",
                            "Pregnancy_Status is empty for male patients.",
                        ],
                    ],
                    "note": "Understanding why data is missing helps us choose the right treatment method (deletion, imputation, or advanced modeling).",
                },
            },
            "mcar": {
                "title": "Missing Completely at Random (MCAR)",
                "kicker": "Handling Missing Data · MCAR",
                "lead": "MCAR means missing values occur randomly and are independent of both observed and unobserved variables.",
                "bullets": [
                    "MCAR has no systematic pattern.",
                    "In practice, MCAR is often an unrealistic assumption in real-world datasets.",
                    "Before assuming MCAR, review data collection workflow and consult domain experts.",
                    "If MCAR is reasonable, simple imputation (mean/median/mode) can work.",
                ],
                "example_table": {
                    "title": "Illustrative Data Example",
                    "kicker": "Handling Missing Data · MCAR",
                    "headers": ["Student ID", "Age", "Score"],
                    "rows": [
                        ["101", "20", "85"],
                        ["102", "Missing", "88"],
                        ["103", "21", "Missing"],
                        ["104", "23", "91"],
                        ["105", "Missing", "76"],
                        ["106", "22", "Missing"],
                    ],
                    "note": "Missing values appear across different rows and columns with no clear pattern, which is consistent with MCAR.",
                },
            },
            "mar": {
                "title": "Missing at Random (MAR)",
                "kicker": "Handling Missing Data · MAR",
                "lead": "MAR means missingness is related to observed features in the dataset, but not to the missing feature itself.",
                "bullets": [
                    "Missingness depends on other observed variables.",
                    "The missing value itself is not the direct cause of the gap.",
                    "More realistic than MCAR in many real-world datasets.",
                    "Model-based imputation that uses observed features usually works better than simple mean/median fills.",
                ],
                "methods_table": {
                    "title": "Imputation Techniques for MAR",
                    "kicker": "Handling Missing Data · MAR",
                    "headers": ["Method", "Idea", "Why it works for MAR"],
                    "rows": [
                        [
                            "kNN Imputation",
                            "Fills missing values using the nearest similar records.",
                            "Uses observed features to find similar patterns.",
                        ],
                        [
                            "MissForest",
                            "Uses random-forest models iteratively to predict missing values.",
                            "Captures non-linear relationships from observed data effectively.",
                        ],
                    ],
                    "note": "Both methods leverage observed features — exactly what MAR assumes is available and informative.",
                },
                "knn_example": {
                    "title": "Example: kNN Imputation (k = 2)",
                    "kicker": "Handling Missing Data · MAR",
                    "lead": "For student C, the two nearest students based on Grade and Attendance are A and B.",
                    "bullets": [
                        "Use observed features (Grade, Attendance) to find the nearest neighbors.",
                        "Impute student C’s missing value from students A and B (e.g. average).",
                        "This works under MAR because missingness is explained by observed patterns, not by the missing value alone.",
                    ],
                    "example_table": {
                        "title": "Example: kNN Imputation (k = 2)",
                        "kicker": "Handling Missing Data · MAR",
                        "headers": ["Student", "Grade", "Attendance (%)", "Study Hours"],
                        "rows": [
                            ["A", "90", "95", "8"],
                            ["B", "88", "92", "7"],
                            ["C", "91", "96", "Missing"],
                            ["D", "70", "75", "3"],
                        ],
                        "note": "Student C is closest to A and B on Grade and Attendance — so we ignore distant student D.",
                    },
                    "result": {
                        "title": "Example: kNN Imputation (k = 2)",
                        "kicker": "Handling Missing Data · MAR",
                        "formula": "Study Hours_{C} = (8 + 7) / 2 = 7.5",
                        "lead": "Average the Study Hours of the two nearest neighbors (A and B) to fill the missing value for student C.",
                        "bullets": [
                            "Neighbor A → Study Hours = 8",
                            "Neighbor B → Study Hours = 7",
                            "Imputed value for C → 7.5",
                        ],
                    },
                },
                "missforest_example": {
                    "title": "Example: MissForest",
                    "kicker": "Handling Missing Data · MAR",
                    "lead": "MissForest trains a Random Forest model using known rows to predict the missing value for student C.",
                    "bullets": [
                        "Treat the incomplete column (Study Hours) as the target.",
                        "Train on complete rows (A, B, D) using observed features such as Grade and Attendance.",
                        "Predict Study Hours for student C from the learned non-linear relationships.",
                        "Iterate across features with missing values until predictions stabilize.",
                    ],
                    "result": {
                        "title": "Example: MissForest",
                        "kicker": "Handling Missing Data · MAR",
                        "formula": "Study Hours_{C} ≈ 7.6",
                        "lead": "The Random Forest prediction for student C’s missing Study Hours is approximately 7.6.",
                        "bullets": [
                            "Uses non-linear patterns learned from complete rows.",
                            "Often more accurate than simple neighbor averages when relationships are complex.",
                        ],
                    },
                    "example_table": {
                        "title": "Example (Table)",
                        "kicker": "Handling Missing Data · MAR",
                        "headers": ["Student", "Grade (Observed)", "Study Hours"],
                        "rows": [
                            ["1", "95", "Missing"],
                            ["2", "92", "Missing"],
                            ["3", "88", "4"],
                            ["4", "75", "6"],
                            ["5", "70", "7"],
                        ],
                        "note": "MAR is common in practice because missingness can often be explained by available variables.",
                    },
                },
            },
            "mnar": {
                "title": "Missing Not at Random (MNAR)",
                "kicker": "Handling Missing Data · MNAR",
                "lead": "MNAR is usually the most complex type because missingness depends on the missing value itself or unobserved variables.",
                "bullets": [
                    "MNAR has a systematic pattern tied to hidden information.",
                    "Standard imputation can be unreliable without additional data or strong domain assumptions.",
                    "Often needs specialized models, sensitivity analysis, or collecting the missing information directly.",
                    "Before treating MNAR like MAR/MCAR, validate the missingness mechanism with domain experts.",
                ],
                "example_table": {
                    "title": "Example (Table)",
                    "kicker": "Handling Missing Data · MNAR",
                    "headers": ["Participant", "Reported Stress Level"],
                    "rows": [
                        ["1", "2"],
                        ["2", "Missing"],
                        ["3", "3"],
                        ["4", "Missing"],
                        ["5", "1"],
                    ],
                    "note": "Missing stress reports may be tied to the stress level itself — people with higher stress may skip the question.",
                },
                "preserve_pattern_table": {
                    "title": "Practical Tip: Preserve Missingness Pattern",
                    "kicker": "Handling Missing Data · MNAR",
                    "headers": ["Participant", "Stress (Imputed)", "Stress_Was_Missing"],
                    "rows": [
                        ["1", "2.0", "0"],
                        ["2", "4.1", "1"],
                        ["3", "3.0", "0"],
                        ["4", "4.3", "1"],
                        ["5", "1.0", "0"],
                    ],
                    "note": "Add an indicator column so the model can still learn from the fact that a value was missing.",
                },
            },
        },
    },
}

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "Build the end-to-end ML workflow: clean data → train → evaluate.",
    "current": "S1",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(slide, MARGIN, Inches(2.7), Inches(11.5), Inches(0.95), s["section_title"], size=44, bold=True, color=PRIMARY)
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.15), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(slide, MARGIN, Inches(2.65), Inches(11.5), Inches(0.9), s["section_title"], size=42, bold=True, color=PRIMARY)
    bar = rect(slide, MARGIN, Inches(3.7), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.05), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    x = MARGIN
    for topic in s["topics"]:
        w = Inches(min(2.35, 0.11 * len(topic) + 1.0))
        if x + w > Inches(12.4):
            break
        soft_card(slide, x, Inches(5.25), w, Inches(0.42), fill=SOFT)
        add_text(slide, x, Inches(5.3), w, Inches(0.32), topic, size=11, color=PRIMARY, align=PP_ALIGN.CENTER)
        x += w + Inches(0.14)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=18)
    content_footer(slide, index, total)


def slide_ml_process_overview(prs, total, index):
    s = SESSION
    mp = ML_PROCESS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, mp["title"], mp["subtitle"])

    col_w = Inches(3.75)
    gap = Inches(0.22)
    top = Inches(2.35)

    for i, step in enumerate(mp["steps"]):
        x = MARGIN + i * (col_w + gap)
        soft_card(slide, x, top, col_w, Inches(4.15), fill=SOFT_2 if i % 2 else SOFT)
        add_text(
            slide,
            x + Inches(0.35),
            top + Inches(0.35),
            Inches(3),
            Inches(0.3),
            f"Step {step['number']}",
            size=12,
            color=SECONDARY,
        )
        add_text(
            slide,
            x + Inches(0.35),
            top + Inches(0.75),
            Inches(3.1),
            Inches(0.55),
            step["title"],
            size=17,
            bold=True,
            color=PRIMARY,
        )
        for bi, bullet in enumerate(step["bullets"]):
            by = top + Inches(1.55) + Inches(bi * 0.5)
            add_text(slide, x + Inches(0.35), by, Inches(0.25), Inches(0.35), "–", size=14, color=SECONDARY)
            add_text(slide, x + Inches(0.6), by, Inches(2.9), Inches(0.4), bullet, size=13, color=MUTED)

    content_footer(slide, index, total)


def slide_ml_process_step(prs, total, index, step):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, f"Step {step['number']}: {step['title']}", "The Machine Learning Process")
    bullets(slide, step["bullets"], top=Inches(2.4), size=20)
    content_footer(slide, index, total)


def slide_topic(prs, total, index, topic_index, topic):
    """Returns number of slides added (1 or more)."""
    content = TOPIC_CONTENT.get(topic)
    if content:
        slide_topic_rich(prs, total, index, content)
        n = 1
        if content.get("workflow"):
            slide_preprocess_workflow(prs, total, index + n, content["workflow"])
            n += 1
        if content.get("detail"):
            slide_topic_detail(prs, total, index + n, content["detail"])
            n += 1
        if content.get("table"):
            slide_topic_table(prs, total, index + n, content["table"])
            n += 1
        for extra in content.get("extra_tables") or []:
            slide_topic_table(prs, total, index + n, extra)
            n += 1
        if content.get("methods"):
            slide_topic_methods(prs, total, index + n, content["methods"])
            n += 1
        for detail in content.get("method_details") or []:
            slide_method_detail(prs, total, index + n, detail)
            n += 1
        if content.get("choosing"):
            slide_choosing_method(prs, total, index + n, content["choosing"])
            n += 1
            if content["choosing"].get("follow_up"):
                slide_choosing_follow_up(prs, total, index + n, content["choosing"]["follow_up"])
                n += 1
        if content.get("norm_example"):
            slide_norm_example(prs, total, index + n, content["norm_example"])
            n += 1
            if content["norm_example"].get("table"):
                slide_norm_table(prs, total, index + n, content["norm_example"]["table"])
                n += 1
        if content.get("std_example"):
            slide_norm_example(prs, total, index + n, content["std_example"])
            n += 1
            if content["std_example"].get("table"):
                slide_norm_table(prs, total, index + n, content["std_example"]["table"])
                n += 1
        if content.get("ordinal"):
            slide_topic_rich(prs, total, index + n, content["ordinal"])
            n += 1
            if content["ordinal"].get("comparison"):
                slide_encoding_comparison(prs, total, index + n, content["ordinal"]["comparison"])
                n += 1
            if content["ordinal"].get("howto"):
                slide_ordinal_howto(prs, total, index + n, content["ordinal"]["howto"])
                n += 1
            if content["ordinal"].get("example_table"):
                slide_topic_table(prs, total, index + n, content["ordinal"]["example_table"])
                n += 1
        if content.get("one_hot"):
            slide_topic_rich(prs, total, index + n, content["one_hot"])
            n += 1
            if content["one_hot"].get("comparison_note"):
                slide_encoding_comparison(
                    prs, total, index + n, content["one_hot"]["comparison_note"]
                )
                n += 1
            if content["one_hot"].get("howto"):
                slide_ordinal_howto(prs, total, index + n, content["one_hot"]["howto"])
                n += 1
            if content["one_hot"].get("example_table"):
                slide_topic_table(prs, total, index + n, content["one_hot"]["example_table"])
                n += 1
            if content["one_hot"].get("dummy_trap"):
                slide_dummy_trap(prs, total, index + n, content["one_hot"]["dummy_trap"])
                n += 1
                if content["one_hot"]["dummy_trap"].get("example_table"):
                    slide_topic_table(
                        prs, total, index + n, content["one_hot"]["dummy_trap"]["example_table"]
                    )
                    n += 1
        if content.get("advanced_encoding"):
            adv = content["advanced_encoding"]
            if adv.get("overview_table"):
                slide_topic_table(prs, total, index + n, adv["overview_table"])
                n += 1
            for extra in adv.get("extra_tables") or []:
                slide_topic_table(prs, total, index + n, extra)
                n += 1
        if content.get("missing_data"):
            slide_missing_intro(prs, total, index + n, content["missing_data"])
            n += 1
            if content["missing_data"].get("types_table"):
                slide_topic_table(prs, total, index + n, content["missing_data"]["types_table"])
                n += 1
            if content["missing_data"].get("causes"):
                causes = content["missing_data"]["causes"]
                slide_missing_causes_intro(prs, total, index + n, causes)
                n += 1
                if causes.get("table"):
                    slide_topic_table(prs, total, index + n, causes["table"])
                    n += 1
            if content["missing_data"].get("mcar"):
                mcar = content["missing_data"]["mcar"]
                slide_missing_mechanism(prs, total, index + n, mcar)
                n += 1
                if mcar.get("example_table"):
                    slide_topic_table(prs, total, index + n, mcar["example_table"])
                    n += 1
            if content["missing_data"].get("mar"):
                mar = content["missing_data"]["mar"]
                slide_missing_mechanism(prs, total, index + n, mar)
                n += 1
                if mar.get("methods_table"):
                    slide_topic_table(prs, total, index + n, mar["methods_table"])
                    n += 1
                if mar.get("knn_example"):
                    knn = mar["knn_example"]
                    slide_missing_mechanism(prs, total, index + n, knn)
                    n += 1
                    if knn.get("example_table"):
                        slide_topic_table(prs, total, index + n, knn["example_table"])
                        n += 1
                    if knn.get("result"):
                        slide_missing_formula(prs, total, index + n, knn["result"])
                        n += 1
                if mar.get("missforest_example"):
                    mf = mar["missforest_example"]
                    slide_missing_mechanism(prs, total, index + n, mf)
                    n += 1
                    if mf.get("result"):
                        slide_missing_formula(prs, total, index + n, mf["result"])
                        n += 1
                    if mf.get("example_table"):
                        slide_topic_table(prs, total, index + n, mf["example_table"])
                        n += 1
            if content["missing_data"].get("mnar"):
                mnar = content["missing_data"]["mnar"]
                slide_missing_mechanism(prs, total, index + n, mnar)
                n += 1
                if mnar.get("example_table"):
                    slide_topic_table(prs, total, index + n, mnar["example_table"])
                    n += 1
                if mnar.get("preserve_pattern_table"):
                    slide_topic_table(prs, total, index + n, mnar["preserve_pattern_table"])
                    n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.7),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(4.35),
        Inches(11.2),
        Inches(0.4),
        s["focus"],
        size=14,
        color=MUTED,
    )
    content_footer(slide, index, total)
    return 1


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    if content.get("headline"):
        add_text(
            slide,
            MARGIN,
            Inches(2.35),
            Inches(12),
            Inches(0.45),
            content["headline"],
            size=22,
            bold=True,
            color=PRIMARY,
        )
    if content.get("body"):
        soft_card(slide, MARGIN, Inches(2.95), Inches(12.0), Inches(1.6), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.4),
            Inches(3.25),
            Inches(11.2),
            Inches(1.1),
            content["body"],
            size=16,
            color=INK,
        )

    stats = content.get("stats") or []
    if stats:
        card_w = Inches(5.7)
        gap = Inches(0.4)
        for i, (label, value) in enumerate(stats):
            x = MARGIN + i * (card_w + gap)
            soft_card(slide, x, Inches(4.85) if content.get("body") else Inches(3.5), card_w, Inches(1.4), fill=SOFT_2)
            y0 = Inches(5.05) if content.get("body") else Inches(3.75)
            add_text(
                slide,
                x + Inches(0.4),
                y0,
                card_w - Inches(0.8),
                Inches(0.3),
                label,
                size=13,
                color=SECONDARY,
            )
            add_text(
                slide,
                x + Inches(0.4),
                y0 + Inches(0.35),
                card_w - Inches(0.8),
                Inches(0.7),
                value,
                size=28 if content.get("body") else 40,
                bold=True,
                color=PRIMARY,
            )
    elif content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(4.9),
            Inches(12),
            Inches(0.4),
            content["note"],
            size=14,
            color=MUTED,
        )

    if content.get("note") and stats:
        add_text(
            slide,
            MARGIN,
            Inches(6.35),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=13,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_detail(prs, total, index, detail):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, detail["title"], detail.get("kicker"))

    cards = detail.get("cards") or []
    bullet_items = detail.get("bullets") or []
    cards_top = Inches(2.3)
    bullets_top = Inches(4.3)

    # If only bullets (no cards), give them more room near the title
    if bullet_items and not cards:
        bullets_top = Inches(2.35)
    if cards and bullet_items:
        # Compact cards so bullets fit
        card_h = Inches(1.35)
        cards_top = Inches(2.25)
        bullets_top = Inches(3.85)
    else:
        card_h = Inches(1.7)

    if cards:
        card_w = Inches(5.7)
        gap = Inches(0.4)
        for i, (label, body) in enumerate(cards):
            x = MARGIN + i * (card_w + gap)
            soft_card(slide, x, cards_top, card_w, card_h, fill=SOFT)
            add_text(
                slide,
                x + Inches(0.35),
                cards_top + Inches(0.2),
                card_w - Inches(0.7),
                Inches(0.3),
                label,
                size=14,
                bold=True,
                color=PRIMARY,
            )
            add_text(
                slide,
                x + Inches(0.35),
                cards_top + Inches(0.55),
                card_w - Inches(0.7),
                Inches(0.65),
                body,
                size=13,
                color=MUTED,
            )

    if bullet_items:
        bullets(slide, bullet_items, top=bullets_top, size=15)

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, table):
    from pptx.util import Pt

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, table["title"], table.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    long_text = any(len(str(cell)) > 28 for row in rows for cell in row)
    row_h = Inches(0.85) if long_text else Inches(0.55)
    table_h = Inches(0.5) + row_h * len(rows)
    shape = slide.shapes.add_table(
        rows=1 + len(rows),
        cols=len(headers),
        left=MARGIN,
        top=Inches(2.3),
        width=Inches(12.0),
        height=min(table_h, Inches(3.6)),
    )
    tbl = shape.table
    align = PP_ALIGN.LEFT if long_text else PP_ALIGN.CENTER

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = align
                for run in p.runs:
                    run.font.size = Pt(12 if long_text else 13)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if table.get("note"):
        note_top = Inches(2.5) + min(table_h, Inches(3.6)) + Inches(0.2)
        soft_card(slide, MARGIN, note_top, Inches(12.0), Inches(0.9), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            note_top + Inches(0.25),
            Inches(11.3),
            Inches(0.5),
            table["note"],
            size=13,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_missing_intro(prs, total, index, missing):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, missing["title"], missing.get("kicker"))

    if missing.get("body"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.35), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.5),
            Inches(11.3),
            Inches(0.95),
            missing["body"],
            size=15,
            color=INK,
        )

    bullets(slide, missing.get("bullets") or [], top=Inches(3.9), size=16)
    content_footer(slide, index, total)


def slide_missing_causes_intro(prs, total, index, causes):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, causes["title"], causes.get("kicker"))

    if causes.get("lead"):
        soft_card(slide, MARGIN, Inches(2.4), Inches(12.0), Inches(1.6), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.4),
            Inches(2.8),
            Inches(11.2),
            Inches(1.0),
            causes["lead"],
            size=18,
            color=INK,
        )

    content_footer(slide, index, total)


def slide_missing_mechanism(prs, total, index, mechanism):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, mechanism["title"], mechanism.get("kicker"))

    if mechanism.get("lead"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.35), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.5),
            Inches(11.3),
            Inches(0.95),
            mechanism["lead"],
            size=15,
            color=INK,
        )

    bullets(slide, mechanism.get("bullets") or [], top=Inches(3.9), size=16)
    content_footer(slide, index, total)


def slide_missing_formula(prs, total, index, result):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, result["title"], result.get("kicker"))

    soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.5), fill=SOFT)
    add_formula(
        slide,
        MARGIN + Inches(0.35),
        Inches(2.6),
        Inches(11.3),
        Inches(0.5),
        result["formula"],
        size=22,
        bold=True,
        color=PRIMARY,
    )
    if result.get("lead"):
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(3.2),
            Inches(11.3),
            Inches(0.4),
            result["lead"],
            size=14,
            color=MUTED,
        )

    bullets(slide, result.get("bullets") or [], top=Inches(4.05), size=16)
    content_footer(slide, index, total)


def slide_dummy_trap(prs, total, index, trap):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, trap["title"], trap.get("kicker"))

    if trap.get("lead"):
        soft_card(slide, MARGIN, Inches(2.2), Inches(12.0), Inches(1.0), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.4),
            Inches(11.3),
            Inches(0.65),
            trap["lead"],
            size=14,
            color=INK,
        )

    bullets(slide, trap.get("steps") or [], top=Inches(3.45), size=16)

    if trap.get("context"):
        add_text(
            slide,
            MARGIN,
            Inches(5.7),
            Inches(12),
            Inches(0.7),
            trap["context"].replace("\n", "  ·  "),
            size=13,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_ordinal_howto(prs, total, index, howto):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, howto["title"], howto.get("kicker"))

    if howto.get("lead"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(0.85), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.45),
            Inches(11.3),
            Inches(0.5),
            howto["lead"],
            size=15,
            bold=True,
            color=PRIMARY,
        )

    steps = howto.get("steps") or []
    bullets(slide, steps, top=Inches(3.4), size=17)
    content_footer(slide, index, total)


def slide_encoding_comparison(prs, total, index, comparison):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, comparison["title"], comparison.get("kicker"))

    if comparison.get("lead"):
        add_text(
            slide,
            MARGIN,
            Inches(2.25),
            Inches(12),
            Inches(0.35),
            comparison["lead"],
            size=16,
            color=MUTED,
        )

    cards = comparison.get("cards") or []
    col_w = Inches(3.75)
    gap = Inches(0.25)
    top = Inches(2.8)
    for i, (title, body) in enumerate(cards):
        x = MARGIN + i * (col_w + gap)
        soft_card(slide, x, top, col_w, Inches(3.5), fill=SOFT if i % 2 == 0 else SOFT_2)
        add_text(
            slide,
            x + Inches(0.3),
            top + Inches(0.35),
            col_w - Inches(0.6),
            Inches(0.3),
            f"0{i + 1}",
            size=12,
            bold=True,
            color=SECONDARY,
        )
        add_text(
            slide,
            x + Inches(0.3),
            top + Inches(0.75),
            col_w - Inches(0.6),
            Inches(0.55),
            title,
            size=16,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            x + Inches(0.3),
            top + Inches(1.5),
            col_w - Inches(0.6),
            Inches(1.6),
            body,
            size=13,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_choosing_method(prs, total, index, choosing):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, choosing["title"], choosing.get("kicker"))

    bullets(slide, choosing.get("bullets") or [], top=Inches(2.3), size=16)

    focus = choosing.get("focus_card") or {}
    if focus:
        soft_card(slide, MARGIN, Inches(4.55), Inches(12.0), Inches(1.85), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.4),
            Inches(4.75),
            Inches(11.2),
            Inches(0.3),
            focus.get("title", ""),
            size=15,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            MARGIN + Inches(0.4),
            Inches(5.15),
            Inches(11.2),
            Inches(0.55),
            focus.get("body", ""),
            size=14,
            color=MUTED,
        )
        if focus.get("note"):
            add_text(
                slide,
                MARGIN + Inches(0.4),
                Inches(5.8),
                Inches(11.2),
                Inches(0.3),
                focus["note"],
                size=13,
                bold=True,
                color=SECONDARY,
            )

    content_footer(slide, index, total)


def slide_norm_table(prs, total, index, table):
    from pptx.util import Pt

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, table["title"], table.get("kicker"))

    soft_card(slide, MARGIN, Inches(2.2), Inches(12.0), Inches(1.15), fill=SOFT)
    add_formula(
        slide,
        MARGIN + Inches(0.35),
        Inches(2.35),
        Inches(11.3),
        Inches(0.35),
        table["formula"],
        size=18,
        bold=True,
        color=PRIMARY,
    )
    if table.get("formula_note"):
        add_formula(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.8),
            Inches(11.3),
            Inches(0.3),
            table["formula_note"],
            size=13,
            bold=False,
            color=MUTED,
        )

    headers = table["headers"]
    rows = table["rows"]
    shape = slide.shapes.add_table(
        rows=1 + len(rows),
        cols=len(headers),
        left=MARGIN,
        top=Inches(3.55),
        width=Inches(12.0),
        height=Inches(2.45),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.bold = c == 1
                    run.font.color.rgb = PRIMARY if c == 1 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if table.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.2),
            Inches(12),
            Inches(0.35),
            table["note"],
            size=13,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_norm_example(prs, total, index, example):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, example["title"], example.get("kicker"))

    add_text(
        slide,
        MARGIN,
        Inches(2.25),
        Inches(12),
        Inches(0.35),
        example["lead"],
        size=16,
        bold=True,
        color=INK,
    )

    cards = example.get("cards") or []
    if cards:
        card_w = Inches(5.7)
        gap = Inches(0.4)
        for i, (label, body) in enumerate(cards):
            x = MARGIN + i * (card_w + gap)
            soft_card(slide, x, Inches(2.7), card_w, Inches(1.05), fill=SOFT)
            add_text(
                slide,
                x + Inches(0.3),
                Inches(2.85),
                card_w - Inches(0.6),
                Inches(0.28),
                label,
                size=13,
                bold=True,
                color=PRIMARY,
            )
            add_text(
                slide,
                x + Inches(0.3),
                Inches(3.2),
                card_w - Inches(0.6),
                Inches(0.35),
                body,
                size=12,
                color=MUTED,
            )

    soft_card(slide, MARGIN, Inches(3.95), Inches(12.0), Inches(1.15), fill=SOFT_2)
    add_formula(
        slide,
        MARGIN + Inches(0.35),
        Inches(4.1),
        Inches(11.3),
        Inches(0.35),
        example["formula"],
        size=18,
        bold=True,
        color=PRIMARY,
    )
    if example.get("formula_note"):
        add_formula(
            slide,
            MARGIN + Inches(0.35),
            Inches(4.55),
            Inches(11.3),
            Inches(0.3),
            example["formula_note"],
            size=13,
            bold=False,
            color=MUTED,
        )

    values = example.get("values") or []
    if values:
        n = len(values)
        cell_w = Inches(12.0) / n
        for i, (raw, scaled) in enumerate(values):
            x = MARGIN + i * cell_w
            soft_card(slide, x + Inches(0.05), Inches(5.3), cell_w - Inches(0.1), Inches(1.15), fill=SOFT)
            add_text(
                slide,
                x + Inches(0.1),
                Inches(5.4),
                cell_w - Inches(0.2),
                Inches(0.3),
                raw,
                size=13,
                color=MUTED,
                align=PP_ALIGN.CENTER,
            )
            add_text(
                slide,
                x + Inches(0.1),
                Inches(5.8),
                cell_w - Inches(0.2),
                Inches(0.4),
                scaled,
                size=18,
                bold=True,
                color=PRIMARY,
                align=PP_ALIGN.CENTER,
            )

    content_footer(slide, index, total)


def slide_choosing_follow_up(prs, total, index, follow_up):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, follow_up["title"], follow_up.get("kicker"))

    cards = follow_up.get("cards") or []
    card_w = Inches(5.7)
    gap = Inches(0.4)
    for i, card in enumerate(cards):
        x = MARGIN + i * (card_w + gap)
        soft_card(slide, x, Inches(2.3), card_w, Inches(3.2), fill=SOFT if i == 0 else SOFT_2)
        add_text(
            slide,
            x + Inches(0.35),
            Inches(2.55),
            card_w - Inches(0.7),
            Inches(0.35),
            card["title"],
            size=16,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            x + Inches(0.35),
            Inches(3.1),
            card_w - Inches(0.7),
            Inches(1.4),
            card["body"],
            size=14,
            color=MUTED,
        )
        if card.get("note"):
            add_text(
                slide,
                x + Inches(0.35),
                Inches(4.7),
                card_w - Inches(0.7),
                Inches(0.55),
                card["note"],
                size=13,
                bold=True,
                color=SECONDARY,
            )

    if follow_up.get("bullets"):
        bullets(slide, follow_up["bullets"], top=Inches(5.75), size=15)

    content_footer(slide, index, total)


def slide_method_detail(prs, total, index, detail):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, detail["title"], detail.get("kicker"))

    add_text(
        slide,
        MARGIN,
        Inches(2.3),
        Inches(12),
        Inches(0.45),
        detail["lead"],
        size=16,
        color=MUTED,
    )

    # Formula card
    soft_card(slide, MARGIN, Inches(2.9), Inches(12.0), Inches(1.55), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.1),
        Inches(11.2),
        Inches(0.3),
        "Formula",
        size=12,
        bold=True,
        color=SECONDARY,
    )
    add_formula(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.45),
        Inches(11.2),
        Inches(0.45),
        detail["formula"],
        size=22,
        bold=True,
        color=PRIMARY,
    )
    if detail.get("formula_note"):
        add_formula(
            slide,
            MARGIN + Inches(0.4),
            Inches(3.95),
            Inches(11.2),
            Inches(0.3),
            detail["formula_note"],
            size=13,
            bold=False,
            color=MUTED,
        )

    bullets(slide, detail.get("bullets") or [], top=Inches(4.75), size=16)
    content_footer(slide, index, total)


def slide_topic_methods(prs, total, index, methods):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, methods["title"], methods.get("kicker"))

    cards = methods.get("cards") or []
    if cards:
        card_w = Inches(5.7)
        gap = Inches(0.4)
        for i, (label, body) in enumerate(cards):
            x = MARGIN + i * (card_w + gap)
            soft_card(slide, x, Inches(2.25), card_w, Inches(1.15), fill=SOFT)
            add_text(
                slide,
                x + Inches(0.35),
                Inches(2.4),
                card_w - Inches(0.7),
                Inches(0.3),
                label,
                size=14,
                bold=True,
                color=PRIMARY,
            )
            add_text(
                slide,
                x + Inches(0.35),
                Inches(2.8),
                card_w - Inches(0.7),
                Inches(0.4),
                body,
                size=13,
                color=MUTED,
            )

    method_rows = methods.get("methods") or []
    top = Inches(3.65) if cards else Inches(2.35)
    for i, (name, desc) in enumerate(method_rows):
        y = top + Inches(i * 0.95)
        soft_card(slide, MARGIN, y, Inches(12.0), Inches(0.85), fill=SOFT_2 if i % 2 else SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            y + Inches(0.12),
            Inches(11.3),
            Inches(0.28),
            name,
            size=15,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            MARGIN + Inches(0.35),
            y + Inches(0.42),
            Inches(11.3),
            Inches(0.35),
            desc,
            size=13,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_preprocess_workflow(prs, total, index, workflow):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, workflow["title"], workflow.get("kicker"))

    if workflow.get("lead"):
        add_text(
            slide,
            MARGIN,
            Inches(2.2),
            Inches(12.0),
            Inches(0.35),
            workflow["lead"],
            size=15,
            color=MUTED,
        )

    steps = workflow.get("steps") or []
    if steps:
        card_w = Inches(2.8)
        gap = Inches(0.2)
        top = Inches(2.75)
        for i, (number, label, body) in enumerate(steps):
            x = MARGIN + i * (card_w + gap)
            soft_card(slide, x, top, card_w, Inches(3.2), fill=SOFT)
            add_text(
                slide,
                x + Inches(0.25),
                top + Inches(0.35),
                card_w - Inches(0.5),
                Inches(0.45),
                number,
                size=28,
                bold=True,
                color=PRIMARY,
                align=PP_ALIGN.CENTER,
            )
            add_text(
                slide,
                x + Inches(0.25),
                top + Inches(1.05),
                card_w - Inches(0.5),
                Inches(0.4),
                label,
                size=18,
                bold=True,
                color=INK,
                align=PP_ALIGN.CENTER,
            )
            add_text(
                slide,
                x + Inches(0.25),
                top + Inches(1.65),
                card_w - Inches(0.5),
                Inches(1.2),
                body,
                size=13,
                color=MUTED,
                align=PP_ALIGN.CENTER,
            )

    content_footer(slide, index, total)


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    n = 1
    if content.get("workflow"):
        n += 1
    if content.get("detail"):
        n += 1
    if content.get("table"):
        n += 1
    n += len(content.get("extra_tables") or [])
    if content.get("methods"):
        n += 1
    n += len(content.get("method_details") or [])
    if content.get("choosing"):
        n += 1
        if content["choosing"].get("follow_up"):
            n += 1
    if content.get("norm_example"):
        n += 1
        if content["norm_example"].get("table"):
            n += 1
    if content.get("std_example"):
        n += 1
        if content["std_example"].get("table"):
            n += 1
    if content.get("ordinal"):
        n += 1
        if content["ordinal"].get("comparison"):
            n += 1
        if content["ordinal"].get("howto"):
            n += 1
        if content["ordinal"].get("example_table"):
            n += 1
    if content.get("one_hot"):
        n += 1
        if content["one_hot"].get("comparison_note"):
            n += 1
        if content["one_hot"].get("howto"):
            n += 1
        if content["one_hot"].get("example_table"):
            n += 1
        if content["one_hot"].get("dummy_trap"):
            n += 1
            if content["one_hot"]["dummy_trap"].get("example_table"):
                n += 1
    if content.get("advanced_encoding"):
        adv = content["advanced_encoding"]
        if adv.get("overview_table"):
            n += 1
        n += len(adv.get("extra_tables") or [])
    if content.get("missing_data"):
        n += 1
        if content["missing_data"].get("types_table"):
            n += 1
        if content["missing_data"].get("causes"):
            n += 1
            if content["missing_data"]["causes"].get("table"):
                n += 1
        if content["missing_data"].get("mcar"):
            n += 1
            if content["missing_data"]["mcar"].get("example_table"):
                n += 1
        if content["missing_data"].get("mar"):
            n += 1
            if content["missing_data"]["mar"].get("methods_table"):
                n += 1
            if content["missing_data"]["mar"].get("knn_example"):
                n += 1
                knn = content["missing_data"]["mar"]["knn_example"]
                if knn.get("example_table"):
                    n += 1
                if knn.get("result"):
                    n += 1
            if content["missing_data"]["mar"].get("missforest_example"):
                n += 1
                mf = content["missing_data"]["mar"]["missforest_example"]
                if mf.get("result"):
                    n += 1
                if mf.get("example_table"):
                    n += 1
        if content["missing_data"].get("mnar"):
            n += 1
            if content["missing_data"]["mnar"].get("example_table"):
                n += 1
            if content["missing_data"]["mnar"].get("preserve_pattern_table"):
                n += 1
    return n


def main() -> None:
    s = SESSION
    mp = ML_PROCESS
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 5 + len(mp["steps"]) + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    n = 5
    slide_ml_process_overview(prs, total, n)
    for step in mp["steps"]:
        n += 1
        slide_ml_process_step(prs, total, n, step)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")


if __name__ == "__main__":
    main()
