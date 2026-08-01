#!/usr/bin/env python3
"""
Build Week 1 · Session 3 presentation using official ETRA Design System.
Content is filled incrementally as teaching blocks are pasted in.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week1-Session3-Classification-Basics.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session3-diagrams"

SESSION = {
    "eyebrow": "Week 1  ·  Session 3",
    "section_number": "03",
    "section_tag": "S03",
    "section_title": "Classical ML",
    "subtitle": "Classification basics — from labels to decision boundaries",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "Classification Basics",
    "topics": [
        "Classification Basics",
        "Logistic regression",
        "K-NN",
        "Decision boundaries",
        "Threshold tuning",
    ],
}

TOPIC_CONTENT: dict = {
    "Classification Basics": {
        "title": "Classification Overview",
        "kicker": "Classification Basics",
        "body": "Classification is a supervised learning task used to assign a new observation to a predefined category based on learned patterns from labeled data.",
        "bullets": [
            "Output is a class label (e.g., 0/1, Yes/No).",
            "Used in spam, diagnosis, fraud, and churn prediction.",
        ],
        "plot_path": "classification-overview.png",
        "note": "Features in → class label out. Probability models come next.",
    },
    "Logistic regression": {
        "title": "Logistic Regression: Definition",
        "kicker": "Binary Classification",
        "body": "Supervised binary classifier that estimates P(y=1|x) from features.",
        "bullets": [
            "Outputs probabilities in [0, 1].",
            "Class label comes from applying a threshold later.",
            "Models features → log-odds of the positive class.",
        ],
        "plot_path": "sigmoid-threshold.png",
        "note": "Maps a linear score to a probability via the sigmoid.",
        "extra_slides": [
            {
                "title": "Why Not Linear Regression for Classification?",
                "kicker": "Motivation for Logistic Regression",
                "body": "Linear regression predicts continuous values rather than discrete classes.",
                "bullets": [
                    "Predictions can fall outside [0, 1] — invalid as probabilities.",
                    "No natural, stable class boundary for labels.",
                ],
                "plot_path": "linear-vs-logistic.png",
                "note": "Sigmoid keeps outputs in [0, 1] and supports a clear cutoff.",
            },
            {
                "title": "Sigmoid Mapping and Decision Threshold",
                "kicker": "Sigmoid & Threshold",
                "body": "Sigmoid turns the linear score z into a probability.",
                "formula": "P(y=1|x) = σ(z) = 1 / (1 + e^{−z}),  z = β_{0} + βᵀx",
                "formula_tex": r"P(y=1\mid x)=\sigma(z)=\dfrac{1}{1+e^{-z}},\quad z=\beta_0+\beta^{\top}x",
                "bullets": [
                    "Default cutoff: P ≥ 0.5 → class 1, else class 0.",
                    "τ is tunable for cost of false positives vs false negatives.",
                ],
                "plot_path": "sigmoid-threshold.png",
                "note": "Tune τ to business/clinical risk — covered in Threshold tuning.",
            },
            {
                "title": "Key Assumptions for Logistic Regression",
                "kicker": "Model Assumptions",
                "body": "Works on log-odds — check independence, collinearity, and linearity in logit space.",
                "formula": "Odds = p / (1 − p),  Logit(p) = log(p / (1 − p))",
                "formula_tex": r"\mathrm{Odds}=\dfrac{p}{1-p},\quad \mathrm{Logit}(p)=\log\!\left(\dfrac{p}{1-p}\right)",
                "layout": "formula_example",
                "bullets": [
                    "Observations should be independent.",
                    "Avoid severe multicollinearity among predictors.",
                    "Relationship ≈ linear in the log-odds space.",
                ],
            },
            {
                "title": "Maximum Likelihood Estimation (MLE)",
                "kicker": "How Logistic Regression Learns",
                "body": "Parameters β are chosen to maximize the likelihood of the observed labels.",
                "bullets": [
                    "Each curve is a different β (different likelihood).",
                    "Best fit: highest likelihood / lowest log-loss.",
                ],
                "plot_path": "mle-sigmoid-curves.png",
                "note": "Pick β that makes the observed labels most probable.",
            },
            {
                "title": "Logistic Regression: Strengths and Limits",
                "kicker": "When to Use It",
                "body": "Strong baseline when classes are roughly linearly separable; weaker on complex shapes.",
                "bullets": [
                    "Strengths: interpretable odds, fast training, calibrated probabilities.",
                    "Limits: linear boundary; struggles with overlap / non-linear shapes.",
                ],
                "plot_path": "logreg-strength-boundary.png",
                "note": "For wiggly boundaries → polynomials, K-NN, or kernels.",
            },
            {
                "title": "Multiclass Extension of Logistic Regression",
                "kicker": "Beyond Binary",
                "body": "Extend binary logistic regression to 3+ classes with Softmax or One-vs-All.",
                "bullets": [
                    "Softmax: one multinomial model for all classes.",
                    "One-vs-All: one binary classifier per class; pick max P.",
                ],
                "plot_path": "multiclass-logistic.png",
                "note": "Final label = argmax_k P(y=k|x).",
            },
        ],
    },
    "K-NN": {
        "title": "K-Nearest Neighbors (K-NN): Core Idea",
        "kicker": "Instance-Based Learning",
        "body": "Non-parametric classifier: label by majority vote among the K nearest training points.",
        "bullets": [
            "Choose K and a distance metric.",
            "Find K nearest neighbors; majority vote wins.",
        ],
        "plot_path": "knn-core-idea.png",
        "note": "Lazy learner — the training set is the model.",
        "extra_slides": [
            {
                "title": "K-NN: Choosing K and Distance Metric",
                "kicker": "Hyperparameters",
                "body": "K and distance control how smooth the decision boundary is.",
                "bullets": [
                    "Small K: flexible, noise-sensitive.",
                    "Large K: smoother; may miss local structure.",
                ],
                "plot_path": "knn-k-tradeoff.png",
                "note": "Pick K with CV · scale features before distance modeling.",
            },
            {
                "title": "K-NN: Strengths and Limitations",
                "kicker": "When to Use It",
                "layout": "table",
                "table": {
                    "headers": ["Strengths", "Limitations"],
                    "rows": [
                        [
                            "Simple and intuitive",
                            "Slow prediction on large datasets",
                        ],
                        [
                            "Works well on small, separated data",
                            "Sensitive to irrelevant features / imbalance",
                        ],
                        [
                            "Native multiclass support",
                            "Curse of dimensionality",
                        ],
                        [
                            "No training phase (lazy baseline)",
                            "Needs feature scaling",
                        ],
                    ],
                },
                "note": "Great baseline; watch cost at inference and high dimensions.",
            },
        ],
    },
    "Decision boundaries": {
        "title": "Decision Boundaries Across Models",
        "kicker": "Geometry of Classification",
        "body": "A decision boundary is the surface where the predicted class flips — its shape depends on the model family.",
        "bullets": [
            "Logistic: linear boundary in feature space (unless you add features).",
            "K-NN: local, piecewise boundary that follows nearest neighbors.",
            "Kernel / RBF-like models: smooth curved boundaries.",
        ],
        "plot_path": "decision-boundaries.png",
        "note": "Match model flexibility to how classes are separated in the data.",
    },
    "Threshold tuning": {
        "title": "Threshold Tuning: Trading FP vs FN",
        "kicker": "Operating Point",
        "body": "The sigmoid gives P(y=1|x); τ turns that probability into a hard class label.",
        "bullets": [
            "Lower τ → more positives (↑ recall, ↑ false positives).",
            "Higher τ → fewer positives (↑ precision, ↑ false negatives).",
            "Default τ = 0.5 is a starting point, not a law.",
        ],
        "plot_path": "threshold-tuning.png",
        "note": "Set τ from costs: missed fraud ≠ false spam alarm.",
    },
}

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "After regression — learn models that predict class labels.",
    "current": "S3",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.15), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.65),
        Inches(11.5),
        Inches(0.9),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.7), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.05), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    # Two rows of topic chips to avoid overlap on long labels
    topics = s["topics"]
    mid = (len(topics) + 1) // 2
    rows = [topics[:mid], topics[mid:]]
    for ri, row in enumerate(rows):
        if not row:
            continue
        y = 5.05 + ri * 0.55
        gap = 0.12
        usable = 12.0
        w = min(2.9, (usable - gap * (len(row) - 1)) / max(len(row), 1))
        total_w = len(row) * w + (len(row) - 1) * gap
        x0 = MARGIN.inches + (usable - total_w) / 2
        for i, topic in enumerate(row):
            x = Inches(x0 + i * (w + gap))
            soft_card(slide, x, Inches(y), Inches(w), Inches(0.42), fill=SOFT)
            add_text(
                slide,
                x,
                Inches(y + 0.05),
                Inches(w),
                Inches(0.32),
                topic,
                size=10,
                color=PRIMARY,
                align=PP_ALIGN.CENTER,
            )

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    """Fit a plot PNG inside a box (same idea as add_diagram)."""
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return None
    aspect = px_w / px_h
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(path),
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    pitch = 0.50 if dense else 0.62
    bsize = 14 if dense else 16

    if content.get("body"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.45),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(3.6)
    else:
        bullet_top = Inches(2.35)
        bsize = 15 if dense else 17

    # Keep bullets above the note/footer band
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    """Definition + formula + optional scatter/line plot."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)
    body_h = 0.95 if has_plot else 1.15

    soft_card(slide, MARGIN, Inches(2.2), col_w, Inches(body_h), fill=SOFT)
    add_text(
        slide,
        Inches(MARGIN.inches + 0.3),
        Inches(2.32),
        text_w,
        Inches(body_h - 0.2),
        content.get("body") or "",
        size=12 if has_plot else 13,
        color=INK,
    )

    # Fold formula_note into footer note to avoid colliding with formula PNG
    footer_note = content.get("note") or ""
    if content.get("formula_note"):
        fn = content["formula_note"]
        footer_note = f"{fn}  ·  {footer_note}" if footer_note else fn
    has_note = bool(footer_note)

    if has_formula:
        formula_top = 2.2 + body_h + 0.1
        formula_card_h = 1.35 if is_frac else 1.1
        formula_box_h = 0.85 if is_frac else 0.55
        soft_card(slide, MARGIN, Inches(formula_top), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(formula_top + 0.08),
            text_w,
            Inches(0.2),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(formula_top + 0.28),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        bullet_top = formula_top + formula_card_h + 0.1
        pitch = 0.42 if (has_note or len(items) >= 2) else 0.52
        bsize = 12
    else:
        bullet_top = 2.2 + body_h + 0.15
        pitch = 0.48 if (has_note or len(items) >= 3) else 0.56
        bsize = 13

    note_band = 6.4 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    if has_plot:
        max_items = min(max_items, 2 if has_formula else 3)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        left=MARGIN,
        width=bullet_w,
        pitch=pitch,
        item_height=Inches(0.4),
    )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if footer_note:
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            footer_note,
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    """Full-width practical example with formula and interpretation."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(2.15)
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.28),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(3.0)
    else:
        formula_top = Inches(2.25)

    has_fnote = bool(content.get("formula_note"))
    formula_h = 1.35 if is_frac else 1.1
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.15),
        Inches(11.2),
        Inches(1.0 if is_frac else 0.75),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )

    # formula_note sits BELOW the formula card (never inside over the PNG)
    cursor = formula_top.inches + formula_h + 0.08
    if has_fnote:
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(cursor),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        cursor += 0.32

    bullet_top = cursor + 0.06
    pitch = 0.42 if dense else 0.52
    bsize = 12 if dense else 14
    note_band = 6.4 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    max_items = min(max_items, 3)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(0.4),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    # Prefer kicker over a long body when the table is dense
    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = 2.25
    if show_body:
        soft_card(slide, MARGIN, Inches(2.15), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.25),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = 2.85

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    # Keep ≥0.15" gap between table bottom and note
    note_top = table_top + table_h + 0.15
    render_note = has_note and (note_top + 0.4) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    """Title + optional short body + full-width diagram."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    diagram_top = Inches(2.15)
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(2.1),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(2.45)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(4.35), fill=SOFT)
    _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + 0.15),
        Inches(11.6),
        Inches(4.05),
        folder=folder,
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.9),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    """Returns number of slides added."""
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(4.25),
        Inches(11.2),
        Inches(0.4),
        "Content coming next — paste the teaching block when ready.",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 4 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
