#!/usr/bin/env python3
"""
Build Week 2 · Session 1 presentation using official ETRA Design System.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week2-Session1-Deep-Learning.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session-w2s1-diagrams"

SESSION = {
    "eyebrow": "Week 2  ·  Session 1",
    "section_number": "07",
    "section_tag": "S07",
    "section_title": "Deep Learning",
    "subtitle": "Neural networks, backpropagation, CNNs, and RNNs",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "Neural Networks · CNNs · RNNs",
    "topics": [
        "Deep Learning",
        "Neural Networks",
        "Perceptron & MLP",
        "Backpropagation",
        "CNNs",
        "RNNs & regularization",
    ],
}

TOPIC_CONTENT: dict = {
    "Deep Learning": {
        "title": "Deep Learning: Introduction",
        "kicker": "Foundations → Optimization → Architectures → Autoencoders",
        "body": "Each phase answers a different question.",
        "bullets": [
            "Four phases: foundations → optimization → architectures → autoencoders.",
            "Each phase has a lab-style checkpoint — pause at phase boundaries.",
            "Focus on deep learning basics — phase breaks are intentional.",
        ],
        "plot_path": "ann-intro.png",
        "note": "Phase 1 focus: understand the “why” and the mathematical mechanics of a single neuron before scaling to layers and data.",
        "extra_slides": [
            {
                "title": "Deep Learning: Introduction",
                "kicker": "Four phases, four questions",
                "body": "Foundations → optimization → architectures → autoencoders. Pause at each boundary.",
                "layout": "diagram",
                "plot_path": "four-phases.png",
                "note": "Phase 1 focus: understand the “why” and the mathematical mechanics of a single neuron before scaling to layers and data.",
            },
            {
                "title": "Phase 1: Foundations & Neural Core",
                "kicker": "Why depth, neurons, and forward pass",
                "body": "Goal: understand what a neuron computes and why stacking layers enables hierarchical features.",
                "formula": "z = w·x + b",
                "formula_tex": r"z = w \cdot x + b",
                "formula_note": "Exit check: explain this — and why nonlinearity is required between layers.",
                "layout": "formula_example",
                "bullets": [
                    "Topics: perceptron, activations, MLP, forward propagation, and loss functions.",
                    "Stacking layers builds hierarchical features only if a nonlinearity sits between them.",
                ],
            },
            {
                "title": "Phase 1 Roadmap",
                "kicker": "Why depth, neurons, and forward pass",
                "layout": "table",
                "table": {
                    "headers": ["Block", "Question it answers"],
                    "rows": [
                        ["Perceptron & MLP", "What does one layer compute?"],
                        ["Forward pass", "How does inference flow through the graph?"],
                        ["Loss & gradients", "What are we optimizing?"],
                    ],
                },
                "note": "Exit check: explain z = w·x + b and why nonlinearity is required between layers.",
            },
            {
                "title": "What Is Deep Learning?",
                "kicker": "Definition in One Sentence — Then the Details",
                "body": "Deep learning learns a parameterized mapping from inputs to outputs by stacking nonlinear transformations (layers).",
                "bullets": [
                    "“Deep” means multiple such transformations — hierarchical features, not one giant linear rule.",
                ],
                "plot_path": "ann-intro.png",
                "note": "Layers rewrite features step by step. Depth is composition, not one giant linear rule.",
            },
            {
                "title": "What Is Deep Learning?",
                "kicker": "Definition in One Sentence — Then the Details",
                "bullets": [
                    "Training = pick a loss that measures mistakes, then adjust weights with gradients so average loss drops on data (generalization to unseen data is the real goal).",
                    "Unlike fixed pipelines of hand-crafted features, the network learns representations that are tuned to the task (edges → textures → parts → objects in vision, etc.).",
                ],
                "note": "Generalization to unseen data is the real goal — not memorizing the training set.",
            },
            {
                "title": "Three Ingredients You Always Have",
                "kicker": "Architecture · Loss · Optimizer + data",
                "layout": "table",
                "table": {
                    "headers": ["Ingredient", "What it means"],
                    "rows": [
                        [
                            "Architecture",
                            "How signals flow (MLP, CNN, RNN, Transformer, …)",
                        ],
                        [
                            "Loss / objective",
                            "What “better” means for your task (MSE, cross-entropy, …)",
                        ],
                        [
                            "Optimizer + data",
                            "How you search parameters and what examples you show",
                        ],
                    ],
                },
                "note": "You always have these three. Change one, and the model’s behavior changes with it.",
            },
            {
                "title": "Three Ingredients You Always Have",
                "kicker": "Architecture · Loss · Optimizer + data",
                "layout": "diagram",
                "plot_path": "three-ingredients.png",
                "note": "A strong architecture with the wrong loss still learns the wrong thing.",
            },
            {
                "title": "AI, Machine Learning, and Deep Learning",
                "kicker": "Subset Relationships — Not Interchangeable Labels",
                "bullets": [
                    "Artificial intelligence (AI): the broad field—rules, search, optimization, learning, planning—anything that aims at intelligent behavior.",
                    "Machine learning (ML): a subfield of AI where behavior improves from data via a learning algorithm (not only hand-written rules).",
                    "Deep learning (DL): a subfield of ML using deep neural networks—multiple learned nonlinear layers; representation learning is central.",
                ],
            },
            {
                "title": "AI, Machine Learning, and Deep Learning",
                "kicker": "Subset Relationships — Not Interchangeable Labels",
                "layout": "diagram",
                "plot_path": "ai-ml-dl.png",
                "note": "DL ⊂ ML ⊂ AI",
            },
            {
                "title": "Containment (conceptual)",
                "kicker": "Subset Relationships — Not Interchangeable Labels",
                "layout": "table",
                "table": {
                    "headers": ["Set", "Contains"],
                    "rows": [
                        ["AI", "Rule systems, classical search, ML, knowledge bases, …"],
                        ["ML", "Linear models, trees, kernel SVMs, shallow nets, deep nets"],
                        ["DL", "CNNs, RNNs/GRU/LSTM, Transformers, GANs, VAEs, …"],
                    ],
                },
                "note": "Saying “we used AI” is vague; “we used supervised DL with a CNN” is precise.",
            },
            {
                "title": "Use-Case Fit: When Deep Learning Shines — and When to Skip It",
                "kicker": "Match Method to Data, Risk, and Budget",
                "bullets": [
                    "DL is strongest when: large data (or strong pre-training), high-dimensional raw inputs (images, text, speech), and compositionality you want learned end-to-end.",
                    "Think twice when: tiny tabular datasets, strict interpretability requirements, tight latency on CPU without optimization, or simple baselines already solve the problem.",
                    "Always compare against strong baselines (logistic regression, gradient boosting, k-NN on engineered features)—DL must earn its complexity on validation metrics you care about.",
                ],
            },
            {
                "title": "Use-Case Fit: When Deep Learning Shines — and When to Skip It",
                "kicker": "Match Method to Data, Risk, and Budget",
                "layout": "table",
                "table": {
                    "headers": ["Situation", "Typical first choice", "Why"],
                    "rows": [
                        [
                            "Large labeled images / audio / text",
                            "Pre-trained DL + fine-tune",
                            "Raw inputs; hierarchical features",
                        ],
                        [
                            "Small tabular (< few k rows)",
                            "Trees / linear + regularization",
                            "Lower variance, fast iteration, interpretability",
                        ],
                        [
                            "Need causal or legal explanation",
                            "Simpler model or hybrid",
                            "Auditability and stability",
                        ],
                        [
                            "Plenty of GPUs + big data",
                            "Larger nets + careful tuning",
                            "Capacity + regularization can scale",
                        ],
                        [
                            "Edge device, millisecond latency",
                            "Small distilled / quantized models",
                            "Memory and power limits",
                        ],
                    ],
                },
            },
        ],
    },
    "Neural Networks": {
        "title": "Why Depth Works",
        "kicker": "Hierarchical Representation Learning",
        "body": "Each layer rewrites the previous layer’s features into a more useful space.",
        "bullets": [
            "Early layers detect simple patterns (edges, local tokens).",
            "Middle layers compose those into textures, parts, phrases.",
            "Late layers assemble task semantics (objects, intent, class scores).",
        ],
        "plot_path": "hierarchical-features.png",
        "note": "Depth is composition. Width is how many features you keep at each level.",
        "extra_slides": [
            {
                "title": "Session Roadmap",
                "kicker": "What we build today",
                "layout": "table",
                "table": {
                    "headers": ["Block", "Question it answers"],
                    "rows": [
                        ["Perceptron & MLP", "What does one layer compute?"],
                        ["Backpropagation", "How do we update the weights?"],
                        ["CNNs", "How do we exploit spatial structure?"],
                        ["RNNs & regularization", "How do we model sequences and generalize?"],
                    ],
                },
                "note": "Exit check: explain z = w·x + b and why nonlinearity is required between layers.",
            },
        ],
    },
    "Perceptron & MLP": {
        "title": "Single Neuron (Perceptron): Math, Weights, and Bias",
        "kicker": "One Affine Map + One Nonlinearity",
        "body": "One weighted sum plus bias, then an activation. The output can be used as-is or passed onward.",
        "formula": "z = wᵀx + b,   a = f(z),   ŷ = a (or further processing)",
        "formula_tex": r"z=w^{\top}x+b,\quad a=f(z),\quad \hat{y}=a",
        "formula_note": "ŷ = a (or further processing)",
        "plot_path": "perceptron.png",
        "extra_slides": [
            {
                "title": "Single Neuron (Perceptron): Math, Weights, and Bias",
                "kicker": "One Affine Map + One Nonlinearity",
                "bullets": [
                    "Affine part (see formula above): z is the weighted sum of inputs plus b—each xᵢ is multiplied by wᵢ.",
                    "The bias b shifts the decision boundary (same orientation, normal w); it does not rotate the separating hyperplane.",
                    "Weights encode which patterns excite the neuron; training adjusts w and b so that on labeled data the loss decreases.",
                    "With no hidden layer and a monotonic activation, the decision boundary in input space is linear (a hyperplane)—powerful but limited.",
                ],
            },
            {
                "title": "Symbols at a Glance",
                "kicker": "One Affine Map + One Nonlinearity",
                "layout": "table",
                "table": {
                    "headers": ["Symbol", "Shape intuition", "Role"],
                    "rows": [
                        ["x", "Feature vector (length d)", "One training example’s inputs"],
                        ["w", "Same length as x", "Direction + strength of sensitivity to each feature"],
                        ["b", "Scalar", "Offset / threshold in logit space"],
                        ["f", "Maps ℝ → ℝ (per neuron)", "Introduces nonlinearity (or identity)"],
                    ],
                },
                "note": "Bias shifts the hyperplane; it does not rotate it. Rotation comes from w.",
            },
            {
                "title": "Weight Initialization: Why Zeros Break Training",
                "kicker": "Symmetry Must Be Broken",
                "bullets": [
                    "If all weights in a layer start equal (especially zero), every hidden unit receives the same gradient and stays identical forever—symmetry never breaks.",
                    "Random initialization makes each unit walk a different path in parameter space so they can specialize as feature detectors.",
                    "Frameworks use He (Kaiming) for ReLU and Xavier (Glorot) for Tanh to keep activation scale stable.",
                ],
            },
            {
                "title": "Weight Initialization: Why Zeros Break Training",
                "kicker": "Symmetry Must Be Broken",
                "body": "He (Kaiming) for ReLU: variance = 2 / n_in. Xavier (Glorot) for Tanh: variance = 2 / (n_in + n_out).",
                "formula": "He: Var = 2/n_in     Xavier: Var = 2/(n_in + n_out)",
                "formula_tex": r"\mathrm{He:}\ \mathrm{Var}=\dfrac{2}{n_{\mathrm{in}}}\qquad \mathrm{Xavier:}\ \mathrm{Var}=\dfrac{2}{n_{\mathrm{in}}+n_{\mathrm{out}}}",
                "formula_note": "Keep activation scale stable across layers",
                "layout": "formula_example",
                "bullets": [
                    "He matches ReLU, which zeros out about half the signals.",
                    "Xavier matches Sigmoid / Tanh, which are symmetric around zero.",
                ],
            },
            {
                "title": "Weight Initialization: Why Zeros Break Training",
                "kicker": "Symmetry Must Be Broken",
                "layout": "table",
                "table": {
                    "headers": ["Init scheme", "Problem it avoids", "Best suited for"],
                    "rows": [
                        [
                            "All zeros",
                            "Hidden units stay identical; no representation diversity",
                            "None (always avoid for weights)",
                        ],
                        [
                            "Xavier (Glorot)",
                            "Vanishing or exploding variance across layers",
                            "Sigmoid / Tanh activations",
                        ],
                        [
                            "He (Kaiming)",
                            "Variance drop due to ReLU zeroing out half the signals",
                            "ReLU / Leaky ReLU activations",
                        ],
                    ],
                },
            },
            {
                "title": "Activation Functions: Sigmoid, Tanh, and ReLU",
                "kicker": "Nonlinearity is the Key to Network Depth",
                "body": "Stacked linear layers collapse to one linear map. Activations prevent that collapse.",
                "formula": "σ(z) = 1/(1+e^{−z}),   tanh(z) = (e^z − e^{−z})/(e^z + e^{−z}),   ReLU(z) = max(0, z)",
                "formula_tex": r"\sigma(z)=\dfrac{1}{1+e^{-z}},\quad \tanh(z)=\dfrac{e^{z}-e^{-z}}{e^{z}+e^{-z}},\quad \mathrm{ReLU}(z)=\max(0,z)",
                "layout": "formula_example",
            },
            {
                "title": "Activation Functions: Sigmoid, Tanh, and ReLU",
                "kicker": "Nonlinearity is the Key to Network Depth",
                "layout": "diagram",
                "plot_path": "sigmoid-tanh.png",
                "note": "f(x)=1/(1+e^{−5x}) vs g(x)=1/(1+e^{−10x})  ·  tanh(x) is zero-centered, range (−1, 1)",
            },
            {
                "title": "Activation Functions: Sigmoid, Tanh, and ReLU",
                "kicker": "Nonlinearity is the Key to Network Depth",
                "bullets": [
                    "Sigmoid: maps inputs to (0, 1)—useful for probabilities, but suffers from vanishing gradients when z is very large or very small (saturation).",
                    "Tanh: zero-centered, maps inputs to (-1, 1)—often nicer than sigmoid for hidden layers but still saturates.",
                    "ReLU: modern default for deep stacks—cheap to compute, avoids vanishing gradient on positive domain; leaky ReLU / ELU are variants.",
                ],
            },
            {
                "title": "Activation Functions: Sigmoid, Tanh, and ReLU",
                "kicker": "Nonlinearity is the Key to Network Depth",
                "layout": "table",
                "table": {
                    "headers": ["Function", "Range", "Primary use-case"],
                    "rows": [
                        [
                            "Sigmoid",
                            "(0, 1)",
                            "Binary classification output layer, gating networks",
                        ],
                        [
                            "Tanh",
                            "(-1, 1)",
                            "Hidden layers (recurrent networks)",
                        ],
                        [
                            "ReLU",
                            "[0, ∞)",
                            "Default choice for hidden layers in MLPs and CNNs",
                        ],
                    ],
                },
            },
            {
                "title": "Sigmoid & Threshold",
                "kicker": "Probability cutoff for classification",
                "body": "Decision threshold τ = 0.5 splits class 0 from class 1 on the sigmoid probability.",
                "formula": "τ = 0.5",
                "formula_tex": r"\tau = 0.5",
                "formula_note": "Predict class 1 when σ(z) ≥ τ",
                "plot_path": "sigmoid-threshold.png",
            },
            {
                "title": "Activation Curves",
                "kicker": "Nonlinearity is the Key to Network Depth",
                "layout": "diagram",
                "plot_path": "activations.png",
                "note": "Flat tails (sigmoid / tanh) → tiny gradients. ReLU is linear on the positive side.",
            },
            {
                "title": "The XOR Problem: Why Hidden Layers Are Logically Necessary",
                "kicker": "Why a single line cannot solve XOR",
                "bullets": [
                    "XOR is not linearly separable in the 2D input plane: no single line separates class 1 from class 0 on the four corners.",
                    "A single perceptron computes a half-space; XOR requires nonlinear mixing of inputs—exactly what a hidden layer provides.",
                    "Hidden units reparameterize inputs into a space where the task becomes linearly separable—this is the core intuition for representation learning.",
                ],
            },
            {
                "title": "The XOR Problem: Why Hidden Layers Are Logically Necessary",
                "kicker": "Why a single line cannot solve XOR",
                "layout": "diagram",
                "plot_path": "xor-separability.png",
                "note": "A hidden layer remaps the inputs so the output layer can cut with a line.",
            },
            {
                "title": "XOR Truth Table",
                "kicker": "Four Points, One Nonlinear Rule",
                "layout": "table",
                "table": {
                    "headers": ["x₁", "x₂", "XOR"],
                    "rows": [
                        ["0", "0", "0"],
                        ["0", "1", "1"],
                        ["1", "0", "1"],
                        ["1", "1", "0"],
                    ],
                },
                "note": "Hidden units reparameterize the inputs so XOR becomes linearly separable.",
            },
            {
                "title": "From Perceptron to MLP (Multi-Layer Perceptron)",
                "kicker": "Depth = Composed Nonlinear Features",
                "body": "Hidden layer: affine map + nonlinearity → features h. Output: another affine map + τ.",
                "formula": "h = φ(W₁x + b₁),   ŷ = τ(W₂h + b₂)",
                "formula_tex": r"h=\varphi(W_1 x+b_1),\quad \hat{y}=\tau(W_2 h+b_2)",
                "formula_note": "φ = hidden activation  ·  τ = output activation",
                "plot_path": "mlp-architecture.png",
            },
            {
                "title": "From Perceptron to MLP (Multi-Layer Perceptron)",
                "kicker": "Depth = Composed Nonlinear Features",
                "bullets": [
                    "Each hidden layer applies an affine map followed by a nonlinearity, producing new features h as functions of the original inputs.",
                    "Width (units per layer) and depth (number of layers) trade off capacity, data needs, compute, and optimization difficulty.",
                    "Universal approximation (broad strokes): one sufficiently wide hidden layer can approximate many continuous functions on compact domains—in practice depth + inductive bias wins.",
                ],
            },
            {
                "title": "From Perceptron to MLP (Multi-Layer Perceptron)",
                "kicker": "Depth = Composed Nonlinear Features",
                "layout": "table",
                "table": {
                    "headers": ["Hyperparameter", "What it controls"],
                    "rows": [
                        [
                            "Hidden width",
                            "How many independent nonlinear features per layer",
                        ],
                        [
                            "Depth",
                            "How many levels of composition (edges→textures→…)",
                        ],
                        [
                            "Skip / residual",
                            "Ease of optimization in very deep nets",
                        ],
                    ],
                },
            },
            {
                "title": "Forward Propagation: The Engine of Inference",
                "kicker": "Neural Network",
                "layout": "diagram",
                "plot_path": "nn-forward.png",
                "note": "Edge thickness reflects weight magnitude.",
            },
            {
                "title": "Forward Propagation: The Engine of Inference",
                "kicker": "The Engine of Inference",
                "bullets": [
                    "Forward pass: feed input x through each layer in order—compute logits or scores at the end (then softmax for class probabilities if needed).",
                    "Implementation-wise this is a computational graph: each op (matmul, ReLU, …) knows how to propagate values forward.",
                    "Training repeats forward passes on mini-batches for efficiency and gradient noise that can help generalization.",
                ],
            },
            {
                "title": "Forward Propagation: The Engine of Inference",
                "kicker": "The Engine of Inference",
                "layout": "table",
                "table": {
                    "headers": ["Step", "Operation", "Output carries…"],
                    "rows": [
                        ["1", "Linear: z = Wx + b", "Weighted sum at each unit"],
                        ["2", "Activation: a = f(z)", "Nonlinear feature for next layer"],
                        ["3", "Repeat / head", "Class logits or regression value"],
                    ],
                },
            },
            {
                "title": "Forward Propagation: The Engine of Inference",
                "kicker": "Linear → Activation → Head",
                "layout": "diagram",
                "plot_path": "forward-pass.png",
                "note": "Training repeats this on mini-batches. The graph also knows how to send gradients backward.",
            },
        ],
    },
    "Backpropagation": {
        "title": "Loss Functions: MSE vs Cross-Entropy",
        "kicker": "Match the Loss to the Output and Task",
        "body": "Regression — Mean Squared Error (MSE)",
        "formula": "L = (1/N) Σ ||ŷᵢ − yᵢ||²",
        "formula_tex": r"\mathcal{L}=\dfrac{1}{N}\sum_i \|\hat{y}_i - y_i\|^2",
        "layout": "formula_example",
        "bullets": [
            "Penalizes large errors heavily (quadratic); differentiable everywhere—classic for continuous targets.",
            "Assumes Gaussian-like noise in a probabilistic story; outliers can dominate—Huber loss is a robust alternative.",
        ],
        "extra_slides": [
            {
                "title": "Loss Functions: MSE vs Cross-Entropy",
                "kicker": "Match the Loss to the Output and Task",
                "body": "Classification — Cross-entropy",
                "formula": "L = −(1/N) Σᵢ Σ_c y_ic log p̂_ic",
                "formula_tex": r"\mathcal{L}=-\dfrac{1}{N}\sum_i\sum_c y_{ic}\log\hat{p}_{ic}",
                "formula_note": "Often with logits + log-softmax for stability",
                "layout": "formula_example",
                "bullets": [
                    "Compares predicted probabilities p̂ to the true distribution (one-hot or soft labels).",
                    "Pairs naturally with softmax at the output; gradients are well-behaved when implemented in log-space.",
                ],
            },
            {
                "title": "Loss Functions: MSE vs Cross-Entropy",
                "kicker": "Match the Loss to the Output and Task",
                "layout": "table",
                "table": {
                    "headers": ["Task", "Common head + loss"],
                    "rows": [
                        ["Regression", "Linear output + MSE / Huber"],
                        ["Binary classification", "Sigmoid + binary cross-entropy"],
                        ["Multi-class", "Softmax + cross-entropy"],
                    ],
                },
                "note": "Often with logits + log-softmax for numerical stability.",
            },
            {
                "title": "Backpropagation: Chain Rule on the Computational Graph",
                "kicker": "Chain along paths",
                "formula": "∂L/∂w = (∂L/∂z) (∂z/∂w)",
                "formula_tex": r"\dfrac{\partial \mathcal{L}}{\partial w}=\dfrac{\partial \mathcal{L}}{\partial z}\,\dfrac{\partial z}{\partial w}",
                "formula_note": "Chain along paths",
                "layout": "formula_example",
                "bullets": [
                    "Backprop applies the chain rule systematically: compute adjoints ∂L/∂ each intermediate tensor from outputs backward.",
                ],
            },
            {
                "title": "Backpropagation: Chain Rule on the Computational Graph",
                "kicker": "Chain along paths",
                "layout": "diagram",
                "plot_path": "backprop-flow.png",
                "note": "Forward computes predictions. Backward carries ∂L/∂w along every path.",
            },
            {
                "title": "Backpropagation: Chain Rule on the Computational Graph",
                "kicker": "Chain along paths",
                "bullets": [
                    "Reverse-mode autodiff (backprop) is efficient when many parameters feed one scalar loss—exactly the neural net case.",
                    "Modern frameworks (PyTorch, JAX, TensorFlow) build the graph (explicit or traced) and implement this reliably—you still debug shapes, numerical stability, and vanishing gradients.",
                ],
                "note": "Understanding backprop = understanding which paths carry gradient and which activations saturate.",
            },
            {
                "title": "Backpropagation: Chain Rule on the Computational Graph",
                "kicker": "Chain along paths",
                "layout": "table",
                "table": {
                    "headers": ["Idea", "Meaning"],
                    "rows": [
                        [
                            "Shared subexpressions",
                            "Same forward values reused; backward visits each edge once",
                        ],
                        [
                            "Topological order",
                            "Visit nodes after their successors—dynamic programming on the graph",
                        ],
                    ],
                },
                "note": "Understanding backprop = understanding which paths carry gradient and which activations saturate.",
            },
            {
                "title": "Phase 2: Optimization & Training Strategy",
                "kicker": "Stable Optimization + Generalization",
                "bullets": [
                    "Phase 2 is about the art of training: diagnosing fit, controlling gradients, picking optimizers, normalizing, and regularizing.",
                    "Good models are not only “big”—they match capacity to data, monitor validation faithfully, and fail visibly when something is wrong.",
                ],
            },
            {
                "title": "Phase 2: Optimization & Training Strategy",
                "kicker": "What you will control in practice",
                "layout": "table",
                "table": {
                    "headers": ["Layer", "Examples"],
                    "rows": [
                        ["Optimization", "SGD, Adam, gradient clipping"],
                        ["Stabilization", "Batch Norm, residual paths, initialization"],
                        ["Generalization", "Dropout, weight decay, early stopping"],
                    ],
                },
            },
            {
                "title": "Training Diagnostics: Underfitting, Good Fit, and Overfitting",
                "kicker": "Read Train vs Validation Behavior",
                "bullets": [
                    "Underfitting: both train and validation error stay high—the model is too simple, features are insufficient, or optimization is stuck (LR too low, wrong loss).",
                    "Good fit: train and validation errors are both low and track each other; small gap is normal if train is slightly better.",
                    "Overfitting: train error keeps dropping while validation worsens or plateaus badly—the model memorizes idiosyncrasies (noise, augment leakage).",
                ],
            },
            {
                "title": "Training Diagnostics: Underfitting, Good Fit, and Overfitting",
                "kicker": "Read Train vs Validation Behavior",
                "layout": "diagram",
                "plot_path": "train-val-curves.png",
                "note": "A small train–val gap is normal. A widening gap is overfitting.",
            },
            {
                "title": "Training Diagnostics: Underfitting, Good Fit, and Overfitting",
                "kicker": "Read Train vs Validation Behavior",
                "layout": "table",
                "table": {
                    "headers": ["Regime", "Train loss", "Val loss / metric", "Typical response"],
                    "rows": [
                        [
                            "Underfitting",
                            "High",
                            "High",
                            "More capacity, better features, longer train, higher LR (carefully)",
                        ],
                        [
                            "Good fit",
                            "Low",
                            "Low, close to train",
                            "Keep regularization; ship or iterate on data",
                        ],
                        [
                            "Overfitting",
                            "Very low",
                            "Worse / diverging gap",
                            "More data, dropout, wd, simpler model, early stop",
                        ],
                    ],
                },
            },
            {
                "title": "Vanishing and Exploding Gradients",
                "kicker": "Depth Multiplies Jacobians Layer by Layer",
                "bullets": [
                    "Vanishing: backprop multiplies many factors < 1 (saturated sigmoid/tanh, poor init)—early layers get tiny updates and learn slowly.",
                    "Exploding: repeated factors > 1 or large weights produce huge parameter updates—loss spikes, NaNs.",
                    "Fixes are layered: architecture (residual paths), activations (ReLU), initialization (He/Xavier), normalization (BN/LN), and optimization (clipping).",
                ],
            },
            {
                "title": "Vanishing and Exploding Gradients",
                "kicker": "Depth Multiplies Jacobians Layer by Layer",
                "layout": "table",
                "table": {
                    "headers": ["Symptom", "Likely cause", "Mitigation"],
                    "rows": [
                        [
                            "Early layers barely move",
                            "Vanishing chain through saturated units",
                            "ReLU, He init, residuals, norm",
                        ],
                        [
                            "Loss → NaN quickly",
                            "Exploding activations or LR too high",
                            "Clip grads, lower LR, check loss scaling",
                        ],
                        [
                            "Stable then sudden blow-up",
                            "Rare but real—mixed precision, bad data batch",
                            "Finite checks, gradient clipping",
                        ],
                    ],
                },
            },
            {
                "title": "Optimizers: SGD and Adam",
                "kicker": "Learning rate, Momentum, and Adaptive Steps",
                "bullets": [
                    "SGD + momentum: accumulates updates in a velocity vector to dampen oscillations in ravines and accelerate along consistent directions.",
                    "Adam: tracks both first-moment (momentum) and second-moment (recent squared gradients) of gradients to adapt learning rates per parameter.",
                    "SGD is often preferred in computer vision for generalization, while Adam is the default starting point for Transformers and general MLPs.",
                ],
                "note": "Learning rate (LR) is the single most important hyperparameter to tune (typically searched on a log grid, e.g., 1e-4, 1e-3, 1e-2).",
            },
            {
                "title": "Optimizers: SGD and Adam",
                "kicker": "Learning rate, Momentum, and Adaptive Steps",
                "formula": "θ ← θ − η ∇_θ L",
                "formula_tex": r"\theta \leftarrow \theta - \eta \nabla_{\theta}\mathcal{L}",
                "formula_note": "η = learning rate  ·  typical log grid: 1e-4, 1e-3, 1e-2",
                "layout": "formula_example",
            },
            {
                "title": "Optimizers: SGD and Adam",
                "kicker": "Learning rate, Momentum, and Adaptive Steps",
                "layout": "table",
                "table": {
                    "headers": ["Optimizer", "Core mechanism", "When to use"],
                    "rows": [
                        [
                            "Vanilla SGD",
                            "Step in negative gradient direction",
                            "Rarely used alone today",
                        ],
                        [
                            "SGD + Momentum",
                            "Adds fraction of previous step direction",
                            "Standard for CV models",
                        ],
                        [
                            "Adam",
                            "Adapts learning rate per parameter based on gradients",
                            "Robust default for text/tabular/MLPs",
                        ],
                    ],
                },
                "note": "Learning rate (LR) is the single most important hyperparameter to tune (typically searched on a log grid, e.g., 1e-4, 1e-3, 1e-2).",
            },
            {
                "title": "Normalization: Batch Normalization",
                "kicker": "Stabilize Activations and Speed Up Deep Training",
                "layout": "diagram",
                "plot_path": "feature-scaling.png",
                "note": "Feature scaling: features brought to a comparable range.",
            },
            {
                "title": "Normalization: Batch Normalization",
                "kicker": "Stabilize Activations and Speed Up Deep Training",
                "body": "Batch Norm (BN): normalizes activations across the mini-batch for each channel, then learns scale γ and shift β.",
                "formula": "x̂ = (x − μ_B) / √(σ²_B + ε),   y = γ x̂ + β",
                "formula_tex": r"\hat{x}=\dfrac{x-\mu_B}{\sqrt{\sigma_B^2+\varepsilon}},\quad y=\gamma\hat{x}+\beta",
                "layout": "formula_example",
                "bullets": [
                    "It stabilizes internal activations, allowing higher learning rates and reducing dependency on precise initialization.",
                    "Alternative is Layer Norm (LN) which normalizes across features per example—standard in Transformers and sequence models.",
                ],
            },
            {
                "title": "Normalization: Batch Normalization",
                "kicker": "Stabilize Activations and Speed Up Deep Training",
                "layout": "table",
                "table": {
                    "headers": ["Method", "Normalization dimension", "Typical use-case"],
                    "rows": [
                        [
                            "Batch Norm",
                            "Across batch and spatial dimensions",
                            "CNN classifiers, ResNets",
                        ],
                        [
                            "Layer Norm",
                            "Across features for a single token/example",
                            "Transformers, RNNs",
                        ],
                    ],
                },
            },
            {
                "title": "Regularization: Dropout, Weight Decay, and Early Stopping",
                "kicker": "Constrain Capacity to Prevent Overfitting",
                "bullets": [
                    "Randomly zero out hidden units during training with probability p (e.g., 0.2-0.5) so neurons cannot co-adapt.",
                ],
                "plot_path": "dropout.png",
            },
            {
                "title": "Regularization: Dropout, Weight Decay, and Early Stopping",
                "kicker": "Constrain Capacity to Prevent Overfitting",
                "bullets": [
                    "Dropout acts as a powerful implicit ensemble method; disabled during inference (evaluation mode).",
                    "Weight Decay (L2 regularization): adds a penalty to the loss proportional to the squared sum of weights to keep weights small and smooth.",
                    "Early stopping: monitor validation loss during training and stop when it ceases to improve for a certain number of epochs (patience).",
                ],
            },
            {
                "title": "Regularization: Dropout, Weight Decay, and Early Stopping",
                "kicker": "Constrain Capacity to Prevent Overfitting",
                "layout": "table",
                "table": {
                    "headers": ["Technique", "Primary action", "Inference behavior"],
                    "rows": [
                        [
                            "Dropout",
                            "Randomly drops connections",
                            "Turned OFF (weights scaled)",
                        ],
                        [
                            "L2 Weight Decay",
                            "Shrinks weights toward zero",
                            "Turned ON (weights are static)",
                        ],
                        [
                            "Early Stopping",
                            "Halts training at lowest val loss",
                            "N/A (model is checkpointed)",
                        ],
                    ],
                },
            },
        ],
    },
    "CNNs": {
        "title": "Phase 3: Specialized Architectures (Vision & Sequences)",
        "kicker": "Images vs Ordered Sequences",
        "bullets": [
            "Images have 2D locality and translation symmetry—convolutions exploit shared weights across space.",
            "Sequences have order—recurrence, causal convolutions, or attention carry context across time or position.",
            "Choosing the right inductive bias beats blindly scaling the wrong architecture.",
        ],
        "extra_slides": [
            {
                "title": "Phase 3: Specialized Architectures (Vision & Sequences)",
                "kicker": "Images vs Ordered Sequences",
                "layout": "table",
                "table": {
                    "headers": ["Modality", "Core structure", "Representative layers"],
                    "rows": [
                        [
                            "Image / video frames",
                            "Spatial grids",
                            "Conv2D, pooling, residual blocks",
                        ],
                        [
                            "Text / speech / sensors",
                            "Ordered tokens or time steps",
                            "RNN/GRU/LSTM, causal conv, self-attention",
                        ],
                    ],
                },
            },
            {
                "title": "Image Data: How Computers “See” Pixels",
                "kicker": "Tensors, Channels, and Local Structure",
                "bullets": [
                    "A color image is usually a tensor of shape H×W×3 (RGB); each entry is an intensity discretized into 8 bits (0–255) before normalization.",
                    "Meaning lives in local neighborhoods: small patches reveal edges; larger contexts reveal objects—hierarchical composition.",
                    "Augmentation (crop, flip, color jitter) teaches invariance; normalization (e.g., ImageNet mean/std) stabilizes optimization.",
                ],
            },
            {
                "title": "Image Data: How Computers “See” Pixels",
                "kicker": "Tensors, Channels, and Local Structure",
                "layout": "diagram",
                "plot_path": "cnn-architecture.png",
            },
            {
                "title": "Image Data: How Computers “See” Pixels",
                "kicker": "Tensors, Channels, and Local Structure",
                "layout": "table",
                "table": {
                    "headers": ["Stage", "What happens"],
                    "rows": [
                        [
                            "Raw sensor",
                            "Quantized intensities per channel per pixel",
                        ],
                        [
                            "Preprocess",
                            "Resize, crop, normalize, augment",
                        ],
                        [
                            "Model input",
                            "Batch × C × H × W tensor fed to stem convolutions",
                        ],
                    ],
                },
            },
            {
                "title": "Convolution: Kernels, Edge Filters, and Feature Maps",
                "kicker": "Sliding Linear Filters + Nonlinearity",
                "bullets": [
                    "A kernel is a small learnable matrix (e.g., 3×3) that responds to local patterns; multiple kernels produce a stack of feature maps.",
                    "Classical edge detectors (Sobel, shown) are fixed kernels; CNNs learn task-specific filters from data.",
                    "Deep stacks widen the receptive field: early layers ≈ edges/textures; late layers ≈ object parts and semantics.",
                ],
            },
            {
                "title": "Convolution: Kernels, Edge Filters, and Feature Maps",
                "kicker": "Sliding Linear Filters + Nonlinearity",
                "layout": "diagram",
                "plot_path": "convolution.png",
            },
            {
                "title": "Convolution: Kernels, Edge Filters, and Feature Maps",
                "kicker": "Sliding Linear Filters + Nonlinearity",
                "layout": "diagram",
                "plot_path": "sobel-y.png",
            },
            {
                "title": "Convolution: Kernels, Edge Filters, and Feature Maps",
                "kicker": "Sliding Linear Filters + Nonlinearity",
                "layout": "table",
                "table": {
                    "headers": ["Hyperparameter", "Effect"],
                    "rows": [
                        ["Kernel size", "Immediate neighborhood size (3×3 common)"],
                        ["Stride", "Downsample spatially when >1"],
                        ["Padding", "Keep spatial size (same padding) or control border"],
                    ],
                },
            },
            {
                "title": "CNN Blocks: From Conv-Pool to ResNet",
                "kicker": "Depth, Regularization, and Residual Learning",
                "layout": "formula_example",
                "formula": "y = F(x) + x",
                "formula_tex": r"y=\mathcal{F}(x)+x",
                "formula_note": "(residual block learns F, identity carries gradient)",
            },
            {
                "title": "CNN Blocks: From Conv-Pool to ResNet",
                "kicker": "Depth, Regularization, and Residual Learning",
                "layout": "diagram",
                "plot_path": "cnn-stack.png",
            },
            {
                "title": "CNN Blocks: From Conv-Pool to ResNet",
                "kicker": "Depth, Regularization, and Residual Learning",
                "layout": "diagram",
                "plot_path": "residual-block.png",
            },
            {
                "title": "CNN Blocks: From Conv-Pool to ResNet",
                "kicker": "Depth, Regularization, and Residual Learning",
                "bullets": [
                    "Standard CNN: Alternates Conv2D, Activation (ReLU), and Pooling (Max Pool) layers to extract spatial features and downsample.",
                    "AlexNet (2012): First large-scale CNN to win ImageNet; used ReLU, dropout, and multi-GPU training.",
                    "ResNet (2015): Introduced skip (shortcut) connections, allowing gradients to flow back easily through identity mappings. Enables training networks with 100+ layers.",
                ],
            },
            {
                "title": "CNN Blocks: From Conv-Pool to ResNet",
                "kicker": "Landmark ideas (simplified)",
                "layout": "table",
                "table": {
                    "headers": ["Model / block", "Idea you should remember"],
                    "rows": [
                        [
                            "Conv + MaxPool",
                            "Feature extraction + spatial downsampling",
                        ],
                        [
                            "AlexNet",
                            "Scale + ReLU + data aug on GPUs",
                        ],
                        [
                            "ResNet",
                            "Residual mapping F(x)+x; skip paths help gradients",
                        ],
                    ],
                },
            },
            {
                "title": "Transfer Learning with Pretrained Models",
                "kicker": "Reuse Features, Adapt the Head",
                "bullets": [
                    "Pretraining on millions of labeled images (e.g. ImageNet) learns general low/mid-level filters (edges, textures) transferable to new domains.",
                    "Typical recipe: replace the final classifier head; freeze the pretrained backbone for small datasets, or fine-tune with a very small learning rate for larger datasets.",
                    "Transfer learning dramatically reduces training time and required labeled samples for a new vision task.",
                ],
            },
            {
                "title": "Transfer Learning with Pretrained Models",
                "kicker": "Reuse Features, Adapt the Head",
                "layout": "diagram",
                "plot_path": "transfer-learning.png",
            },
            {
                "title": "Transfer Learning with Pretrained Models",
                "kicker": "Reuse Features, Adapt the Head",
                "layout": "table",
                "table": {
                    "headers": ["Target Dataset Size", "Strategy"],
                    "rows": [
                        [
                            "Small target dataset",
                            "Freeze backbone, train head only",
                        ],
                        [
                            "Medium dataset",
                            "Fine-tune top blocks with a smaller learning rate",
                        ],
                        [
                            "Large target dataset",
                            "Full fine-tune of all layers with a small learning rate",
                        ],
                    ],
                },
            },
        ],
    },
    "RNNs & regularization": {
        "title": "Sequential Data & Vanilla RNNs",
        "kicker": "One Transition, Unfolded Through Time",
        "layout": "formula_example",
        "formula": "hₜ = φ(W_hh hₜ₋₁ + W_xh xₜ + b)",
        "formula_tex": r"h_t=\varphi(W_{hh}h_{t-1}+W_{xh}x_t+b)",
        "extra_slides": [
            {
                "title": "Sequential Data & Vanilla RNNs",
                "kicker": "One Transition, Unfolded Through Time",
                "layout": "diagram",
                "plot_path": "rnn-unfold.png",
            },
            {
                "title": "Sequential Data & Vanilla RNNs",
                "kicker": "One Transition, Unfolded Through Time",
                "bullets": [
                    "Sequence modeling: text, audio, and sensor streams have order; shuffling them breaks their semantic meaning.",
                    "Recurrent Neural Networks (RNN): maintain a recurrent state vector hₜ that acts as a memory summary of the past.",
                    "The same weights W_hh and W_xh are reused across all timesteps (parameter efficiency), but long-range dependencies suffer from vanishing/exploding gradients.",
                ],
            },
            {
                "title": "Sequential Data & Vanilla RNNs",
                "kicker": "One Transition, Unfolded Through Time",
                "layout": "table",
                "table": {
                    "headers": ["Feature", "Limitation"],
                    "rows": [
                        [
                            "Reuses parameters",
                            "Sequential execution (cannot parallelize steps)",
                        ],
                        [
                            "Tracks state hₜ",
                            "Cannot hold long-term memories in practice",
                        ],
                    ],
                },
            },
            {
                "title": "LSTM: Gates for Long-Term Memory",
                "kicker": "Forget, Input, Output — Cell State Flows Linearly",
                "layout": "diagram",
                "plot_path": "lstm-gates.png",
            },
            {
                "title": "LSTM: Gates for Long-Term Memory",
                "kicker": "Forget, Input, Output — Cell State Flows Linearly",
                "bullets": [
                    "Cell state cₜ can accumulate information with an additive path—mitigates vanishing signal compared with plain tanh RNNs.",
                    "Forget gate fₜ decides how much past cell content to erase; input gate iₜ and candidate g̃ control new information written.",
                    "Output gate oₜ filters what becomes hidden state hₜ exposed to the next layer or timestep.",
                ],
            },
            {
                "title": "LSTM: Gates for Long-Term Memory",
                "kicker": "Forget, Input, Output — Cell State Flows Linearly",
                "layout": "table",
                "table": {
                    "headers": ["Gate", "Role (intuition)"],
                    "rows": [
                        ["fₜ", "Erase irrelevant history from the cell"],
                        ["iₜ", "Allow new candidate information in"],
                        ["oₜ", "Expose part of the cell as hₜ"],
                    ],
                },
            },
            {
                "title": "Phase 4: Generative Models & Reconstruction",
                "kicker": "Generative Models and Representation Learning",
                "bullets": [
                    "Phase 4 explores unsupervised representation learning and basic generative frameworks.",
                    "Autoencoders are the foundation here: they learn to compress inputs to a low-dimensional bottleneck (latent space) and reconstruct them.",
                ],
            },
            {
                "title": "Phase 4: Generative Models & Reconstruction",
                "kicker": "Generative Models and Representation Learning",
                "layout": "table",
                "table": {
                    "headers": ["Concept", "Input", "Goal"],
                    "rows": [
                        [
                            "Autoencoder",
                            "Unlabeled data x",
                            "Compress and reconstruct x̂ ≈ x",
                        ],
                    ],
                },
            },
            {
                "title": "Autoencoders (AE): Compression and Reconstruction",
                "kicker": "Encoder → Bottleneck → Decoder",
                "layout": "diagram",
                "plot_path": "autoencoder.png",
            },
            {
                "title": "Autoencoders (AE): Compression and Reconstruction",
                "kicker": "Encoder → Bottleneck → Decoder",
                "bullets": [
                    "Encoder f maps input x to a low-dimensional code z = f(x); decoder g maps z back to x̂ ≈ x.",
                    "The bottleneck forces a compressed representation—useful for denoising (train on noisy→clean), anomaly detection (high recon error = abnormal).",
                    "Linear AE with MSE relates to PCA when constraints align—nonlinear AE learns curved manifolds.",
                ],
            },
            {
                "title": "Autoencoders (AE): Compression and Reconstruction",
                "kicker": "Encoder → Bottleneck → Decoder",
                "layout": "table",
                "table": {
                    "headers": ["Variant", "Idea"],
                    "rows": [
                        [
                            "Denoising AE",
                            "Corrupt input, reconstruct clean—robust features",
                        ],
                        [
                            "Sparse AE",
                            "Penalty on activations—encourages informative sparse codes",
                        ],
                        [
                            "Contractive AE",
                            "Penalize sensitivity of code to small input changes—smoother latent map",
                        ],
                    ],
                },
            },
        ],
    },
}

TAKEAWAYS = [
    "A neuron is z = wᵀx + b then a = f(z). Depth needs nonlinearity between layers.",
    "Backprop is the chain rule on a graph. Watch vanishing gradients, LR, and train vs val.",
    "CNNs share filters across space; RNNs share weights across time; regularize so they generalize.",
]

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "Week 1 classical ML is done — now stacked nonlinear models.",
    "current": "S7",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.05), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.5),
        Inches(11.5),
        Inches(0.85),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.5), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(3.85), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    topics = s["topics"]
    n_cols = 3
    chip_w = Inches(3.7)
    chip_h = Inches(0.48)
    gap_x = Inches(0.22)
    start_y = Inches(4.7)
    for i, topic in enumerate(topics):
        row, col = divmod(i, n_cols)
        x = MARGIN + col * (chip_w + gap_x)
        y = start_y + Inches(row * 0.68)
        soft_card(slide, x, y, chip_w, chip_h, fill=SOFT)
        add_text(slide, x, y + Inches(0.06), chip_w, Inches(0.36), topic, size=12, color=PRIMARY, align=PP_ALIGN.CENTER)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def content_title(slide, title: str, subtitle: str | None = None):
    """Title that shrinks on long headlines so it does not cover the kicker."""
    long = len(title or "") > 46
    size = 24 if long else 32
    title_h = 0.80 if long else 0.62
    add_text(
        slide,
        MARGIN,
        Inches(1.12),
        Inches(11.2),
        Inches(title_h),
        title,
        size=size,
        bold=True,
        color=PRIMARY,
    )
    if subtitle:
        add_text(
            slide,
            MARGIN,
            Inches(1.12 + title_h),
            Inches(11.2),
            Inches(0.34),
            subtitle,
            size=14 if long else 15,
            color=MUTED,
        )


def content_top(title: str) -> float:
    return 2.42 if len(title or "") > 46 else 2.18


def _chars_per_line(width_in: float, size: float) -> int:
    return max(22, int(width_in * (108 / max(size, 10))))


def _line_count(text: str, width_in: float, size: float) -> int:
    cpl = _chars_per_line(width_in, size)
    words = (text or "").split()
    if not words:
        return 1
    lines, cur = 1, 0
    for word in words:
        need = len(word) + (1 if cur else 0)
        if cur + need > cpl:
            lines += 1
            cur = len(word)
        else:
            cur += need
    return lines


def _bullet_box(items: list[str], width_in: float, size: float) -> tuple[float, float]:
    lines = max((_line_count(t, width_in, size) for t in items), default=1)
    item_h = max(0.42, min(1.15, lines * (size / 72) * 1.35 + 0.10))
    return item_h, item_h + 0.10


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
        if px_w <= 0 or px_h <= 0:
            return None
        aspect = px_w / px_h
        src = path
        if aspect < 0.5:
            from io import BytesIO

            crop_h = max(1, min(px_h, int(px_w / 2.4)))
            cropped = im.crop((0, 0, px_w, crop_h))
            buf = BytesIO()
            cropped.save(buf, format="PNG")
            buf.seek(0)
            src = buf
            px_w, px_h = cropped.size
            aspect = px_w / max(px_h, 1)
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(src) if src is path else src,
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    bsize = 14 if dense else 16

    if content.get("body"):
        top = content_top(content["title"])
        soft_card(slide, MARGIN, Inches(top), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(top + 0.20),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(top + 1.35)
        bsize = 14 if dense else 15
    else:
        bullet_top = Inches(content_top(content["title"]))
        bsize = 14 if dense else 15

    item_h, pitch = _bullet_box(items, 11.2, bsize)
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(item_h),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    has_body = bool((content.get("body") or "").strip())
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)

    y_after = content_top(content["title"])
    if has_body:
        soft_card(slide, MARGIN, Inches(y_after), col_w, Inches(1.2), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(y_after + 0.15),
            text_w,
            Inches(0.95),
            content.get("body") or "",
            size=13,
            color=INK,
        )
        y_after = y_after + 1.35

    if has_formula:
        formula_card_h = 1.65 if is_frac else 1.35
        formula_box_h = 0.95 if is_frac else 0.65
        soft_card(slide, MARGIN, Inches(y_after), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(y_after + 0.10),
            text_w,
            Inches(0.22),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(y_after + 0.35),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        if content.get("formula_note"):
            add_text(
                slide,
                Inches(MARGIN.inches + 0.3),
                Inches(y_after + 0.35 + formula_box_h + 0.02),
                text_w,
                Inches(0.28),
                content["formula_note"],
                size=10,
                color=MUTED,
            )
        bullet_top = y_after + formula_card_h + 0.12
        bsize = 12
    else:
        bullet_top = y_after
        bsize = 13

    item_h, pitch = _bullet_box(items, bullet_w.inches, bsize)
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch)) if items else 0
    if has_formula and has_plot:
        max_items = min(max_items, 2)
    if items:
        bullets(
            slide,
            items[:max_items],
            top=Inches(bullet_top),
            size=bsize,
            left=MARGIN,
            width=bullet_w,
            pitch=pitch,
            item_height=Inches(item_h),
        )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(content_top(content["title"]))
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(body_top.inches + 0.13),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(body_top.inches + 0.85)
    else:
        formula_top = body_top

    formula_h = 1.55 if is_frac else 1.25
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.12),
        Inches(11.2),
        Inches(0.95 if is_frac else 0.7),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )
    if content.get("formula_note"):
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(formula_top.inches + formula_h - 0.35),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    bullet_top = formula_top.inches + formula_h + 0.12
    bsize = 13 if dense else 15
    item_h, pitch = _bullet_box(items, 11.2, bsize)
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch)) if items else 0
    if has_note:
        max_items = min(max_items, 3)
    if items:
        bullets(
            slide,
            items[:max_items],
            top=Inches(bullet_top),
            size=bsize,
            pitch=pitch,
            width=Inches(11.2),
            item_height=Inches(item_h),
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = content_top(content["title"])
    if show_body:
        soft_card(slide, MARGIN, Inches(table_top), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(table_top + 0.10),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = table_top + 0.70

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    note_top = table_top + table_h + 0.08
    render_note = has_note and (note_top + 0.45) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    diagram_top = Inches(content_top(content["title"]))
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(diagram_top.inches - 0.05),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(diagram_top.inches + 0.30)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(4.35), fill=SOFT)
    _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + 0.15),
        Inches(11.6),
        Inches(4.05),
        folder=folder,
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.9),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def slide_takeaways(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Takeaways", f"{index:02d}")
    title_block(slide, "3 Ideas to Keep", "Leave with these")
    for i, text in enumerate(TAKEAWAYS):
        y = Inches(2.35) + Inches(i * 1.25)
        soft_card(slide, MARGIN, y, Inches(12.0), Inches(1.08), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            y + Inches(0.28),
            Inches(0.5),
            Inches(0.5),
            str(i + 1),
            size=22,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            MARGIN + Inches(0.95),
            y + Inches(0.28),
            Inches(10.6),
            Inches(0.6),
            text,
            size=16,
            color=INK,
        )
    content_footer(slide, index, total)


def slide_close(prs, total, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))
    add_text(
        slide,
        MARGIN,
        Inches(2.9),
        Inches(12),
        Inches(0.8),
        "Thank you",
        size=40,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN,
        Inches(3.8),
        Inches(12),
        Inches(0.4),
        "Labs: code/15- Deep Learning",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    content_footer(slide, index, total)


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 6 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    n += 1
    slide_takeaways(prs, total, n)
    n += 1
    slide_close(prs, total, n)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
