#!/usr/bin/env python3
"""
Build Week 3 · Session 1 presentation using official ETRA Design System.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week3-Session1-NLP.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session-w3s1-diagrams"

SESSION = {
    "eyebrow": "Week 3  ·  Session 1",
    "section_number": "08",
    "section_tag": "S08",
    "section_title": "NLP",
    "subtitle": "NLP Workshop: Foundations to Seq2Seq",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "Foundations · Tokenization · Seq2Seq",
    "topics": [
        "NLP Fundamentals",
        "Regex & Cleaning",
        "Tokenization",
        "POS & NER",
        "Language Modeling",
        "Embeddings & RNNs",
        "Seq2Seq & Attention",
    ],
}

TOPIC_CONTENT: dict = {
    "NLP Fundamentals": {
        "title": "NLP Fundamentals and Challenges",
        "kicker": "Why Natural Language Is Difficult for Machines",
        "bullets": [
            "Natural language is ambiguous, context-dependent, and full of idioms.",
            "Meaning often depends on pragmatics, domain, and cultural background.",
            "The same sentence can map to multiple valid interpretations.",
        ],
        "plot_path": "linguistic-levels.png",
        "note": "NLP systems must model form, meaning, and context simultaneously.",
        "extra_slides": [
            {
                "title": "Ambiguity in Language: Practical Examples",
                "layout": "table",
                "table": {
                    "headers": ["Expression", "Possible Interpretation A", "Possible Interpretation B"],
                    "rows": [
                        ["The bank is closed", "Financial institution is closed", "River bank area is inaccessible"],
                        ["It's cold", "Low temperature", "Emotionally distant behavior"],
                        ["The chicken is ready to eat", "Food is ready to be eaten", "Animal is ready to eat food"],
                    ],
                },
                "note": "Robust NLP requires disambiguation using surrounding textual and situational context.",
            },
        ],
    },
    "Regex & Cleaning": {
        "title": "Regular Expressions (Regex): Core Idea",
        "bullets": [
            "Regex is a compact pattern language for matching and transforming text.",
            "Typical uses include validation, extraction, and rule-based cleanup.",
            "Useful for emails, phone numbers, logs, and normalization pipelines.",
        ],
        "extra_slides": [
            {
                "title": "Regular Expressions (Regex): Core Idea",
                "layout": "table",
                "table": {
                    "headers": ["Pattern", "Meaning", "Example Match"],
                    "rows": [
                        ["\\d+", "one or more digits", "2026"],
                        ["\\w+", "one or more word chars", "nlp_model1"],
                        ["^[A-Z]", "starts with uppercase letter", "Hello"],
                        ["\\S+", "one or more non-space chars", "token123"],
                    ],
                },
            },
            {
                "title": "Text Cleaning — Conceptual Steps & Python Code",
                "kicker": "Conceptual stages matched with Python regex implementations",
                "layout": "table",
                "table": {
                    "headers": ["Stage", "Python Implementation", "Goal"],
                    "rows": [
                        ["1. Decode & Clean HTML", "html.unescape(text) & re.sub(r'<[^>]+>', '', text)", "Keep readable text only"],
                        ["2. Remove Contacts/URLs", "re.sub(r'https?://\\S+|www\\.\\S+', '', text)", "Drop URL metadata"],
                        ["3. Remove Mentions/Socials", "re.sub(r'@\\w+|#\\w+', '', text)", "Remove handles & hashtags"],
                        ["4. Drop Non-ASCII & Emoji", "text.encode('ascii', 'ignore').decode('ascii')", "Drop emoji/special symbols"],
                        ["5. Normalize Whitespace", "re.sub(r'\\s+', ' ', text).strip()", "Normalize extra spaces"],
                    ],
                },
                "note": "A regex-based cleaning function is usually the first stage in an NLP preprocessing pipeline.",
            },
            {
                "title": "Text Cleaning — Worked Example",
                "kicker": "Trace of raw input transitioning to clean text",
                "layout": "table",
                "table": {
                    "headers": ["Stage", "Operation performed", "Text snapshot"],
                    "rows": [
                        ["0. Raw Input", "Original raw string", "<p>John: \"AI is amazing!!!! 🤖🔥\" Visit: https://ai.com @john #AI</p>"],
                        ["1. Clean HTML", "Decode entities & remove tags", "John: \"AI is amazing!!!! 🤖🔥\" Visit: https://ai.com @john #AI"],
                        ["2. Remove URLs", "Remove web addresses", "John: \"AI is amazing!!!! 🤖🔥\" Visit: @john #AI"],
                        ["3. Remove Socials", "Remove @mentions and #hashtags", "John: \"AI is amazing!!!! 🤖🔥\" Visit:"],
                        ["4. Drop Emoji/Symbols", "ASCII-only encoding & cleanup", "John AI is amazing!!!! Visit"],
                    ],
                },
                "note": "Text cleaning produces standard, noise-free input for tokenizer components.",
            },
        ],
    },
    "Tokenization": {
        "title": "Tokenization Granularities (Word, Char, Subword)",
        "layout": "table",
        "table": {
            "headers": ["Granularity", "Strength", "Tradeoff", "Example output (for 'running')"],
            "rows": [
                ["Word / Space-based", "Human-readable tokens", "Weak on rare/compound words (OOV)", "['running']"],
                ["Character-based", "No unknown words (no OOV)", "Very long sequences, loss of meaning", "['r', 'u', 'n', 'n', 'i', 'n', 'g']"],
                ["Subword-based", "Handles rare words, compact vocab", "Requires tokenizer training/tuning", "['run', '##ning']"],
            ],
        },
        "note": "Subword tokenization is the industry standard for modern LLMs and neural NLP models.",
        "extra_slides": [
            {
                "title": "Tokenization Views (Code Examples)",
                "kicker": "Same sentence, word / char / subword splits side by side",
                "layout": "diagram",
                "plot_path": "hello-world-tokens.png",
            },
            {
                "title": "Tokenization Views (Code Examples)",
                "kicker": "Same sentence, word / char / subword splits side by side",
                "layout": "table",
                "table": {
                    "headers": ["Tokenizer Type", "Python Implementation", "Result on \"I'm running to the store!\""],
                    "rows": [
                        ["Word", "word_tokenize(text)", "['I', \"'m\", 'running', 'to', 'the', 'store', '!']"],
                        ["Character", "list(text)", "['I', \"'\", 'm', ' ', 'r', 'u', 'n', 'n', 'i', 'n', 'g', ...]"],
                        ["Subword", "model_tokenizer.tokenize(text)", "['I', \"'m\", 'run', '##ning', 'to', 'the', 'store', '!']"],
                        ["Sentence", "sent_tokenize(text)", "[\"I'm running to the store!\"]"],
                    ],
                },
            },
            {
                "title": "Subword Tokenization (LLMs) — Why It Wins",
                "kicker": "How subwords solve core vocabulary challenges",
                "layout": "diagram",
                "plot_path": "hello-world-tokens.png",
            },
            {
                "title": "Subword Tokenization (LLMs) — Why It Wins",
                "kicker": "How subwords solve core vocabulary challenges",
                "layout": "table",
                "table": {
                    "headers": ["Word Tokenizer Issue", "Subword Solution", "Example Worked Out"],
                    "rows": [
                        ["Huge vocabulary size", "Reuse frequent morphemes/pieces", "unbelievably → un + ##believe + ##ably"],
                        ["Out-of-Vocabulary (OOV) words", "Break unknown forms into known units", "ChatGPT → ['Chat', 'G', 'PT']"],
                        ["Morphological variants", "Shared roots reduce feature sparsity", "running/runs/runner share 'run'"],
                    ],
                },
            },
            {
                "title": "BPE Walkthrough & Inference Example",
                "kicker": "How Byte Pair Encoding learns merges and splits unseen words",
                "layout": "diagram",
                "plot_path": "hello-world-tokens.png",
            },
            {
                "title": "BPE Walkthrough & Inference Example",
                "kicker": "How Byte Pair Encoding learns merges and splits unseen words",
                "bullets": [
                    "Walkthrough: start with chars (l, o, w, e, s, t, r) -> merge l+o (lo) -> merge lo+w (low) -> compact vocabulary.",
                    "Inference on 'lowered the lowest curtain': splits rare words to ['low', 'er', 'ed'] and ['low', 'est'] while keeping 'the' and 'curtain' whole.",
                ],
            },
            {
                "title": "BPE Walkthrough & Inference Example",
                "kicker": "How Byte Pair Encoding learns merges and splits unseen words",
                "layout": "table",
                "table": {
                    "headers": ["Word", "Subword Output", "Explanation"],
                    "rows": [
                        ["lowered", "['low', 'er', 'ed']", "Decomposes into learned root and suffix subwords"],
                        ["lowest", "['low', 'est']", "Decomposes into base 'low' and suffix 'est'"],
                    ],
                },
                "note": "BPE repeatedly merges frequent adjacent units until target vocabulary size is reached.",
            },
            {
                "title": "Stop Words: Purpose & Tradeoffs",
                "kicker": "Filtering low-information tokens from texts",
                "layout": "diagram",
                "plot_path": "hello-world-tokens.png",
            },
            {
                "title": "Stop Words: Purpose & Tradeoffs",
                "kicker": "Filtering low-information tokens from texts",
                "bullets": [
                    "Function words (the, is, and) dominate word counts but carry minimal topic signal.",
                    "Removing stop words reduces dimensionality for TF-IDF & Bag-of-Words.",
                    "Critical Rule: NEVER remove stop words for language models (GPT/BERT), translation, or sentiment analysis where negation ('not') and syntax carry the core meaning.",
                ],
            },
            {
                "title": "Stop Words: Purpose & Tradeoffs",
                "kicker": "Filtering low-information tokens from texts",
                "layout": "table",
                "table": {
                    "headers": ["Task Category", "Stop Word Policy", "Reasoning"],
                    "rows": [
                        ["Search & Topic Modeling", "Remove", "Improves speed and focus on key words"],
                        ["Large Language Models", "Keep", "Syntactic and semantic flow is vital"],
                        ["Sentiment Analysis", "Keep", "Negations like 'not' reverse polarity"],
                    ],
                },
            },
            {
                "title": "Stemming vs Lemmatization",
                "kicker": "Normalizing words to their base forms",
                "layout": "table",
                "table": {
                    "headers": ["Aspect", "Stemming", "Lemmatization"],
                    "rows": [
                        ["Mechanism", "Rule-based suffix chopping (heuristic)", "Dictionary lookup + morphological analysis"],
                        ["Speed", "Very fast, computationally cheap", "Slower, requires resource lookup"],
                        ["Output", "Can be non-words (e.g., study -> studi)", "Always dictionary words (e.g., study -> study)"],
                        ["POS Awareness", "No", "Yes (e.g., saw -> see or saw depending on POS)"],
                    ],
                },
                "note": "Lemmatization preserves semantics better than aggressive stemming.",
            },
        ],
    },
    "POS & NER": {
        "title": "Part-of-Speech (POS) Tagging",
        "kicker": "Grammar Labels per Token",
        "layout": "diagram",
        "plot_path": "hello-world-tokens.png",
        "extra_slides": [
            {
                "title": "Part-of-Speech (POS) Tagging",
                "kicker": "Grammar Labels per Token",
                "bullets": [
                    "Universal categories include NOUN, VERB, ADJ, ADV, DET, ADP, PROPN, NUM, …",
                    "POS feeds lemmatization quality and downstream relation extraction.",
                    "SpaCy exposes coarse pos_ and fine-grained tag_ plus dependency dep_.",
                ],
            },
            {
                "title": "Part-of-Speech (POS) Tagging",
                "kicker": "Grammar Labels per Token",
                "layout": "table",
                "table": {
                    "headers": ["POS", "Role", "Quick examples"],
                    "rows": [
                        ["NOUN / PROPN", "Things and names", "city, Tesla"],
                        ["VERB / AUX", "Actions and helpers", "run, is"],
                        ["ADJ / ADV", "Modifiers", "quick, quickly"],
                        ["DET / ADP", "Structure words", "the, in"],
                    ],
                },
            },
            {
                "title": "Named Entity Recognition (NER)",
                "kicker": "Typed Spans over Text",
                "bullets": [
                    "Labels include PERSON, ORG, GPE, DATE, MONEY, PERCENT, PRODUCT, EVENT, …",
                    "Useful for indexing, compliance redaction, financial news graphs, and search facets.",
                    "Small models err on edge cases — always spot-check domain text.",
                ],
            },
            {
                "title": "Named Entity Recognition (NER)",
                "kicker": "Typed Spans over Text",
                "layout": "table",
                "table": {
                    "headers": ["Label", "Examples"],
                    "rows": [
                        ["PERSON", "Tim Cook, Jensen Huang"],
                        ["ORG / GPE", "Apple, Texas"],
                        ["MONEY / PERCENT", "$430 billion, 3.5%"],
                    ],
                },
            },
            {
                "title": "Full SpaCy Pipeline Functionally",
                "bullets": [
                    "One doc object: tokens, lemmas, POS, entities, noun_chunks, sents.",
                    "Typical export: clean token list, lemma bag, entity list, POS histogram.",
                    "Use length-based heuristics (e.g. longest sentence) only as weak importance cues.",
                ],
            },
        ],
    },
    "Language Modeling": {
        "title": "NLP Language Modeling with N-grams",
        "kicker": "From Count-Based Prediction to Evaluation",
        "bullets": [
            "Language models estimate probabilities over word sequences.",
            "N-gram models predict the next token from a limited context window.",
            "This module connects N-gram intuition to modern LLM generation.",
        ],
        "plot_path": "google-ngram.png",
        "extra_slides": [
            {
                "title": "Formal Language Modeling Objective",
                "formula": "P(w₁,…,w_T) = Π P(w_t | w₁,…,w_{t-1})",
                "formula_tex": r"P(w_1,\ldots,w_T)=\prod_{t=1}^{T} P(w_t \mid w_1,\ldots,w_{t-1})",
                "layout": "formula_example",
                "bullets": [
                    "A language model can score full sentences or predict the next word.",
                    "Exact estimation over full histories is intractable for real corpora.",
                    "Approximation strategies include N-grams, smoothing, and neural LMs.",
                ],
            },
            {
                "title": "Worked Bigram Example",
                "body": "Given counts: C(\"want to\") = 608, C(\"want\") = 927",
                "formula": "P(\"to\" | \"want\") = 608 / 927 ≈ 0.656",
                "formula_tex": r"P(\mathrm{to}\mid\mathrm{want})=\dfrac{608}{927}\approx 0.656",
                "layout": "formula_example",
                "bullets": [
                    "Interpretation: in this corpus, \"to\" follows \"want\" about 65.6% of the time.",
                    "Count quality depends strongly on corpus domain and size.",
                ],
            },
            {
                "title": "Why Use Log Probabilities",
                "bullets": [
                    "Multiplying many small probabilities causes numerical underflow.",
                    "Log transform converts multiplication into stable addition.",
                    "Sequence scoring becomes computationally robust and efficient.",
                ],
            },
            {
                "title": "Why Use Log Probabilities",
                "layout": "table",
                "table": {
                    "headers": ["Original Space", "Log Space"],
                    "rows": [
                        ["P = p1 * p2 * ... * pk", "log P = log p1 + log p2 + ... + log pk"],
                        ["Very small numbers", "Numerically stable sums"],
                    ],
                },
            },
            {
                "title": "Perplexity (PP): Interpretation",
                "formula": "PP(W) = P(w₁,…,w_T)^{-1/T}",
                "formula_tex": r"\mathrm{PP}(W)=P(w_1,\ldots,w_T)^{-\frac{1}{T}}=\exp\left(-\frac{1}{T}\sum_{t=1}^{T}\log P(w_t\mid h_t)\right)",
                "layout": "formula_example",
                "bullets": [
                    "Perplexity is the average branching uncertainty of the model.",
                    "Lower PP means better predictive confidence on unseen text.",
                    "Use identical test sets when comparing models.",
                ],
                "note": "Raw sentence probabilities are length-sensitive; perplexity normalizes by token count.",
            },
        ],
    },
    "Embeddings & RNNs": {
        "title": "Static vs Contextualized Word Representations",
        "layout": "table",
        "table": {
            "headers": ["Property", "Static Embeddings (Word2Vec/GloVe)", "Contextualized Embeddings (ELMo/BERT/GPT)"],
            "rows": [
                ["Vector per word", "One fixed vector", "Different vectors per context"],
                ["Polysemy handling", "Weak", "Strong"],
                ["Context direction", "Usually local/global corpus only", "Bidirectional or autoregressive sequence context"],
                ["Task transfer", "Moderate", "High with pretraining + fine-tuning"],
            ],
        },
        "note": "Contextualization is a core enabler for modern LLM quality.",
        "extra_slides": [
            {
                "title": "How Contextualized Embeddings Work",
                "kicker": "Words as vectors in embedding space",
                "bullets": [
                    "Static embeddings: one vector per word type.",
                    "Contextual: vector depends on surrounding tokens (ELMo, BERT-style).",
                ],
                "plot_path": "embedding-space.png",
            },
            {
                "title": "How Contextualized Embeddings Work",
                "kicker": "Architecture Layers",
                "bullets": [
                    "Token/base embedding layer maps input ids to dense vectors.",
                    "Context encoder layers (often Transformer blocks) enrich each token with sequence context.",
                    "Task head uses contextual features for prediction.",
                ],
            },
            {
                "title": "How Contextualized Embeddings Work",
                "kicker": "Layer Semantics",
                "bullets": [
                    "Lower layers capture local syntax/patterns.",
                    "Middle layers improve sense disambiguation.",
                ],
            },
            {
                "title": "RNN Mechanism and Weight Sharing",
                "formula": "h_t = f(W_x x_t + W_h h_{t-1} + b)",
                "formula_tex": r"h_t = f(W_x x_t + W_h h_{t-1} + b)",
                "plot_path": "rnn-unfold.png",
            },
            {
                "title": "Common RNN Input/Output Patterns",
                "layout": "table",
                "table": {
                    "headers": ["Pattern", "Mapping", "Example Task"],
                    "rows": [
                        ["One-to-One", "single input -> single output", "basic regression/classification"],
                        ["One-to-Many", "single input -> sequence output", "image captioning"],
                        ["Many-to-One", "sequence input -> single output", "sentiment classification"],
                        ["Many-to-Many", "sequence input -> sequence output", "machine translation"],
                    ],
                },
                "note": "Sequence-aware patterns are where recurrent models provide clear value.",
            },
            {
                "title": "RNN Family: Vanilla, GRU, BiRNN, LSTM",
                "layout": "diagram",
                "plot_path": "gru-cell.png",
            },
            {
                "title": "RNN Family: Vanilla, GRU, BiRNN, LSTM",
                "layout": "table",
                "table": {
                    "headers": ["Model", "Vanishing Gradient Risk", "Key Strength", "Typical Limitation"],
                    "rows": [
                        ["Vanilla RNN", "High", "Simple architecture", "Poor long-range memory"],
                        ["GRU", "Lower", "Efficient gating with fewer parameters", "Can still degrade on very long contexts"],
                        ["BiRNN", "Medium", "Uses past and future context", "Higher compute and memory"],
                        ["LSTM", "Low (relative)", "Strong long-term dependency handling", "Heavier than GRU"],
                    ],
                },
            },
        ],
    },
    "Seq2Seq & Attention": {
        "title": "NLP Seq2Seq for Neural Machine Translation",
        "kicker": "Encoder-Decoder Modeling, Decoding Strategies, and Evaluation",
        "layout": "diagram",
        "plot_path": "encoding-comparison.png",
        "extra_slides": [
            {
                "title": "NLP Seq2Seq for Neural Machine Translation",
                "kicker": "Encoder-Decoder Modeling, Decoding Strategies, and Evaluation",
                "bullets": [
                    "Seq2Seq maps variable-length input sequences to variable-length outputs.",
                    "Encoder-decoder models were foundational for neural machine translation.",
                    "This module covers training, decoding, bottlenecks, attention, and metrics.",
                ],
            },
            {
                "title": "Inference: Greedy Decoding vs Beam Search",
                "layout": "diagram",
                "plot_path": "seq2seq-attention.png",
            },
            {
                "title": "Inference: Greedy Decoding vs Beam Search",
                "layout": "table",
                "table": {
                    "headers": ["Method", "Decision Rule", "Strength", "Risk"],
                    "rows": [
                        ["Greedy", "Pick top token each step", "Fast and simple", "May miss globally better sequence"],
                        ["Beam Search", "Track top-k hypotheses each step", "Better global sequence quality", "Higher compute; can become generic/bland"],
                    ],
                },
                "note": "Typical beam sizes are moderate (e.g., 4-10) to balance quality and cost.",
            },
            {
                "title": "Information Bottleneck in Basic Seq2Seq",
                "layout": "diagram",
                "plot_path": "seq2seq-attention.png",
            },
            {
                "title": "Information Bottleneck in Basic Seq2Seq",
                "bullets": [
                    "Compressing a full source sentence into one fixed vector can lose detail.",
                    "Longer/complex inputs worsen the bottleneck effect.",
                    "Decoder needs different source details at different output steps.",
                ],
                "note": "This motivates attention over all encoder hidden states.",
            },
            {
                "title": "Attention as the Bottleneck Solution",
                "layout": "diagram",
                "plot_path": "attention-architecture.png",
            },
            {
                "title": "Attention as the Bottleneck Solution",
                "bullets": [
                    "Decoder attends to relevant encoder positions at each generation step.",
                    "Dynamic alignment improves translation adequacy and fluency.",
                    "Attention laid the foundation for transformer-dominant NMT systems.",
                ],
            },
            {
                "title": "NMT Evaluation Metrics: BLEU and ROUGE",
                "layout": "table",
                "table": {
                    "headers": ["Metric", "Orientation", "What It Measures", "Common Limitation"],
                    "rows": [
                        ["BLEU", "Precision-oriented", "n-gram overlap from candidate to reference", "Weak semantic/syntactic sensitivity"],
                        ["ROUGE-N", "Recall-oriented", "n-gram overlap from reference to candidate", "Can reward lexical overlap over meaning"],
                        ["F1 (with overlap metrics)", "Balance precision/recall", "Harmonic tradeoff view", "Still limited on deep semantics"],
                    ],
                },
            },
            {
                "title": "Transformer Encoder-Decoder Overview",
                "kicker": "Bridge to Week 4 GenAI — encoder-decoder intuition before BERT/GPT",
                "layout": "diagram",
                "plot_path": "attention-architecture.png",
            },
            {
                "title": "Transformer Encoder-Decoder Overview",
                "kicker": "Bridge to Week 4 GenAI — encoder-decoder intuition before BERT/GPT",
                "layout": "diagram",
                "plot_path": "encoding-comparison.png",
            },
            {
                "title": "Transformer Encoder-Decoder Overview",
                "kicker": "Bridge to Week 4 GenAI — encoder-decoder intuition before BERT/GPT",
                "bullets": [
                    "Encoder builds contextual token representations from source text.",
                    "Decoder generates target tokens autoregressively with attention to encoder states.",
                    "Foundation architecture for translation, summarization, and many LLM pipelines.",
                ],
                "note": "Attention mechanisms reduce reliance on recurrence for sequence modeling.",
            },
        ],
    },
}

TAKEAWAYS = [
    "Language is ambiguous — cleaning and tokenization choices persist through the pipeline.",
    "N-grams and RNNs model sequences; attention removes the seq2seq bottleneck.",
    "Subword tokenizers and encoder–decoder attention are the bridge to Transformers.",
]

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "Week 3 Session 1 — NLP from raw text to Seq2Seq.",
    "current": "S8",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.05), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.5),
        Inches(11.5),
        Inches(0.85),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.5), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(3.85), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    topics = s["topics"]
    n_cols = 3
    chip_w = Inches(3.7)
    chip_h = Inches(0.48)
    gap_x = Inches(0.22)
    start_y = Inches(4.7)
    for i, topic in enumerate(topics):
        row, col = divmod(i, n_cols)
        x = MARGIN + col * (chip_w + gap_x)
        y = start_y + Inches(row * 0.68)
        soft_card(slide, x, y, chip_w, chip_h, fill=SOFT)
        add_text(slide, x, y + Inches(0.06), chip_w, Inches(0.36), topic, size=12, color=PRIMARY, align=PP_ALIGN.CENTER)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def content_title(slide, title: str, subtitle: str | None = None):
    """Title that shrinks on long headlines so it does not cover the kicker."""
    long = len(title or "") > 46
    size = 24 if long else 32
    title_h = 0.80 if long else 0.62
    add_text(
        slide,
        MARGIN,
        Inches(1.12),
        Inches(11.2),
        Inches(title_h),
        title,
        size=size,
        bold=True,
        color=PRIMARY,
    )
    if subtitle:
        add_text(
            slide,
            MARGIN,
            Inches(1.12 + title_h),
            Inches(11.2),
            Inches(0.34),
            subtitle,
            size=14 if long else 15,
            color=MUTED,
        )


def content_top(title: str) -> float:
    return 2.42 if len(title or "") > 46 else 2.18


def _chars_per_line(width_in: float, size: float) -> int:
    return max(22, int(width_in * (108 / max(size, 10))))


def _line_count(text: str, width_in: float, size: float) -> int:
    cpl = _chars_per_line(width_in, size)
    words = (text or "").split()
    if not words:
        return 1
    lines, cur = 1, 0
    for word in words:
        need = len(word) + (1 if cur else 0)
        if cur + need > cpl:
            lines += 1
            cur = len(word)
        else:
            cur += need
    return lines


def _bullet_box(items: list[str], width_in: float, size: float) -> tuple[float, float]:
    lines = max((_line_count(t, width_in, size) for t in items), default=1)
    item_h = max(0.42, min(1.15, lines * (size / 72) * 1.35 + 0.10))
    return item_h, item_h + 0.10


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
        if px_w <= 0 or px_h <= 0:
            return None
        aspect = px_w / px_h
        src = path
        if aspect < 0.5:
            from io import BytesIO

            crop_h = max(1, min(px_h, int(px_w / 2.4)))
            cropped = im.crop((0, 0, px_w, crop_h))
            buf = BytesIO()
            cropped.save(buf, format="PNG")
            buf.seek(0)
            src = buf
            px_w, px_h = cropped.size
            aspect = px_w / max(px_h, 1)
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(src) if src is path else src,
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    bsize = 14 if dense else 16

    if content.get("body"):
        top = content_top(content["title"])
        soft_card(slide, MARGIN, Inches(top), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(top + 0.20),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(top + 1.35)
        bsize = 14 if dense else 15
    else:
        bullet_top = Inches(content_top(content["title"]))
        bsize = 14 if dense else 15

    item_h, pitch = _bullet_box(items, 11.2, bsize)
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(item_h),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    has_body = bool((content.get("body") or "").strip())
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)

    y_after = content_top(content["title"])
    if has_body:
        soft_card(slide, MARGIN, Inches(y_after), col_w, Inches(1.2), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(y_after + 0.15),
            text_w,
            Inches(0.95),
            content.get("body") or "",
            size=13,
            color=INK,
        )
        y_after = y_after + 1.35

    if has_formula:
        formula_card_h = 1.65 if is_frac else 1.35
        formula_box_h = 0.95 if is_frac else 0.65
        soft_card(slide, MARGIN, Inches(y_after), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(y_after + 0.10),
            text_w,
            Inches(0.22),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(y_after + 0.35),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        if content.get("formula_note"):
            add_text(
                slide,
                Inches(MARGIN.inches + 0.3),
                Inches(y_after + 0.35 + formula_box_h + 0.02),
                text_w,
                Inches(0.28),
                content["formula_note"],
                size=10,
                color=MUTED,
            )
        bullet_top = y_after + formula_card_h + 0.12
        bsize = 12
    else:
        bullet_top = y_after
        bsize = 13

    item_h, pitch = _bullet_box(items, bullet_w.inches, bsize)
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch)) if items else 0
    if has_formula and has_plot:
        max_items = min(max_items, 2)
    if items:
        bullets(
            slide,
            items[:max_items],
            top=Inches(bullet_top),
            size=bsize,
            left=MARGIN,
            width=bullet_w,
            pitch=pitch,
            item_height=Inches(item_h),
        )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(content_top(content["title"]))
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(body_top.inches + 0.13),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(body_top.inches + 0.85)
    else:
        formula_top = body_top

    formula_h = 1.55 if is_frac else 1.25
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.12),
        Inches(11.2),
        Inches(0.95 if is_frac else 0.7),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )
    if content.get("formula_note"):
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(formula_top.inches + formula_h - 0.35),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    bullet_top = formula_top.inches + formula_h + 0.12
    bsize = 13 if dense else 15
    item_h, pitch = _bullet_box(items, 11.2, bsize)
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch)) if items else 0
    if has_note:
        max_items = min(max_items, 3)
    if items:
        bullets(
            slide,
            items[:max_items],
            top=Inches(bullet_top),
            size=bsize,
            pitch=pitch,
            width=Inches(11.2),
            item_height=Inches(item_h),
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = content_top(content["title"])
    if show_body:
        soft_card(slide, MARGIN, Inches(table_top), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(table_top + 0.10),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = table_top + 0.70

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    note_top = table_top + table_h + 0.08
    render_note = has_note and (note_top + 0.45) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    content_title(slide, content["title"], content.get("kicker"))

    diagram_top = Inches(content_top(content["title"]))
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(diagram_top.inches - 0.05),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(diagram_top.inches + 0.30)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(4.35), fill=SOFT)
    _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + 0.15),
        Inches(11.6),
        Inches(4.05),
        folder=folder,
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.9),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def slide_takeaways(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Takeaways", f"{index:02d}")
    title_block(slide, "3 Ideas to Keep", "Leave with these")
    for i, text in enumerate(TAKEAWAYS):
        y = Inches(2.35) + Inches(i * 1.25)
        soft_card(slide, MARGIN, y, Inches(12.0), Inches(1.08), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            y + Inches(0.28),
            Inches(0.5),
            Inches(0.5),
            str(i + 1),
            size=22,
            bold=True,
            color=PRIMARY,
        )
        add_text(
            slide,
            MARGIN + Inches(0.95),
            y + Inches(0.28),
            Inches(10.6),
            Inches(0.6),
            text,
            size=16,
            color=INK,
        )
    content_footer(slide, index, total)


def slide_close(prs, total, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))
    add_text(
        slide,
        MARGIN,
        Inches(2.9),
        Inches(12),
        Inches(0.8),
        "Thank you",
        size=40,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN,
        Inches(3.8),
        Inches(12),
        Inches(0.4),
        "Labs: code/16- NLP",
        size=16,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    content_footer(slide, index, total)


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 6 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    n += 1
    slide_takeaways(prs, total, n)
    n += 1
    slide_close(prs, total, n)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
