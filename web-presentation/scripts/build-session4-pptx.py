#!/usr/bin/env python3
"""
Build Week 1 · Session 4 presentation using official ETRA Design System.
Content is filled incrementally as teaching blocks are pasted in.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week1-Session4-Naive-Bayes-Trees-Evaluation.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session4-diagrams"

SESSION = {
    "eyebrow": "Week 1  ·  Session 4",
    "section_number": "04",
    "section_tag": "S04",
    "section_title": "Classical ML",
    "subtitle": "Naive Bayes, trees, and evaluation metrics",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "Naive Bayes, Trees, Evaluation",
    "topics": [
        "Naive Bayes",
        "Decision trees",
        "Random Forest",
        "Precision / recall / F1",
    ],
}

TOPIC_CONTENT: dict = {
    "Naive Bayes": {
        "title": "Naive Bayes: Bayes' Theorem",
        "kicker": "Core Probability Rule",
        "body": "Bayes’ theorem updates our belief about A after observing evidence B.",
        "formula": "P(A|B) = P(B|A) P(A) / P(B)",
        "formula_tex": r"P(A\mid B)=\dfrac{P(B\mid A)\,P(A)}{P(B)}",
        "formula_note": "Posterior ∝ Likelihood × Prior",
        "layout": "formula_example",
        "bullets": [
            "P(A|B): posterior probability after observing evidence B.",
            "P(A): prior belief before seeing evidence.",
            "P(B|A): likelihood of evidence if class A is true.",
        ],
        "note": "P(B): evidence probability (normalization term). This theorem updates beliefs when new evidence appears.",
        "extra_slides": [
            {
                "title": "Bayes’ Theorem — Visual Map",
                "kicker": "Core Probability Rule",
                "body": "Prior × Likelihood → Posterior (normalized by evidence).",
                "layout": "diagram",
                "plot_path": "bayes-theorem.png",
                "note": "Same rule as the formula slide — different view of the update.",
            },
            {
                "title": "Why Is It Called 'Naive'?",
                "kicker": "Conditional Independence Assumption",
                "body": "Naive Bayes assumes features are conditionally independent given the class label.",
                "bullets": [
                    "Simplifies parameter estimation and speeds up training.",
                    "Works well in text classification and high-dimensional sparse data.",
                    "Can still perform well even when independence is not fully true.",
                ],
                "plot_path": "naive-independence.png",
                "note": "When comparing classes for the same X, P(X) is constant and can be ignored.",
            },
            {
                "title": "Naive Bayes: Practical Notes",
                "kicker": "Variants & Implementation",
                "layout": "table",
                "table": {
                    "headers": ["Variant", "Typical Data", "Key Idea"],
                    "rows": [
                        [
                            "Gaussian NB",
                            "Continuous",
                            "Each feature ~ Gaussian given the class",
                        ],
                        [
                            "Multinomial NB",
                            "Count-based",
                            "Common for document word counts",
                        ],
                        [
                            "Bernoulli NB",
                            "Binary",
                            "Presence/absence features",
                        ],
                    ],
                },
                "note": "Handle zero-frequency with smoothing (e.g., Laplace). Scale/preprocess by feature type and variant.",
            },
        ],
    },
    "Decision trees": {
        "title": "Decision Tree Classification (CART)",
        "kicker": "Recursive Partitioning",
        "body": "Decision Tree Classification predicts class labels by recursively splitting feature space into purer class regions.",
        "bullets": [
            "Internal nodes represent feature-based decisions.",
            "Branches represent decision outcomes.",
            "Leaf nodes output final class prediction.",
        ],
        "plot_path": "decision-tree-split.png",
        "note": "Decision Tree Split · Root → branches → leaf nodes · split by feature threshold.",
        "extra_slides": [
            {
                "title": "How Splits Are Chosen in Classification Trees",
                "kicker": "Gini · Entropy · Information Gain",
                "body": "Splits are chosen greedily to maximize information gain.",
                "formula": "Gain = Impurity(parent) − weighted impurity(children)",
                "formula_tex": r"\mathrm{Gain}=I(\mathrm{parent})-\sum_k \frac{|S_k|}{|S|}I(S_k)",
                "formula_note": "Gini / Entropy measure impurity — lower is better",
                "bullets": [
                    "Choose the split that best separates classes at each node.",
                    "Stop splitting using rules like max depth or minimum samples.",
                    "Gini impurity and entropy measure node purity — lower is better.",
                ],
                "layout": "formula_example",
                "note": "Better splits reduce impurity and increase class purity in child nodes.",
            },
            {
                "title": "Decision Tree: Classification vs Regression",
                "kicker": "Same Structure · Different Objectives",
                "layout": "table",
                "table": {
                    "headers": ["Aspect", "Classification Tree", "Regression Tree"],
                    "rows": [
                        ["Target", "Categorical class", "Continuous value"],
                        ["Split criterion", "Gini / Entropy", "MSE / MAE"],
                        [
                            "Leaf output",
                            "Class label or class probability",
                            "Numeric mean/median",
                        ],
                    ],
                },
                "note": "Both use recursive partitioning (Root → branches → leaves) but optimize different objectives.",
            },
            {
                "title": "Gini Impurity and Entropy",
                "kicker": "Node Purity Measures",
                "body": "Gini (CART default) and entropy both measure node impurity — lower is better.",
                "layout": "diagram",
                "plot_path": "gini-entropy-formulas.png",
                "note": "Better splits reduce impurity and increase class purity in child nodes.",
            },
        ],
    },
    "Random Forest": {
        "title": "Random Forest Classification",
        "kicker": "Ensemble Learning by Bagging",
        "body": "Train many trees on bootstrap samples, then combine predictions by majority vote.",
        "formula": "ŷ_RF(x) = mode{ŷ_1(x), ŷ_2(x), …, ŷ_T(x)}",
        "formula_tex": r"\hat{y}_{RF}(x)=\operatorname{mode}\{\hat{y}_1(x),\hat{y}_2(x),\ldots,\hat{y}_T(x)\}",
        "formula_note": "Majority vote across T trees",
        "bullets": [
            "Sample bootstrap data and train many trees.",
            "Each tree predicts a class for the new sample.",
        ],
        "plot_path": "random-forest-ensemble.png",
        "note": "Final class = majority vote (mode). Regression RF uses the mean instead.",
        "extra_slides": [
            {
                "title": "Why Random Forest Often Outperforms a Single Tree",
                "kicker": "Bias–Variance Tradeoff",
                "layout": "table",
                "table": {
                    "headers": ["Property", "Single Decision Tree", "Random Forest"],
                    "rows": [
                        ["Variance", "High", "Lower (averaging effect)"],
                        ["Overfitting risk", "Higher", "Lower"],
                        ["Interpretability", "High", "Moderate"],
                        [
                            "Predictive robustness",
                            "Sensitive to data noise",
                            "More stable",
                        ],
                    ],
                },
                "note": "Bagging + feature randomness → less variance without much extra bias.",
            },
        ],
    },
    "Precision / recall / F1": {
        "title": "Classification Errors: FP and FN",
        "kicker": "Error Types",
        "body": "Not all mistakes are equal — False Positives and False Negatives carry different costs.",
        "bullets": [
            "False Positive (Type I): predict positive, actual is negative.",
            "False Negative (Type II): predict negative, actual is positive.",
            "Error impact depends on domain costs (healthcare vs spam).",
        ],
        "plot_path": "fp-fn-errors.png",
        "note": "Model evaluation should consider error trade-offs, not accuracy alone.",
        "extra_slides": [
            {
                "title": "Confusion Matrix and Accuracy",
                "kicker": "Actual vs Predicted Counts",
                "body": "Confusion matrix counts TP, FP, FN, TN — Accuracy is overall correctness.",
                "formula": "Accuracy = (TP + TN) / (TP + TN + FP + FN)",
                "formula_tex": r"\mathrm{Accuracy}=\dfrac{TP+TN}{TP+TN+FP+FN}",
                "formula_note": "Correct predictions ÷ all predictions",
                "bullets": [
                    "TP / TN: correct positives / negatives.",
                    "FP / FN: false alarms / missed detections.",
                ],
                "plot_path": "confusion-matrix-accuracy.png",
                "note": "High accuracy can hide poor minority-class performance when data is imbalanced.",
            },
        ],
    },
}

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "After classification basics — probabilistic models, trees, and metrics.",
    "current": "S4",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.15), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.65),
        Inches(11.5),
        Inches(0.9),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.7), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.05), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    x = MARGIN
    for topic in s["topics"]:
        w = Inches(min(2.8, 0.12 * len(topic) + 1.1))
        if x + w > Inches(12.4):
            break
        soft_card(slide, x, Inches(5.25), w, Inches(0.42), fill=SOFT)
        add_text(slide, x, Inches(5.3), w, Inches(0.32), topic, size=11, color=PRIMARY, align=PP_ALIGN.CENTER)
        x += w + Inches(0.14)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    """Fit a plot PNG inside a box (same idea as add_diagram)."""
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return None
    aspect = px_w / px_h
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(path),
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    pitch = 0.50 if dense else 0.62
    bsize = 14 if dense else 16

    if content.get("body"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.45),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(3.6)
    else:
        bullet_top = Inches(2.35)
        bsize = 15 if dense else 17

    # Keep bullets above the note/footer band
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    """Definition + formula + optional scatter/line plot."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)

    soft_card(slide, MARGIN, Inches(2.2), col_w, Inches(1.2), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.3),
        Inches(2.35),
        text_w,
        Inches(0.95),
        content.get("body") or "",
        size=13,
        color=INK,
    )

    if has_formula:
        formula_card_h = 1.65 if is_frac else 1.35
        formula_box_h = 0.95 if is_frac else 0.65
        soft_card(slide, MARGIN, Inches(3.55), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(3.65),
            text_w,
            Inches(0.22),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(3.9),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        if content.get("formula_note"):
            note_y = 3.9 + formula_box_h + 0.05
            add_text(
                slide,
                Inches(MARGIN.inches + 0.3),
                Inches(note_y),
                text_w,
                Inches(0.3),
                content["formula_note"],
                size=10,
                color=MUTED,
            )
        bullet_top = 3.55 + formula_card_h + 0.12
        pitch = 0.46 if (has_note or len(items) >= 2) else 0.55
        bsize = 12
    else:
        bullet_top = 3.55
        pitch = 0.50 if (has_note or len(items) >= 3) else 0.58
        bsize = 13

    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    # With formula+plot keep at most 2 bullets to avoid crowding
    if has_formula and has_plot:
        max_items = min(max_items, 2)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        left=MARGIN,
        width=bullet_w,
        pitch=pitch,
        item_height=Inches(0.42),
    )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    """Full-width practical example with formula and interpretation."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(2.15)
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.28),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(3.0)
    else:
        formula_top = Inches(2.25)

    formula_h = 1.55 if is_frac else 1.25
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.12),
        Inches(11.2),
        Inches(0.95 if is_frac else 0.7),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )
    if content.get("formula_note"):
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(formula_top.inches + formula_h - 0.35),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    bullet_top = formula_top.inches + formula_h + 0.12
    pitch = 0.46 if dense else 0.56
    bsize = 13 if dense else 15
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    if has_note:
        max_items = min(max_items, 3)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(0.42),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    # Prefer kicker over a long body when the table is dense
    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = 2.25
    if show_body:
        soft_card(slide, MARGIN, Inches(2.15), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.25),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = 2.85

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    # If note would still collide, drop the note rather than overlap the table
    note_top = table_top + table_h + 0.08
    render_note = has_note and (note_top + 0.45) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    """Title + optional short body + full-width diagram."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    diagram_top = Inches(2.15)
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(2.1),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(2.45)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(4.35), fill=SOFT)
    _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + 0.15),
        Inches(11.6),
        Inches(4.05),
        folder=folder,
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.9),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    """Returns number of slides added."""
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(4.25),
        Inches(11.2),
        Inches(0.4),
        "Content coming next — paste the teaching block when ready.",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 4 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
