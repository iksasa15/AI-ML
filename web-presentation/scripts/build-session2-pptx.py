#!/usr/bin/env python3
"""
Build Week 1 · Session 2 presentation using official ETRA Design System.
Content is filled incrementally as teaching blocks are pasted in.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))
from etra_brand import (  # noqa: E402
    INK,
    MARGIN,
    MUTED,
    PRIMARY,
    SECONDARY,
    SLIDE_H,
    SLIDE_W,
    SOFT,
    SOFT_2,
    WHITE,
    add_formula,
    add_text,
    bullets,
    content_footer,
    content_header,
    gradient_fill,
    is_fraction_formula,
    logo,
    paint_light,
    rect,
    right_rail,
    soft_card,
    title_block,
)

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "pdf-exports" / "Week1-Session2-Regression-Models.pptx"
PLOTS = ROOT / "public" / "assets" / "plots"
DIAGRAMS = ROOT / "public" / "assets" / "session2-diagrams"

SESSION = {
    "eyebrow": "Week 1  ·  Session 2",
    "section_number": "02",
    "section_tag": "S02",
    "section_title": "Regression Models",
    "subtitle": "From linear fits to regularized models",
    "trainer_line": "AI & Machine Learning Bootcamp",
    "focus": "Regression Models",
    "topics": [
        "Linear regression",
        "OLS & R²",
        "Residual diagnostics",
        "Polynomial regression",
        "SVR",
        "Decision Tree",
        "Random Forest",
        "Ridge / Lasso",
    ],
}

TOPIC_CONTENT: dict = {
    "Linear regression": {
        "title": "What is Simple Linear Regression?",
        "kicker": "Regression · Simple Linear Regression",
        "body": "Simple Linear Regression models the relationship between one independent variable x and one dependent variable y using a straight line.",
        "formula": "y = b_{0} + b_{1}x",
        "formula_note": "b₀ = intercept (value of y when x = 0)  ·  b₁ = slope (change in y per unit change in x)",
        "bullets": [
            "x = input feature · y = target to predict",
            "Finds the straight line that best fits the data.",
        ],
        "plot": "slide-19-1.png",
        "note": "One feature → one straight-line fit. Multiple features come next.",
        "extra_slides": [
            {
                "title": "Practical Example",
                "kicker": "Simple Linear Regression",
                "body": "Suppose x = Hours of Study and y = Exam Score.",
                "formula": "y = 40 + 5x",
                "layout": "formula_example",
                "bullets": [
                    "Slope b₁ = 5: +1 study hour → +5 score points.",
                    "Intercept b₀ = 40: score when hours = 0.",
                    "Example: 3 hours → ŷ = 40 + 5(3) = 55.",
                ],
            },
            {
                "title": "Meaning of Each Symbol",
                "kicker": "Simple Linear Regression",
                "layout": "table",
                "table": {
                    "headers": ["Symbol", "Meaning"],
                    "rows": [
                        ["y", "Predicted output (dependent variable)."],
                        ["x", "Input feature (independent variable)."],
                        ["b₀", "Intercept: predicted value of y when x = 0."],
                        ["b₁", "Slope: expected change in y when x increases by 1 unit."],
                    ],
                },
                "note": "The regression line is the best-fit line that minimizes prediction errors.",
            },
            {
                "title": "Multiple Linear Regression",
                "kicker": "Linear Regression · Multiple Features",
                "body": "Multiple Linear Regression describes the relationship between one dependent variable y and two or more independent variables.",
                "formula": "y = b_{0} + b_{1}x_{1} + b_{2}x_{2} + ⋯ + b_{n}x_{n}",
                "formula_note": "Each bⱼ is the effect of feature xⱼ holding other features fixed.",
                "layout": "formula_example",
                "bullets": [
                    "Predict y from several inputs at once.",
                    "Still linear in the coefficients b₀ … bₙ.",
                ],
                "note": "OLS still minimizes the sum of squared residuals.",
            },
            {
                "title": "House Price Example",
                "kicker": "Multiple Linear Regression",
                "body": "Predict house price from size and bedrooms.",
                "formula": "y = 50 + 2.5 x_{1} + 15 x_{2}",
                "formula_note": "y = price  ·  x₁ = size (m²)  ·  x₂ = bedrooms",
                "layout": "formula_example",
                "bullets": [
                    "+1 m² → price +2.5 (other variables fixed).",
                    "+1 bedroom → price +15 (other variables fixed).",
                ],
                "note": "Interpret each coefficient holding other variables constant.",
            },
            {
                "title": "House Price Calculations",
                "kicker": "Multiple Linear Regression · Example",
                "layout": "table",
                "table": {
                    "headers": [
                        "House",
                        "Size x₁ (m²)",
                        "Bedrooms x₂",
                        "Predicted price y ($1000)",
                    ],
                    "rows": [
                        ["A", "100", "2", "50 + 2.5(100) + 15(2) = 330"],
                        ["B", "120", "3", "50 + 2.5(120) + 15(3) = 395"],
                        ["C", "80", "1", "50 + 2.5(80) + 15(1) = 265"],
                    ],
                },
                "note": "Same formula, different inputs → different predicted prices.",
            },
        ],
    },
    "OLS & R²": {
        "title": "What is OLS?",
        "kicker": "Ordinary Least Squares (OLS)",
        "body": "Ordinary Least Squares is the most common method to fit a linear regression line by minimizing prediction errors.",
        "bullets": [
            "Chooses the line with the smallest total squared error.",
            "Residual = actual y − predicted ŷ.",
            "Simple linear case has a closed-form solution for b₀, b₁.",
        ],
        "note": "Smaller squared errors → better fit on the training data.",
        "extra_slides": [
            {
                "title": "Sum of Squared Errors",
                "kicker": "Ordinary Least Squares (OLS)",
                "body": "OLS chooses parameters by minimizing the total squared difference between actual and predicted values.",
                "formula": "min Σ_{i = 1 … n} (y_{i} − ŷ_{i})²",
                "formula_tex": r"\min \sum_{i=1}^{n} (y_i - \hat{y}_i)^2",
                "formula_note": "Sum of squared residuals over all n training points",
                "layout": "formula_example",
                "bullets": [
                    "Actual value: yᵢ",
                    "Predicted value: ŷᵢ",
                    "Residual error: eᵢ = yᵢ − ŷᵢ",
                ],
                "note": "Squared errors make all errors positive, penalize large errors, and provide a clear optimization objective.",
            },
            {
                "title": "Gradient Descent",
                "kicker": "OLS · Optimization",
                "body": "Loss curve L(w) — minimum at the bottom of the bowl.",
                "formula": "min_{w} L(w)",
                "formula_note": "We search for the parameter w that makes the loss as small as possible.",
                "bullets": [
                    "Start with a guess for w, then move downhill.",
                    "Stop when updates become tiny (converged).",
                ],
                "plot_path": "ols-loss-curve.png",
                "note": "Converged: w = 2.00 · L(w) = 0.40",
            },
            {
                "title": "R-Squared (R²)",
                "kicker": "Goodness of Fit",
                "body": "R² measures how much variation in y is explained by the model inputs.",
                "formula": "R² = 1 − (SS_{res}) / (SS_{tot})",
                "formula_tex": r"R^{2} = 1 - \dfrac{SS_{res}}{SS_{tot}}",
                "formula_note": "SS_res = residual SS  ·  SS_tot = total SS",
                "bullets": [
                    "Near 1 → strong explanatory power.",
                    "Near 0 → little better than predicting the mean.",
                ],
                "plot": "slide-22-1.png",
                "note": "High R² on training data ≠ guaranteed good new-data predictions.",
            },
            {
                "title": "Interpreting R²",
                "kicker": "Goodness of Fit",
                "body": "Larger R² means more of the variation in y is captured by the model.",
                "bullets": [
                    "R² = 1 → perfect fit on that sample.",
                    "R² = 0 → mean-level baseline performance.",
                ],
                "plot": "r2-comparison.png",
                "note": "Validate generalization with hold-out metrics, not R² alone.",
            },
            {
                "title": "Residual Sum of Squares",
                "kicker": "R² · Building Blocks",
                "body": "SS_res measures the unexplained error — how far predictions are from the actual values.",
                "formula": "SS_{res} = Σ_{i = 1 … n} (y_{i} − ŷ_{i})²",
                "formula_note": "Sum of squared residuals (prediction errors)",
                "layout": "formula_example",
                "bullets": [
                    "Smaller SS_res → predictions closer to the data.",
                    "Appears in the numerator of the R² formula.",
                ],
            },
            {
                "title": "Total Sum of Squares",
                "kicker": "R² · Building Blocks",
                "body": "SS_tot measures the total variation in y around its mean — the baseline scatter to explain.",
                "formula": "SS_{tot} = Σ_{i = 1 … n} (y_{i} − ȳ)²",
                "formula_note": "Sum of squared deviations from the mean ȳ",
                "layout": "formula_example",
                "bullets": [
                    "This is the variance (up to a constant factor) of the observed y values.",
                    "R² asks: what fraction of SS_tot is left unexplained by SS_res?",
                ],
                "note": "R² = 1 − SS_res / SS_tot  ·  explained variation = SS_tot − SS_res",
            },
            {
                "title": "What Do These Terms Mean?",
                "kicker": "R² · Building Blocks",
                "layout": "table",
                "table": {
                    "headers": ["Term", "Description"],
                    "rows": [
                        [
                            "SS_res",
                            "Residual Sum of Squares: squared distance between actual values and predictions.",
                        ],
                        [
                            "SS_tot",
                            "Total Sum of Squares: squared distance between actual values and the mean of y.",
                        ],
                    ],
                },
                "note": "R² compares unexplained error (SS_res) with total variation (SS_tot).",
            },
            {
                "title": "Why R² Alone Is Not Enough",
                "kicker": "Adjusted R-Squared",
                "body": "In multiple regression, adding predictors often increases R² even when the new variables are not useful.",
                "bullets": [
                    "Ordinary R² never decreases when you add a predictor (it stays the same or goes up).",
                    "That can reward bloated models that fit noise instead of signal.",
                    "We need a metric that balances fit against model complexity.",
                ],
                "note": "Next: Adjusted R² applies a penalty for extra predictors.",
            },
            {
                "title": "Adjusted R-Squared",
                "kicker": "Goodness of Fit · Model Comparison",
                "body": "Adjusted R² penalizes unnecessary predictors.",
                "formula": "R²_{adj} = 1 − (1 − R²)(n − 1)/(n − p − 1)",
                "formula_tex": r"R^{2}_{adj} = 1 - \dfrac{(1 - R^{2})(n - 1)}{n - p - 1}",
                "formula_note": "n = samples  ·  p = predictors",
                "layout": "formula_example",
                "bullets": [
                    "Better for comparing models with different p.",
                    "Can fall when R² rises if the new variable is weak.",
                ],
            },
            {
                "title": "Adjusted R² — Quick Example",
                "kicker": "Adjusted R-Squared",
                "layout": "table",
                "table": {
                    "headers": [
                        "Model",
                        "Predictors",
                        "R²",
                        "Adjusted R²",
                        "Interpretation",
                    ],
                    "rows": [
                        ["A", "x₁, x₂", "0.82", "0.81", "Strong and efficient model"],
                        [
                            "B",
                            "x₁, x₂, x₃, x₄",
                            "0.83",
                            "0.79",
                            "Higher R², but worse after penalty",
                        ],
                    ],
                },
                "note": "Model B looks better by R², but Adjusted R² shows the extra variables are likely not helpful.",
            },
        ],
    },
    "Residual diagnostics": {
        "title": "Assumptions of Linear Regression",
        "kicker": "Overview",
        "layout": "table",
        "table": {
            "headers": ["Assumption", "Meaning", "If Violated"],
            "rows": [
                ["Linearity", "y vs each x is roughly linear", "Biased predictions"],
                ["Homoscedasticity", "Error variance roughly constant", "Unstable SEs"],
                ["Normal Errors", "Residuals ≈ normal", "Weaker inference"],
                ["Independence", "Errors are independent", "Biased tests"],
                ["Low Multicollinearity", "Predictors not highly redundant", "Unstable coeffs"],
                ["Limited Outliers", "Extremes do not dominate fit", "Distorted line"],
            ],
        },
        "note": "Next slides: practical checks and residual plots.",
        "extra_slides": [
            {
                "title": "Assumptions with Practical Examples",
                "kicker": "Residual Diagnostics",
                "layout": "diagram",
                "body": "Each panel shows a practical picture of one assumption — and what “good” looks like.",
                "plot_path": "assumption-examples-grid.png",
                "note": "Checks: scatter/residuals · residuals vs fitted · histogram/Q-Q · design/DW · corr/VIF · Cook/leverage.",
            },
            {
                "title": "Assumptions — Quick Reference",
                "kicker": "Practical Examples",
                "layout": "table",
                "table": {
                    "headers": ["Assumption", "Example", "Check"],
                    "rows": [
                        ["Linearity", "Sales ≈ linear in ad spend", "Scatter / residual plot"],
                        ["Homoscedasticity", "Similar errors high & low", "Residuals vs fitted"],
                        ["Normal Errors", "Most residuals near zero", "Histogram / Q-Q"],
                        ["Independence", "Independent customers", "Design / Durbin-Watson"],
                        ["Low Multicollinearity", "Ads not near-duplicates", "Corr / VIF"],
                        ["Outlier Control", "One house ≠ whole model", "Cook / leverage"],
                    ],
                },
                "note": "If violated: transforms, robust methods, or another model.",
            },
            {
                "title": "Residual Patterns to Spot",
                "kicker": "Assumption Checks",
                "body": "Common residual-plot signatures for assumption problems.",
                "bullets": [
                    "Good: random cloud around zero.",
                    "Fan / curve / outlier → investigate.",
                ],
                "plot_path": "residual-assumptions.png",
            },
            {
                "title": "Residuals vs Fitted Values",
                "kicker": "Practical Check · Homoscedasticity",
                "body": "Look for random scatter around zero with roughly constant spread.",
                "bullets": [
                    "No curve → linearity looks OK.",
                    "Fan shape → possible heteroscedasticity.",
                ],
                "plot": "slide-25-1.png",
            },
            {
                "title": "Building a Regression Model",
                "kicker": "Goal",
                "layout": "table",
                "table": {
                    "headers": ["Method", "Core Idea", "Best Use Case"],
                    "rows": [
                        ["All-in", "Use all predictors", "Strong prior knowledge"],
                        ["Backward", "Start full, remove weakest", "Many candidates"],
                        ["Forward", "Start empty, add strongest", "Compact model"],
                        ["Stepwise", "Add and remove dynamically", "Flexible search"],
                    ],
                },
                "note": "Match the method to candidate count and prior knowledge.",
            },
            {
                "title": "Model-Building Flow",
                "kicker": "All-in · Backward · Forward · Stepwise",
                "body": "Different search paths for choosing predictors under a significance level.",
                "bullets": [
                    "Backward: prune a full model.",
                    "Forward: grow from empty.",
                    "Stepwise: add and remove as you go.",
                ],
                "plot_path": "model-building-flow.png",
            },
            {
                "title": "Backward Elimination",
                "kicker": "Model-Building Methods · Steps",
                "body": "Start full, then drop the least useful predictors one at a time.",
                "bullets": [
                    "Set SL (e.g. 0.05) and fit all predictors.",
                    "Remove highest p if p > SL; refit and repeat.",
                ],
                "note": "Best with many candidate predictors.",
            },
            {
                "title": "Forward Selection",
                "kicker": "Model-Building Methods · Steps",
                "body": "Grow the model from scratch by adding the strongest predictors.",
                "bullets": [
                    "Start empty; add lowest p below SL.",
                    "Stop when no remaining variable qualifies.",
                ],
                "note": "Best when you want a compact model from the start.",
            },
            {
                "title": "Bidirectional (Stepwise)",
                "kicker": "Model-Building Methods · Steps",
                "body": "Add and remove predictors dynamically each iteration.",
                "bullets": [
                    "Forward step: add significant variables.",
                    "Backward step: drop non-significant ones.",
                ],
                "note": "More adaptive than pure forward or backward alone.",
            },
            {
                "title": "Mini Example — Backward Elimination",
                "kicker": "Model-Building Methods",
                "layout": "table",
                "table": {
                    "headers": ["Iteration", "Candidate Predictors", "Selected Action"],
                    "rows": [
                        ["Start", "x₁, x₂, x₃, x₄", "Full model"],
                        ["1", "highest p-value: x₄ = 0.42", "Remove x₄"],
                        ["2", "highest p-value: x₃ = 0.11", "Remove x₃"],
                        ["Final", "x₁, x₂ with p < 0.05", "Keep final model"],
                    ],
                },
                "note": "With SL = 0.05, variables with p > SL are removed until only significant predictors remain.",
            },
            {
                "title": "Significance Level and p-value",
                "kicker": "Model Selection · Decision Rule",
                "body": "Compare p to alpha to decide if a predictor is statistically significant.",
                "formula": "0.03 < 0.05  ⇒  Reject H_{0}",
                "formula_note": "Example: p = 0.03 · alpha = 0.05",
                "bullets": [
                    "p < alpha → reject H₀ (significant).",
                    "p ≥ alpha → fail to reject H₀.",
                ],
                "plot_path": "pvalue-alpha.png",
                "note": "Predictor is significant at the 5% level in this example.",
            },
            {
                "title": "Key Concepts: H₀, p-value, Alpha",
                "kicker": "Significance Level and p-value",
                "layout": "table",
                "table": {
                    "headers": ["Concept", "Meaning"],
                    "rows": [
                        ["Null Hypothesis (H₀)", "No effect / no relationship"],
                        [
                            "p-value",
                            "Probability of observing current result (or more extreme) if H₀ is true",
                        ],
                        [
                            "Significance Level (alpha)",
                            "Maximum tolerated Type I error probability",
                        ],
                    ],
                },
                "note": "In model building, predictors with p < alpha are kept as statistically significant.",
            },
            {
                "title": "Why Is SL = 0.05 Common?",
                "kicker": "Standard Choice",
                "body": "SL = 0.05 balances sensitivity and false-positive risk.",
                "formula": "SL = 1 − Confidence Level",
                "formula_note": "95% confidence → SL = 0.05",
                "layout": "formula_example",
                "bullets": [
                    "Convention, not a law — adjust for domain risk.",
                    "Report effect size and CIs with p-values.",
                ],
            },
            {
                "title": "SL vs Confidence Level",
                "kicker": "Why Is SL = 0.05 Common?",
                "layout": "table",
                "table": {
                    "headers": ["SL", "Confidence Level", "Strictness"],
                    "rows": [
                        ["0.10", "90%", "Less strict"],
                        ["0.05", "95%", "Standard"],
                        ["0.01", "99%", "More strict"],
                    ],
                },
                "note": "Lower SL → higher confidence requirement → stricter evidence to keep a predictor.",
            },
            {
                "title": "How to Choose in Practice",
                "kicker": "Significance Level",
                "layout": "table",
                "table": {
                    "headers": ["Scenario", "Suggested SL"],
                    "rows": [
                        [
                            "High-stakes decisions (medicine, safety, finance)",
                            "0.01",
                        ],
                        ["General modeling and reporting", "0.05"],
                        ["Exploratory analysis", "0.10"],
                    ],
                },
                "note": "Start with SL = 0.05, then adjust based on domain risk, dataset size, and model purpose.",
            },
        ],
    },
    "Polynomial regression": {
        "title": "Polynomial Regression",
        "kicker": "Definition",
        "body": "Polynomial Regression extends linear regression to model non-linear relationships by adding polynomial terms.",
        "formula": "y = b_{0} + b_{1}x + b_{2}x² + b_{3}x³ + ⋯",
        "formula_note": "Still linear in the coefficients",
        "layout": "formula_example",
        "bullets": [
            "Adds polynomial terms to capture curves.",
            "Still linear in parameters b₀, b₁, …",
        ],
        "note": "Transform x → x², x³, … then fit linear regression.",
        "extra_slides": [
            {
                "title": "Why Polynomial Regression?",
                "kicker": "Motivation",
                "body": "Many real-world relationships are curved — a straight line underfits them.",
                "bullets": [
                    "Population growth · disease progression",
                    "Sales trends · rainfall vs crop yield",
                ],
            },
            {
                "title": "How Polynomial Regression Works",
                "kicker": "Polynomial Regression",
                "body": "y = b₀ + b₁x + b₂x² + ⋯ + b_d x^d",
                "layout": "table",
                "table": {
                    "headers": ["Step", "Description"],
                    "rows": [
                        ["1", "Start with input feature x"],
                        ["2", "Create transformed terms: x², x³, …"],
                        ["3", "Fit linear regression on transformed features"],
                        ["4", "Estimate coefficients (b₀, b₁, b₂, …)"],
                        ["5", "Use the fitted equation to predict y"],
                    ],
                },
                "note": "d is the polynomial degree.",
            },
            {
                "title": "Choosing the Polynomial Degree",
                "kicker": "Model Complexity",
                "body": "Degree controls complexity — too low underfits, too high overfits.",
                "bullets": [
                    "Pick degree with CV / validation error.",
                    "Also compare adjusted R².",
                ],
                "plot_path": "poly-degree-tradeoff.png",
                "note": "Balanced degree → better bias-variance trade-off.",
            },
            {
                "title": "When to Use Polynomial Regression",
                "kicker": "Good Use Cases",
                "body": "Data shows clear curved trends · more flexibility than a straight line · you still want interpretable behavior.",
                "layout": "table",
                "table": {
                    "headers": ["Application", "Example"],
                    "rows": [
                        ["Forecasting", "Demand and sales trends"],
                        ["Environment", "Rainfall vs crop yield"],
                        ["Medicine", "Disease / progression patterns"],
                        ["Engineering", "Stress-strain and similar curves"],
                    ],
                },
                "note": "Choose polynomial regression when curves matter, but full black-box models are not required yet.",
            },
            {
                "title": "Limitations of Polynomial Regression",
                "kicker": "Trade-offs",
                "layout": "table",
                "table": {
                    "headers": ["Limitation", "Explanation"],
                    "rows": [
                        [
                            "Overfitting risk",
                            "High-degree models may memorize noise",
                        ],
                        [
                            "Outlier sensitivity",
                            "Extreme points can shift the curve strongly",
                        ],
                        [
                            "Extrapolation weakness",
                            "Predictions outside training range can be unreliable",
                        ],
                        [
                            "Limited flexibility for very complex patterns",
                            "Some tasks are better handled by tree-based or neural models",
                        ],
                    ],
                },
                "note": "Use polynomial regression as a baseline, then compare with alternatives as complexity grows.",
            },
            {
                "title": "Polynomial Regression Visual Explanation",
                "kicker": "Linear vs Polynomial Fits",
                "body": "Compare linear vs polynomial fits as degree changes.",
                "bullets": [
                    "Too low degree → high bias (underfit).",
                    "Too high degree → high variance (overfit).",
                ],
                "plot": "slide-39-1.png",
            },
            {
                "title": "Effect of Polynomial Degree",
                "kicker": "Visual Explanation",
                "body": "Higher degree = more flexible curve — useful, then noisy.",
                "bullets": [
                    "Pick degree with validation curves.",
                    "Do not trust training error alone.",
                ],
                "plot": "slide-39-2.png",
            },
        ],
    },
    "SVR": {
        "title": "Support Vector Machine (SVM) and SVR",
        "kicker": "What is SVR?",
        "body": "SVR is the regression version of SVM: a flat function with a tolerance margin (ε-tube).",
        "bullets": [
            "SVM: widest margin; support vectors define it.",
            "SVR ignores small errors inside the ε-tube.",
            "Used in forecasting, engineering curves, disease progression.",
        ],
        "plot_path": "svr-epsilon-tube.png",
        "note": "SVM separates classes; SVR predicts continuous y with an ε-insensitive zone.",
        "extra_slides": [
            {
                "title": "SVR Formulation and Equations",
                "kicker": "Model Form",
                "body": "SVR predicts with a (possibly kernelized) linear function in a feature space.",
                "formula": "f(x) = wᵀ φ(x) + b",
                "formula_note": "φ(x) maps input to a (possibly) higher-dimensional space",
                "layout": "formula_example",
                "bullets": [
                    "φ(x) maps input to a (possibly) higher-dimensional space.",
                    "With kernels, we avoid computing φ(x) explicitly.",
                ],
            },
            {
                "title": "Soft-Margin SVR Objective",
                "kicker": "Optimization Objective",
                "body": "Keep the function flat while penalizing points outside the ε-tube.",
                "formula": "min ½∥w∥² + C Σ (ξᵢ + ξᵢ*)",
                "formula_note": "C trades flatness vs errors outside the tube",
                "layout": "formula_example",
                "bullets": [
                    "½∥w∥² → flatter function.",
                    "Larger C → tighter fit; smaller C → smoother.",
                ],
            },
            {
                "title": "SVR Constraints",
                "kicker": "Soft-Margin SVR",
                "body": "Stay within ε of yᵢ, or pay with a slack variable.",
                "formula": "yᵢ − (wᵀφ(xᵢ) + b) ≤ ε + ξᵢ",
                "formula_note": "Also: (wᵀφ(xᵢ)+b) − yᵢ ≤ ε + ξᵢ*  ·  ξᵢ, ξᵢ* ≥ 0",
                "bullets": [
                    "Inside tube → no penalty.",
                    "Outside tube → slack cost.",
                ],
                "plot_path": "svr-epsilon-tube.png",
            },
            {
                "title": "SVR Hyperparameters",
                "kicker": "Tuning Knobs",
                "layout": "table",
                "table": {
                    "headers": ["Hyperparameter", "Role"],
                    "rows": [
                        [
                            "epsilon (ε)",
                            "Width of the no-penalty tube around prediction",
                        ],
                        [
                            "C",
                            "Penalty strength for points outside the tube",
                        ],
                        [
                            "Kernel",
                            "Controls curve shape (linear, RBF, polynomial)",
                        ],
                    ],
                },
                "note": "Tune ε, C, and the kernel together — changing one often changes how the others behave.",
            },
            {
                "title": "Why SVR?",
                "kicker": "Compared with Traditional Regression",
                "layout": "table",
                "table": {
                    "headers": ["Aspect", "Traditional", "SVR"],
                    "rows": [
                        ["Errors", "Minimize all residuals", "Ignore inside ε-tube"],
                        ["Outliers", "Can be sensitive", "More robust if tuned"],
                        ["Non-linear", "Hand-built features", "Kernels directly"],
                        ["Large data", "Usually faster", "Can get expensive"],
                    ],
                },
                "note": "Scale features · start with RBF · tune C, ε, kernel params.",
            },
            {
                "title": "SVR Visual Intuition",
                "kicker": "Why SVR?",
                "body": "ε-tube ignores small residuals; larger deviations pull the fit.",
                "bullets": [
                    "Inside tube → no loss.",
                    "Kernels capture curves without hand features.",
                ],
                "plot_path": "svr-epsilon-tube.png",
            },
            {
                "title": "SVR Fit Example",
                "kicker": "Why SVR?",
                "body": "Tuned tube + kernel follows curves while resisting small noise.",
                "bullets": [
                    "Scale features first.",
                    "Tune C, ε, kernel with CV.",
                ],
                "plot": "slide-42-2.png",
            },
        ],
    },
    "Decision Tree": {
        "title": "Decision Tree Regression (CART)",
        "kicker": "What is Decision Tree Regressor?",
        "body": "Decision Tree Regressor predicts continuous values by splitting feature space into rule-based regions and assigning a constant value in each region.",
        "bullets": [
            "Captures non-linear relationships.",
            "Easy to interpret if tree depth is controlled.",
            "Handles complex feature interactions.",
        ],
        "plot_path": "decision-tree-split.png",
        "note": "Decision Tree Split: recursive partitioning by feature threshold · Root → branches → leaf nodes.",
        "extra_slides": [
            {
                "title": "Core Concept of Regression Trees",
                "kicker": "Tree Structure",
                "body": "A regression tree is built from decision rules that partition the input space into regions with constant predictions.",
                "bullets": [
                    "Internal node: decision rule (e.g., x ≤ 4.5).",
                    "Branch: outcome of the rule.",
                    "Leaf node: predicted value (usually region mean).",
                ],
            },
            {
                "title": "CART Perspective",
                "kicker": "Classification and Regression Trees",
                "layout": "table",
                "table": {
                    "headers": ["Tree Type", "Target Type", "Typical Criterion"],
                    "rows": [
                        [
                            "Classification Tree",
                            "Categorical target",
                            "Gini / Entropy",
                        ],
                        [
                            "Regression Tree",
                            "Continuous target",
                            "MSE (or MAE)",
                        ],
                    ],
                },
                "note": "This section focuses on the Regression Tree part of CART.",
            },
            {
                "title": "How Splitting Happens (MSE Criterion)",
                "kicker": "Node MSE",
                "body": "CART picks the split that maximally reduces MSE (regression) or Gini/entropy (classification).",
                "formula": "MSE(S) = (1 / |S|) Σ_{i ∈ S} (y_{i} − ȳ_S)²",
                "formula_tex": r"\mathrm{MSE}(S) = \dfrac{1}{|S|} \sum_{i \in S} (y_i - \bar{y}_S)^2",
                "formula_note": "Average squared distance from the node mean ȳ_S",
                "layout": "formula_example",
                "bullets": [
                    "Greedy splits are fast but not globally optimal.",
                    "Depth and min-samples-leaf control overfitting on small datasets.",
                ],
            },
            {
                "title": "Split Score",
                "kicker": "MSE Criterion",
                "body": "A candidate split is scored by the size-weighted average MSE of the left and right child nodes.",
                "formula": "MSE_split = (|S_L|/|S|) MSE(S_L) + (|S_R|/|S|) MSE(S_R)",
                "formula_note": "S_L / S_R = left and right subsets after the candidate split",
                "layout": "formula_example",
                "bullets": [
                    "Choose the split with the minimum weighted MSE.",
                    "Larger child nodes influence the score more (weighted by |S_L| and |S_R|).",
                ],
            },
            {
                "title": "Prediction at Leaf",
                "kicker": "Regression Tree",
                "body": "Once a region is pure enough to stop splitting, the leaf predicts a constant value for every point that lands there.",
                "formula": "ŷ_leaf = (1 / |S_leaf|) Σ_{i ∈ S_leaf} y_{i}",
                "formula_tex": r"\hat{y}_{leaf} = \dfrac{1}{|S_{leaf}|} \sum_{i \in S_{leaf}} y_i",
                "formula_note": "Usually the mean of the training targets in that leaf region",
                "layout": "formula_example",
                "bullets": [
                    "Same prediction for all samples that reach the same leaf.",
                    "Deeper trees create smaller regions and more varied leaf values.",
                ],
            },
            {
                "title": "Decision Tree Regression Algorithm",
                "kicker": "Simple Steps",
                "body": "Recursive partitioning by feature threshold · Root → branches → leaf nodes.",
                "layout": "table",
                "table": {
                    "headers": ["Step", "Action"],
                    "rows": [
                        ["1", "Start with all data at the root node."],
                        ["2", "Evaluate candidate splits across features."],
                        ["3", "Compute weighted MSE for each split."],
                        ["4", "Select the split with minimum error."],
                        ["5", "Repeat recursively on child nodes."],
                        ["6", "Stop based on rules (max depth, min samples)."],
                        ["7", "Predict using leaf mean value."],
                    ],
                },
                "note": "Decision Tree Split: greedy recursive partitioning until a stopping rule fires.",
            },
            {
                "title": "Decision Tree Fit (Visual)",
                "kicker": "Algorithm · Static Plot",
                "body": "A fitted regression tree partitions x into intervals and predicts a constant in each leaf region.",
                "bullets": [
                    "Piecewise-constant prediction surface.",
                    "More depth → more steps and finer partitions.",
                ],
                "plot": "slide-46-2.png",
            },
            {
                "title": "Important Hyperparameters",
                "kicker": "Decision Tree Regression",
                "layout": "table",
                "table": {
                    "headers": ["Hyperparameter", "Effect"],
                    "rows": [
                        ["max_depth", "Controls model complexity"],
                        [
                            "min_samples_split",
                            "Minimum samples required to split a node",
                        ],
                        [
                            "min_samples_leaf",
                            "Minimum samples in each leaf",
                        ],
                        [
                            "max_leaf_nodes",
                            "Limits number of terminal regions",
                        ],
                    ],
                },
                "note": "These settings reduce overfitting and improve generalization.",
            },
        ],
    },
    "Random Forest": {
        "title": "Random Forest Regression",
        "kicker": "What is Random Forest Regressor?",
        "body": "Ensemble of many trees — average their predictions.",
        "formula": "ŷ_{RF}(x) = (1 / N_{trees}) Σ_{t = 1}^{N_{trees}} ŷ_{t}(x)",
        "formula_tex": r"\hat{y}_{RF}(x) = \dfrac{1}{N_{trees}} \sum_{t=1}^{N_{trees}} \hat{y}_t(x)",
        "formula_note": "Average of individual tree predictions",
        "bullets": [
            "Bootstrap samples → many trees.",
            "Averaging reduces variance vs one deep tree.",
        ],
        "plot_path": "random-forest-ensemble.png",
        "note": "Also see fitted RF curve: slide-46-3.",
        "extra_slides": [
            {
                "title": "Decision Tree vs Random Forest",
                "kicker": "Regression Comparison",
                "layout": "table",
                "table": {
                    "headers": ["Aspect", "Decision Tree", "Random Forest"],
                    "rows": [
                        ["Model", "Single tree", "Many trees averaged"],
                        ["Split", "MSE per split", "MSE trees, then average"],
                        ["Overfit risk", "Higher", "Lower"],
                        ["Stability", "Sensitive", "More stable"],
                        ["Prediction", "One tree output", "Mean of trees"],
                        ["Interpretability", "High", "Medium"],
                    ],
                },
                "note": "Tree for interpretability · Forest for stronger performance.",
            },
            {
                "title": "Feature Importance (Random Forest)",
                "kicker": "Decision Tree vs Random Forest",
                "body": "Importance scores show how much each feature helped reduce error.",
                "bullets": [
                    "Higher bar → more contribution.",
                    "Guide for insight — not causal proof.",
                ],
                "plot": "slide-48-1.png",
            },
            {
                "title": "Evaluating Regression Models",
                "kicker": "Key Performance Metrics",
                "layout": "table",
                "table": {
                    "headers": ["Metric", "Formula", "Tells Us"],
                    "rows": [
                        ["MAE", "(1/n) Σ |yᵢ − ŷᵢ|", "Average error size"],
                        ["MSE", "(1/n) Σ (yᵢ − ŷᵢ)²", "Penalizes large errors"],
                        ["RMSE", "√MSE", "Error on y scale"],
                        ["R²", "1 − SS_res/SS_tot", "Variance explained"],
                    ],
                },
                "note": "Use more than one metric (e.g. RMSE + R²).",
            },
        ],
    },
    "Ridge / Lasso": {
        "title": "Regularization Methods (Why Needed?)",
        "kicker": "Overfitting Problem",
        "body": "Too flexible → fits noise instead of true patterns.",
        "formula": "min_β  Σ (y_{i} − ŷ_{i})² + λ · Ω(β)",
        "formula_note": "data loss + λ × penalty",
        "bullets": [
            "Low train error, weak test performance.",
            "High variance / unstable predictions.",
        ],
        "plot_path": "regularization-path.png",
        "note": "Penalize large coefficients to control complexity.",
        "extra_slides": [
            {
                "title": "Regularization Objective",
                "kicker": "Data Loss + Penalty",
                "body": "Trade off fit vs small coefficients.",
                "formula": "min_β  data loss + λ · penalty",
                "formula_note": "loss = Σ(yᵢ − ŷᵢ)²  ·  penalty = Ω(β)",
                "layout": "formula_example",
                "bullets": [
                    "λ = 0 → ordinary least squares.",
                    "Larger λ → stronger shrinkage · tune with CV.",
                ],
            },
            {
                "title": "Ridge, Lasso, and Elastic Net",
                "kicker": "Regularization Families",
                "body": "Different penalties shrink coefficients differently as λ grows.",
                "bullets": [
                    "Ridge (L2): shrink all — good with correlated features.",
                    "Lasso (L1): can zero coeffs — feature selection.",
                ],
                "plot_path": "regularization-path.png",
                "note": "Elastic Net blends L1 + L2 for correlated sparse settings.",
            },
            {
                "title": "Without Regularization",
                "kicker": "Ordinary Least Squares",
                "body": "Minimize squared error only — no coefficient penalty.",
                "formula": "min_β  Σ (y_{i} − ŷ_{i})²",
                "formula_note": "Same as λ = 0",
                "layout": "formula_example",
                "bullets": [
                    "No shrinkage.",
                    "Can overfit with many / correlated features.",
                ],
            },
            {
                "title": "Ridge Regression (L2)",
                "kicker": "Penalty: sum of squared coefficients",
                "body": "L2 penalty shrinks coefficients; rarely zeros them.",
                "formula": "min_β  Σ (y_{i} − ŷ_{i})² + λ Σ βⱼ²",
                "formula_note": "λ ≥ 0 controls shrinkage",
                "layout": "formula_example",
                "bullets": [
                    "Smooth shrinkage of all coefficients.",
                    "Good when many features correlate.",
                ],
            },
            {
                "title": "Lasso Regression (L1)",
                "kicker": "Penalty: sum of absolute coefficients",
                "body": "L1 penalty can set some coefficients exactly to zero.",
                "formula": "min_β  Σ (y_{i} − ŷ_{i})² + λ Σ |βⱼ|",
                "formula_note": "Larger λ → sparser model",
                "layout": "formula_example",
                "bullets": [
                    "Built-in feature selection.",
                    "Best when few predictors truly matter.",
                ],
            },
            {
                "title": "Elastic Net (L1 + L2)",
                "kicker": "Blend of Lasso and Ridge",
                "body": "Combine L1 + L2 for sparsity with more stable correlation behavior.",
                "formula": "min_β  Σ (y_{i} − ŷ_{i})² + λ_{1} Σ |βⱼ| + λ_{2} Σ βⱼ²",
                "formula_note": "λ₁ → sparsity  ·  λ₂ → ridge shrinkage",
                "layout": "formula_example",
                "bullets": [
                    "Good for correlated feature groups.",
                    "Tune λ₁, λ₂ (or mix ratio) with CV.",
                ],
            },
            {
                "title": "When to Use Each Method",
                "kicker": "Ridge · Lasso · Elastic Net",
                "body": "Match the penalty to the feature structure you expect.",
                "bullets": [
                    "Ridge: many small/medium useful features.",
                    "Lasso: few important features expected.",
                    "Elastic Net: correlated features + need sparsity.",
                ],
            },
            {
                "title": "Regularization Comparison Table",
                "kicker": "Ridge vs Lasso vs Elastic Net",
                "layout": "table",
                "table": {
                    "headers": [
                        "Method",
                        "Penalty Type",
                        "Main Behavior",
                        "Feature Selection",
                    ],
                    "rows": [
                        [
                            "Ridge",
                            "L2 (sum βⱼ²)",
                            "Shrinks coefficients smoothly",
                            "No",
                        ],
                        [
                            "Lasso",
                            "L1 (sum |βⱼ|)",
                            "Shrinks and can set some coefficients to zero",
                            "Yes",
                        ],
                        [
                            "Elastic Net",
                            "L1 + L2",
                            "Combines shrinkage + selection",
                            "Yes",
                        ],
                    ],
                },
                "note": "Scale features before Ridge/Lasso/Elastic Net — penalties depend on coefficient magnitude.",
            },
            {
                "title": "Regression Models — Pros & Cons",
                "kicker": "Session Summary",
                "layout": "table",
                "table": {
                    "headers": ["Model", "Advantages", "Disadvantages"],
                    "rows": [
                        [
                            "Linear",
                            "Simple, fast, interpretable",
                            "Linear only; outlier-sensitive",
                        ],
                        [
                            "Polynomial",
                            "Captures curves",
                            "Degree sensitive; can overfit",
                        ],
                        [
                            "SVR",
                            "Kernels; robust if tuned",
                            "Needs scaling; slow on big data",
                        ],
                        [
                            "Decision Tree",
                            "Interpretable; no scaling",
                            "Overfits; unstable",
                        ],
                        [
                            "Random Forest",
                            "Accurate; lower variance",
                            "Less interpretable; heavier",
                        ],
                    ],
                },
                "note": "Interpretability → Linear/Tree · Performance → Random Forest.",
            },
            {
                "title": "Model Summary (Name + Equation)",
                "kicker": "Session Summary",
                "layout": "table",
                "table": {
                    "headers": ["Model", "Core Equation"],
                    "rows": [
                        ["Linear Regression", "ŷ = b₀ + b₁x"],
                        [
                            "Polynomial Regression",
                            "ŷ = b₀ + b₁x + b₂x² + ⋯ + b_d x^d",
                        ],
                        ["SVR", "f(x) = wᵀ φ(x) + b"],
                        [
                            "Decision Tree Regression",
                            "ŷ_leaf = mean(y in S_leaf)",
                        ],
                        [
                            "Random Forest Regression",
                            "ŷ_RF(x) = (1/N_trees) Σ ŷ_t(x)",
                        ],
                    ],
                },
                "note": "If interpretability is priority: Linear Regression or Decision Tree. If predictive performance is priority: Random Forest.",
            },
            {
                "title": "Visual Comparison of Regression Models",
                "kicker": "Same Dataset · Different Fits",
                "body": "Same data · five model families side by side.",
                "bullets": [
                    "Linear/Ridge: fast · watch curves.",
                    "Tree/RF: flexible · weak extrapolation.",
                ],
                "plot": "slide-54-1.png",
                "note": "Visual output from Python code.",
            },
            {
                "title": "Model Families at a Glance",
                "kicker": "Session Closing",
                "layout": "table",
                "table": {
                    "headers": ["Model", "Strength", "Watch Out"],
                    "rows": [
                        [
                            "Linear / Ridge",
                            "Interpretable, fast",
                            "Nonlinear patterns",
                        ],
                        [
                            "SVR",
                            "Margin-based, kernels",
                            "Slow on huge data",
                        ],
                        [
                            "Tree / RF",
                            "Nonlinear, little scaling",
                            "Extrapolation",
                        ],
                    ],
                },
                "note": "Visual output from Python code is shown on the previous slide.",
            },
        ],
    },
}

BIG_PICTURE = {
    "title": "Where Are We in the Bootcamp?",
    "focus": "After preprocessing — learn models that predict continuous targets.",
    "current": "S2",
    "weeks": [
        {
            "label": "Week 1",
            "sessions": [
                ("S1", "Foundations & Preprocessing"),
                ("S2", "Regression Models"),
                ("S3", "Classification Basics"),
                ("S4", "Naive Bayes & Trees"),
                ("S5", "SVM & Kernels"),
                ("S6", "Clustering & PCA"),
            ],
        },
        {"label": "Week 2", "sessions": [("S7", "Deep Learning")]},
        {
            "label": "Week 3",
            "sessions": [
                ("S8", "NLP Fundamentals"),
                ("S9", "Tokenization"),
                ("S10", "Text Analysis & NER"),
                ("S11", "Language Modeling"),
                ("S12", "Embeddings & RNNs"),
                ("S13", "Seq2Seq & NMT"),
            ],
        },
        {
            "label": "Week 4",
            "sessions": [
                ("S14", "Generative AI"),
                ("S15", "RAG Systems"),
                ("S16", "MLOps"),
            ],
        },
    ],
}


def slide_title(prs, total):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.2), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.7),
        Inches(11.5),
        Inches(0.95),
        s["section_title"],
        size=44,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.8), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.15), Inches(11), Inches(0.5), s["subtitle"], size=17, color=MUTED)
    add_text(slide, MARGIN, Inches(6.7), Inches(6), Inches(0.3), "ETRA", size=12, bold=True, color=PRIMARY)
    add_text(
        slide,
        Inches(8.5),
        Inches(6.7),
        Inches(4),
        Inches(0.3),
        s["trainer_line"],
        size=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def slide_big_picture(prs, total, index):
    bp = BIG_PICTURE
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  Big picture", f"{index:02d}")
    title_block(slide, bp["title"], bp["focus"])

    weeks = bp["weeks"]
    col_w = Inches(2.9)
    gap = Inches(0.18)
    top = Inches(2.35)

    for wi, week in enumerate(weeks):
        x = MARGIN + wi * (col_w + gap)
        add_text(slide, x, top, col_w, Inches(0.3), week["label"], size=12, bold=True, color=PRIMARY)
        rect(slide, x, top + Inches(0.32), col_w - Inches(0.2), Inches(0.012), SOFT)

        for si, (sid, title) in enumerate(week["sessions"]):
            y = top + Inches(0.5) + Inches(si * 0.58)
            current = sid == bp["current"]
            if current:
                rect(slide, x, y + Inches(0.08), Inches(0.05), Inches(0.28), PRIMARY)
            add_text(
                slide,
                x + Inches(0.16),
                y,
                Inches(0.45),
                Inches(0.4),
                sid,
                size=11,
                bold=current,
                color=PRIMARY if current else SECONDARY,
                anchor=MSO_ANCHOR.MIDDLE,
            )
            add_text(
                slide,
                x + Inches(0.6),
                y,
                col_w - Inches(0.75),
                Inches(0.4),
                title,
                size=11,
                bold=current,
                color=INK if current else MUTED,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    content_footer(slide, index, total)


def slide_section_divider(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    logo(slide, height=Inches(0.4))

    add_text(slide, MARGIN, Inches(2.15), Inches(11), Inches(0.35), s["eyebrow"], size=14, color=SECONDARY)
    add_text(
        slide,
        MARGIN,
        Inches(2.65),
        Inches(11.5),
        Inches(0.9),
        s["section_title"],
        size=42,
        bold=True,
        color=PRIMARY,
    )
    bar = rect(slide, MARGIN, Inches(3.7), Inches(1.4), Inches(0.06), PRIMARY)
    gradient_fill(bar, PRIMARY, SECONDARY, 0)
    add_text(slide, MARGIN, Inches(4.05), Inches(11), Inches(0.4), s["focus"], size=16, color=MUTED)

    x = MARGIN
    for topic in s["topics"]:
        w = Inches(min(2.8, 0.12 * len(topic) + 1.1))
        if x + w > Inches(12.4):
            break
        soft_card(slide, x, Inches(5.25), w, Inches(0.42), fill=SOFT)
        add_text(slide, x, Inches(5.3), w, Inches(0.32), topic, size=11, color=PRIMARY, align=PP_ALIGN.CENTER)
        x += w + Inches(0.14)

    content_footer(slide, index, total)


def slide_agenda(prs, total, index):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, s["focus"], "What we will cover in this session")
    bullets(slide, s["topics"], top=Inches(2.35), size=20)
    content_footer(slide, index, total)


def _add_plot(slide, name, left, top, width, max_height, *, folder=None):
    """Fit a plot PNG inside a box (same idea as add_diagram)."""
    from PIL import Image

    base = folder or PLOTS
    path = base / name
    if not path.is_file() and folder is None:
        path = DIAGRAMS / name
    if not path.is_file():
        return None
    with Image.open(path) as im:
        px_w, px_h = im.size
    if px_w <= 0 or px_h <= 0:
        return None
    aspect = px_w / px_h
    max_w_in = width.inches
    max_h_in = max_height.inches
    fit_w = min(max_w_in, max_h_in * aspect)
    fit_h = fit_w / aspect
    if fit_h > max_h_in:
        fit_h = max_h_in
        fit_w = fit_h * aspect
    x = left.inches + (max_w_in - fit_w) / 2
    return slide.shapes.add_picture(
        str(path),
        Inches(x),
        top,
        width=Inches(fit_w),
        height=Inches(fit_h),
    )


def slide_topic_rich(prs, total, index, content):
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    dense = len(items) >= 4 or (len(items) >= 3 and has_note)
    pitch = 0.50 if dense else 0.62
    bsize = 14 if dense else 16

    if content.get("body"):
        soft_card(slide, MARGIN, Inches(2.25), Inches(12.0), Inches(1.15), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.45),
            Inches(11.3),
            Inches(0.8),
            content["body"],
            size=15,
            color=INK,
        )
        bullet_top = Inches(3.6)
    else:
        bullet_top = Inches(2.35)
        bsize = 15 if dense else 17

    # Keep bullets above the note/footer band
    note_y = 6.35 if has_note else 6.7
    max_items = max(1, int((note_y - bullet_top.inches - 0.1) / pitch))
    bullets(
        slide,
        items[:max_items],
        top=bullet_top,
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.45),
            Inches(12),
            Inches(0.3),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_linear_intro(prs, total, index, content):
    """Definition + formula + optional scatter/line plot."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    plot_name = content.get("plot") or content.get("plot_path")
    has_plot = bool(plot_name)
    has_formula = bool(content.get("formula"))
    has_note = bool(content.get("note"))
    is_frac = has_formula and (
        is_fraction_formula(content["formula"]) or bool(content.get("formula_tex"))
    )

    col_w = Inches(6.3) if has_plot else Inches(12.0)
    text_w = Inches(5.8) if has_plot else Inches(11.3)
    bullet_w = Inches(5.6) if has_plot else Inches(11.2)

    soft_card(slide, MARGIN, Inches(2.2), col_w, Inches(1.2), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.3),
        Inches(2.35),
        text_w,
        Inches(0.95),
        content.get("body") or "",
        size=13,
        color=INK,
    )

    if has_formula:
        formula_card_h = 1.65 if is_frac else 1.35
        formula_box_h = 0.95 if is_frac else 0.65
        soft_card(slide, MARGIN, Inches(3.55), col_w, Inches(formula_card_h), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.3),
            Inches(3.65),
            text_w,
            Inches(0.22),
            "Formula",
            size=11,
            bold=True,
            color=SECONDARY,
        )
        add_formula(
            slide,
            Inches(MARGIN.inches + 0.3),
            Inches(3.9),
            text_w,
            Inches(formula_box_h),
            content["formula"],
            size=18 if is_frac else 20,
            bold=True,
            color=PRIMARY,
            formula_tex=content.get("formula_tex"),
        )
        if content.get("formula_note"):
            note_y = 3.9 + formula_box_h + 0.05
            add_text(
                slide,
                Inches(MARGIN.inches + 0.3),
                Inches(note_y),
                text_w,
                Inches(0.3),
                content["formula_note"],
                size=10,
                color=MUTED,
            )
        bullet_top = 3.55 + formula_card_h + 0.12
        pitch = 0.46 if (has_note or len(items) >= 2) else 0.55
        bsize = 12
    else:
        bullet_top = 3.55
        pitch = 0.50 if (has_note or len(items) >= 3) else 0.58
        bsize = 13

    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    # With formula+plot keep at most 2 bullets to avoid crowding
    if has_formula and has_plot:
        max_items = min(max_items, 2)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        left=MARGIN,
        width=bullet_w,
        pitch=pitch,
        item_height=Inches(0.42),
    )

    plot_folder = DIAGRAMS if content.get("plot_path") else PLOTS
    if plot_name:
        soft_card(slide, Inches(7.15), Inches(2.2), Inches(5.45), Inches(4.15), fill=SOFT)
        _add_plot(
            slide,
            plot_name,
            Inches(7.35),
            Inches(2.4),
            Inches(5.05),
            Inches(3.75),
            folder=plot_folder,
        )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.28),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_formula_example(prs, total, index, content):
    """Full-width practical example with formula and interpretation."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    items = content.get("bullets") or []
    has_note = bool(content.get("note"))
    is_frac = is_fraction_formula(content.get("formula") or "") or bool(
        content.get("formula_tex")
    )
    dense = len(items) >= 3 or (len(items) >= 2 and has_note)

    body_top = Inches(2.15)
    if content.get("body"):
        soft_card(slide, MARGIN, body_top, Inches(12.0), Inches(0.7), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.28),
            Inches(11.3),
            Inches(0.45),
            content["body"],
            size=14,
            color=INK,
        )
        formula_top = Inches(3.0)
    else:
        formula_top = Inches(2.25)

    formula_h = 1.55 if is_frac else 1.25
    soft_card(slide, MARGIN, formula_top, Inches(12.0), Inches(formula_h), fill=SOFT)
    add_formula(
        slide,
        Inches(MARGIN.inches + 0.4),
        Inches(formula_top.inches + 0.12),
        Inches(11.2),
        Inches(0.95 if is_frac else 0.7),
        content["formula"],
        size=24 if is_frac else 26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
        formula_tex=content.get("formula_tex"),
    )
    if content.get("formula_note"):
        add_text(
            slide,
            Inches(MARGIN.inches + 0.4),
            Inches(formula_top.inches + formula_h - 0.35),
            Inches(11.2),
            Inches(0.28),
            content["formula_note"],
            size=12,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )

    bullet_top = formula_top.inches + formula_h + 0.12
    pitch = 0.46 if dense else 0.56
    bsize = 13 if dense else 15
    note_band = 6.45 if has_note else 6.75
    max_items = max(1, int((note_band - bullet_top - 0.05) / pitch))
    if has_note:
        max_items = min(max_items, 3)
    bullets(
        slide,
        items[:max_items],
        top=Inches(bullet_top),
        size=bsize,
        pitch=pitch,
        width=Inches(11.2),
        item_height=Inches(0.42),
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.55),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_topic_table(prs, total, index, content):
    from pptx.util import Pt

    s = SESSION
    table = content["table"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    headers = table["headers"]
    rows = table["rows"]
    has_note = bool(content.get("note"))
    n_rows = 1 + len(rows)

    # Prefer kicker over a long body when the table is dense
    show_body = bool(content.get("body")) and len(rows) <= 4
    compact = len(rows) >= 5 or (show_body and has_note and len(rows) >= 4)
    body_size = 11 if compact else 13
    header_size = 10 if compact else 12

    table_top = 2.25
    if show_body:
        soft_card(slide, MARGIN, Inches(2.15), Inches(12.0), Inches(0.55), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(2.25),
            Inches(11.3),
            Inches(0.38),
            content["body"],
            size=13,
            color=INK,
        )
        table_top = 2.85

    footer_limit = 6.55
    note_reserve = 0.55 if has_note else 0.0
    available = footer_limit - table_top - note_reserve
    row_h = min(0.52 if not compact else 0.40, available / n_rows)
    row_h = max(0.30, row_h)
    table_h = row_h * n_rows

    # If note would still collide, drop the note rather than overlap the table
    note_top = table_top + table_h + 0.08
    render_note = has_note and (note_top + 0.45) <= 6.9

    shape = slide.shapes.add_table(
        rows=n_rows,
        cols=len(headers),
        left=MARGIN,
        top=Inches(table_top),
        width=Inches(12.0),
        height=Inches(table_h),
    )
    tbl = shape.table

    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(header_size)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = "Helvetica"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(body_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = PRIMARY if c == 0 else INK
                    run.font.name = "Helvetica"
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if r % 2 else SOFT_2

    if render_note:
        soft_card(slide, MARGIN, Inches(note_top), Inches(12.0), Inches(0.45), fill=SOFT)
        add_text(
            slide,
            MARGIN + Inches(0.35),
            Inches(note_top + 0.08),
            Inches(11.3),
            Inches(0.32),
            content["note"],
            size=12,
            color=MUTED,
        )

    content_footer(slide, index, total)


def slide_full_diagram(prs, total, index, content):
    """Title + optional short body + full-width diagram."""
    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, content["title"], content.get("kicker"))

    diagram_top = Inches(2.15)
    if content.get("body"):
        add_text(
            slide,
            MARGIN,
            Inches(2.1),
            Inches(12.0),
            Inches(0.35),
            content["body"],
            size=13,
            color=MUTED,
        )
        diagram_top = Inches(2.45)

    plot_name = content.get("plot_path") or content.get("plot")
    folder = DIAGRAMS if content.get("plot_path") else PLOTS
    soft_card(slide, MARGIN, diagram_top, Inches(12.0), Inches(4.35), fill=SOFT)
    _add_plot(
        slide,
        plot_name,
        Inches(MARGIN.inches + 0.2),
        Inches(diagram_top.inches + 0.15),
        Inches(11.6),
        Inches(4.05),
        folder=folder,
    )

    if content.get("note"):
        add_text(
            slide,
            MARGIN,
            Inches(6.9),
            Inches(12),
            Inches(0.25),
            content["note"],
            size=11,
            color=MUTED,
        )

    content_footer(slide, index, total)


def _render_content_slide(prs, total, index, content):
    layout = content.get("layout")
    if layout == "diagram":
        slide_full_diagram(prs, total, index, content)
    elif layout == "table" or (content.get("table") and layout not in ("formula_example", "diagram")):
        slide_topic_table(prs, total, index, content)
    elif layout == "formula_example":
        slide_formula_example(prs, total, index, content)
    elif content.get("formula") or content.get("plot") or content.get("plot_path"):
        slide_linear_intro(prs, total, index, content)
    else:
        slide_topic_rich(prs, total, index, content)


def slide_topic(prs, total, index, topic_index, topic):
    """Returns number of slides added."""
    content = TOPIC_CONTENT.get(topic)
    if content:
        _render_content_slide(prs, total, index, content)
        n = 1
        for extra in content.get("extra_slides") or []:
            _render_content_slide(prs, total, index + n, extra)
            n += 1
        return n

    s = SESSION
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_light(slide)
    right_rail(slide)
    content_header(slide, f"{s['section_tag']}  ·  {s['section_title']}", f"{index:02d}")
    title_block(slide, topic, f"Topic {topic_index} of {len(s['topics'])}")
    soft_card(slide, MARGIN, Inches(2.45), Inches(12.0), Inches(3.7), fill=SOFT)
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(3.55),
        Inches(11.2),
        Inches(0.55),
        topic,
        size=26,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        MARGIN + Inches(0.4),
        Inches(4.25),
        Inches(11.2),
        Inches(0.4),
        "Content coming next — paste the teaching block when ready.",
        size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    content_footer(slide, index, total)
    return 1


def _topic_slide_count(topic: str) -> int:
    content = TOPIC_CONTENT.get(topic)
    if not content:
        return 1
    return 1 + len(content.get("extra_slides") or [])


def main() -> None:
    s = SESSION
    topic_pages = sum(_topic_slide_count(t) for t in s["topics"])
    total = 4 + topic_pages
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, total)
    n = 2
    slide_big_picture(prs, total, n)
    n = 3
    slide_section_divider(prs, total, n)
    n = 4
    slide_agenda(prs, total, n)
    for i, topic in enumerate(s["topics"], start=1):
        n += 1
        added = slide_topic(prs, total, n, i, topic)
        n += added - 1

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print("Brand: ETRA Design System v1.0")
    print("Topics:", " · ".join(s["topics"]))


if __name__ == "__main__":
    main()
